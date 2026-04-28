from __future__ import annotations

import io
import json
import logging
import re
import ssl
import sys
import tempfile
import time
import zipfile
from collections import Counter
from datetime import date, datetime
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

import requests

from minimum_workflow.contracts import ExtractionResult

try:
    import pythoncom
except ImportError:
    pythoncom = None

try:
    import win32com.client as win32_client
except ImportError:
    win32_client = None


# 先覆盖常见文本编码，后续可按实际样本继续扩充。
TEXT_ENCODINGS = ("utf-8-sig", "utf-8", "gb18030")
MINERU_API_BASE = "https://mineru.net/api/v4"
MINERU_REQUEST_TIMEOUT = 60
MINERU_DOWNLOAD_TIMEOUT = 120
MINERU_UPLOAD_TIMEOUT = 300
MINERU_RATE_LIMIT_MAX_RETRIES = 3
MINERU_RATE_LIMIT_RETRY_DELAY_SECONDS = 10
MINERU_HTTP_MAX_RETRIES = 3
MINERU_HTTP_BACKOFF_SECONDS = 1.0
BID_TABLE_HEADERS = ("发布日期", "项目名称", "采购人", "预算金额", "截止时间")
DOCX_NS = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
DOCX_TAG_PREFIX = f"{{{DOCX_NS['w']}}}"
# 目录里图片数量远高于正式文档图片，先按“文档图片白名单、其余默认纯照片跳过”的保守口径处理。
DOCUMENT_IMAGE_KEYWORDS = (
    "证书",
    "报告",
    "检验",
    "检测",
    "检定",
    "校准",
    "资质",
    "发票",
    "合同",
    "协议",
    "报价",
    "方案",
    "政策",
    "通知",
    "公告",
    "公示",
    "条例",
    "规划",
    "清单",
    "名录",
    "参数",
    "手册",
    "说明书",
    "联系人",
    "名片",
    "表",
    "登记",
    "批复",
    "函",
)
PHOTO_SKIP_KEYWORDS = (
    "照片",
    "图片",
    "现场",
    "航拍",
    "实拍",
    "合影",
    "活动",
    "宣传图",
    "效果图",
    "封面图",
    "配图",
    "img",
    "dji",
)
PAGED_IMAGE_NAME_KEYWORDS = ("封面", "目录", "页", "page", "scan", "扫描")
WECHAT_IMAGE_KEYWORDS = ("微信图片", "mmexport")
WEBPAGE_SHELL_SIGNAL_KEYWORDS = ("行业搜索引擎", "低空经济资源网", "QUICK LINKS")
WEBPAGE_NAVIGATION_LINE_SET = {
    "首页",
    "新闻资讯",
    "国内资讯",
    "国际资讯",
    "政策法规",
    "法律法规",
    "国家政策",
    "地方政策",
    "行业标准",
    "产品库",
    "evtol",
    "飞行汽车",
    "无人机",
    "科普研学",
    "直升机",
    "固定翼飞机",
    "平台方案",
    "企业库",
    "低空制造业",
    "低空运营业",
    "低空基建与信息服务业",
    "低空配套业",
    "应用场景",
    "场景分类",
    "低空案例",
    "低空报告",
    "行业报告",
    "低空词典",
    "专栏文章",
    "会展赛事",
    "caac考证",
    "考试介绍",
    "招生简章",
    "招生范围",
    "指导价格",
    "考点查询",
    "我要报名",
    "问题解答",
    "飞手资源库",
    "往期招生",
    "专题专区",
    "院校专题",
    "关于我们",
    "公司新闻",
    "联系我们",
    "公众号",
    "微信小程序",
    "contact us",
    "quick links",
}
WEBPAGE_FOOTER_LINE_PREFIXES = (
    "扫码咨询",
    "前一个：",
    "后一个：",
    "免责声明",
    "公众号",
    "微信小程序",
    "联系我们",
    "CONTACT US",
    "电话：",
    "手机：",
    "邮箱：",
    "地址：",
    "快速链接",
    "QUICK LINKS",
)
WEBPAGE_CREATED_AT_PATTERN = re.compile(r"^创建时间[:：]")
WEBPAGE_EMPTY_SOURCE_PATTERN = re.compile(r"^来源[:：]\s*$")
WEBPAGE_TRAFFIC_PATTERN = re.compile(r"^浏览量[:：]")


def normalize_preview(text: str, limit: int = 300) -> str:
    # 清洗掉 HTML 注释（如拆分链路的 "<!-- 分片：xxx -->"）与行首 markdown 标题符，避免它们污染摘要。
    cleaned = re.sub(r"<!--.*?-->", " ", text or "", flags=re.DOTALL)
    cleaned = re.sub(r"(?m)^\s*#+\s*", "", cleaned)
    return re.sub(r"\s+", " ", cleaned).strip()[:limit]



def build_mineru_ssl_context() -> ssl.SSLContext:
    context = ssl.create_default_context()
    if sys.platform.startswith("win"):
        context.load_default_certs()
    return context


class MinerUSSLAdapter(requests.adapters.HTTPAdapter):
    def init_poolmanager(self, connections: int, maxsize: int, block: bool = False, **pool_kwargs: Any) -> None:
        pool_kwargs["ssl_context"] = build_mineru_ssl_context()
        super().init_poolmanager(connections, maxsize, block=block, **pool_kwargs)


# MinerU 走 requests 时显式接入系统证书链，避免 Windows 环境下因缺少本地 CA 而直接握手失败。
def create_mineru_session() -> requests.Session:
    session = requests.Session()
    session.mount("https://", MinerUSSLAdapter())
    return session


# 网页转存 Markdown 常会把站点导航和页脚一起带进正文；这里只在命中特定站点壳层特征时做保守清洗。
def normalize_text_line(line: str) -> str:
    return re.sub(r"\s+", " ", re.sub(r"^#+\s*", "", line or "")).strip()


def looks_like_webpage_shell_text(text: str) -> bool:
    signal_hits = sum(1 for keyword in WEBPAGE_SHELL_SIGNAL_KEYWORDS if keyword in text)
    if "原文: http" in text and signal_hits >= 1:
        return True
    return signal_hits >= 2


def is_webpage_navigation_line(line: str) -> bool:
    normalized = normalize_text_line(line).lower()
    if not normalized:
        return False
    return normalized in WEBPAGE_NAVIGATION_LINE_SET


def is_webpage_footer_start(line: str) -> bool:
    normalized = normalize_text_line(line)
    if not normalized:
        return False
    return any(normalized.startswith(prefix) for prefix in WEBPAGE_FOOTER_LINE_PREFIXES)


def is_webpage_created_at_line(line: str) -> bool:
    return bool(WEBPAGE_CREATED_AT_PATTERN.match(normalize_text_line(line)))


def is_webpage_empty_source_line(line: str) -> bool:
    return bool(WEBPAGE_EMPTY_SOURCE_PATTERN.match(normalize_text_line(line)))


def is_webpage_traffic_line(line: str) -> bool:
    return bool(WEBPAGE_TRAFFIC_PATTERN.match(normalize_text_line(line)))


def collapse_blank_lines(lines: list[str]) -> str:
    collapsed: list[str] = []
    previous_blank = False
    for raw_line in lines:
        stripped_line = raw_line.rstrip()
        if stripped_line:
            collapsed.append(stripped_line)
            previous_blank = False
            continue
        if collapsed and not previous_blank:
            collapsed.append("")
            previous_blank = True
    while collapsed and not collapsed[-1]:
        collapsed.pop()
    return "\n".join(collapsed)


def _split_markdown_frontmatter(text: str) -> tuple[str, str]:
    cleaned_text = (text or "").replace("\r", "")
    if not cleaned_text.startswith("---\n"):
        return "", cleaned_text

    closing_index = cleaned_text.find("\n---\n", 4)
    if closing_index == -1:
        return "", cleaned_text

    frontmatter_end = closing_index + len("\n---\n")
    frontmatter = cleaned_text[:frontmatter_end].rstrip()
    body_text = cleaned_text[frontmatter_end:].lstrip("\n")
    return frontmatter, body_text


def _infer_webpage_source_title(source_path: Path, lines: list[str]) -> str:
    for line in lines[:5]:
        normalized = normalize_text_line(line)
        if normalized and not normalized.startswith("原文:"):
            return normalized
    return normalize_text_line(source_path.stem)


def clean_webpage_shell_text(source_path: Path, text: str) -> str:
    frontmatter_text, body_text = _split_markdown_frontmatter(text)
    shell_candidate = body_text if frontmatter_text else text
    if not looks_like_webpage_shell_text(shell_candidate):
        return text

    lines = shell_candidate.replace("\r", "").split("\n")
    prefix_length = 2 if len(lines) > 1 and normalize_text_line(lines[1]).startswith("原文:") else 1
    source_title = _infer_webpage_source_title(source_path, lines)
    content_start_index = -1

    for index in range(prefix_length, len(lines)):
        if is_webpage_created_at_line(lines[index]):
            content_start_index = index + 1
            break

    if content_start_index == -1:
        anchor_index = -1
        for index in range(prefix_length, len(lines)):
            normalized = normalize_text_line(lines[index])
            if normalized and normalized == source_title:
                anchor_index = index
                break
        if anchor_index > prefix_length:
            navigation_hits = sum(1 for line in lines[prefix_length:anchor_index] if is_webpage_navigation_line(line))
            if navigation_hits >= 8:
                content_start_index = anchor_index + 1

    if content_start_index > prefix_length:
        lines = lines[:prefix_length] + [""] + lines[content_start_index:]

    cleaned_lines: list[str] = []
    footer_started = False
    for index, line in enumerate(lines):
        if index >= prefix_length + 1 and is_webpage_empty_source_line(line):
            footer_started = True
            break
        if index >= prefix_length + 1 and is_webpage_footer_start(line):
            footer_started = True
            break
        if index >= prefix_length + 1 and is_webpage_traffic_line(line):
            continue
        cleaned_lines.append(line)

    cleaned = collapse_blank_lines(cleaned_lines if footer_started else cleaned_lines)
    if not cleaned:
        return text
    if frontmatter_text:
        return f"{frontmatter_text}\n\n{cleaned}"
    return cleaned


def has_meaningful_text(text: str, minimum_length: int = 30) -> bool:
    compact = re.sub(r"\s+", "", text)
    return len(compact) >= minimum_length


def try_markitdown_office_fallback(source_path: Path, file_type: str, local_result: ExtractionResult) -> ExtractionResult:
    if local_result.extraction_status == "已提取文本":
        return local_result

    try:
        from minimum_workflow.sample_docx_extract_to_md import try_extract_with_markitdown
    except Exception as exc:
        logging.debug("markitdown导入失败: %s", exc)
        return local_result

    markitdown_result = try_extract_with_markitdown(source_path, file_type)
    if markitdown_result is None:
        return local_result
    if local_result.note:
        markitdown_result.note = f"{markitdown_result.note} {local_result.note}".strip()
    return markitdown_result


def extract_text(
    source_path: Path,
    file_type: str,
    *,
    enable_ocr: bool = False,
    ocr_token: str | None = None,
) -> ExtractionResult:
    if not source_path.exists():
        return ExtractionResult(
            extractor_name="none",
            extraction_status="源文件不存在",
            extracted_text="",
            preview_text="",
            text_length=0,
            page_count=None,
            source_encoding="",
            note="源文件不存在，当前样例无法抽取。",
        )

    if file_type in {"txt", "markdown", "json", "csv", "log"}:
        return extract_plain_text(source_path)
    if file_type == "word":
        if source_path.suffix.lower() == ".docx":
            local_result = extract_docx_text(source_path)
            return try_markitdown_office_fallback(source_path, file_type, local_result)
        local_result = ExtractionResult(
            extractor_name="word:none",
            extraction_status="待人工复核",
            extracted_text="",
            preview_text="",
            text_length=0,
            page_count=None,
            source_encoding="",
            note="当前仅接入 docx 真实抽取，.doc 仍需后续补充解析链路。",
        )
        return try_markitdown_office_fallback(source_path, file_type, local_result)
    if file_type == "excel":
        local_result = extract_excel_text(source_path)
        return try_markitdown_office_fallback(source_path, file_type, local_result)
    if file_type == "presentation":
        local_result = extract_presentation_text(source_path)
        return try_markitdown_office_fallback(source_path, file_type, local_result)
    if file_type == "pdf":
        return extract_pdf_text(source_path, enable_ocr=enable_ocr, ocr_token=ocr_token)
    if file_type == "image":
        if should_skip_image_file(source_path):
            return skip_image_as_photo(source_path)
        if enable_ocr and ocr_token:
            try:
                ocr_result = extract_image_with_mineru_ocr(source_path, ocr_token)
            except Exception as exc:
                return extract_image_with_ocr_placeholder(source_path, extra_note=f"已尝试真实 OCR，但调用失败：{exc}")
            return finalize_image_result_after_ocr(source_path, ocr_result)
        return extract_image_with_ocr_placeholder(source_path)

    return ExtractionResult(
        extractor_name="none",
        extraction_status="待人工复核",
        extracted_text="",
        preview_text="",
        text_length=0,
        page_count=None,
        source_encoding="",
        note="当前文件类型尚未接入真实文本抽取。",
    )


def extract_plain_text(source_path: Path) -> ExtractionResult:
    for encoding in TEXT_ENCODINGS:
        try:
            text = source_path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue

        return build_text_extraction_result(source_path, text, encoding)

    return ExtractionResult(
        extractor_name="text:none",
        extraction_status="待人工复核",
        extracted_text="",
        preview_text="",
        text_length=0,
        page_count=None,
        source_encoding="",
        note="文本文件编码暂未识别，需后续人工复核或补编码策略。",
    )


def build_text_extraction_result(source_path: Path, text: str, encoding: str) -> ExtractionResult:
    if not text.strip():
        return ExtractionResult(
            extractor_name=f"text:{encoding}",
            extraction_status="待人工复核",
            extracted_text=text,
            preview_text="",
            text_length=len(text),
            page_count=None,
            source_encoding=encoding,
            note="文本文件为空或仅包含空白内容。",
        )

    cleaned_text = clean_webpage_shell_text(source_path, text)
    cleaned_webpage_shell = cleaned_text != text
    extra_metadata = build_bid_summary_metadata(source_path, cleaned_text)
    extracted_text = cleaned_text
    note_parts = ["已完成文本文件读取。"]
    if cleaned_webpage_shell:
        note_parts.append("已清洗网页导航壳层。")
    if extra_metadata:
        extracted_text = build_bid_summary_text(cleaned_text, extra_metadata["招投标表格Markdown"])
        note_parts.append("已识别为招投标汇总列表，已生成结构化表格。")

    preview = normalize_preview(extracted_text)
    return ExtractionResult(
        extractor_name=f"text:{encoding}",
        extraction_status="已提取文本",
        extracted_text=extracted_text,
        preview_text=preview,
        text_length=len(extracted_text),
        page_count=None,
        source_encoding=encoding,
        note=" ".join(note_parts),
        extra_metadata=extra_metadata,
    )


def looks_like_bid_summary_text(source_path: Path, text: str) -> bool:
    source_name = source_path.stem
    if "招投标" in source_name or "采购" in source_name or "汇总" in source_name:
        return True
    header_line = next((line for line in text.splitlines() if all(header in line for header in BID_TABLE_HEADERS)), "")
    if not header_line:
        return False
    data_line_count = sum(1 for line in text.splitlines() if re.match(r"^20\d{2}-\d{2}-\d{2}", line.strip()))
    return data_line_count >= 3


def split_bid_table_line(line: str) -> list[str]:
    parts = [part.strip() for part in re.split(r"\t+|\s{2,}", line.strip()) if part.strip()]
    if len(parts) < 5:
        return []
    publish_date = parts[0]
    deadline = parts[-1]
    budget = parts[-2]
    purchaser = parts[-3]
    project_name = " ".join(parts[1:-3]).strip()
    if not project_name:
        return []
    return [publish_date, project_name, purchaser, budget, deadline]


def parse_bid_summary_rows(text: str) -> list[list[str]]:
    rows: list[list[str]] = []
    header_seen = False
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if not header_seen and all(header in line for header in BID_TABLE_HEADERS):
            header_seen = True
            continue
        if not header_seen:
            continue
        if not re.match(r"^20\d{2}-\d{2}-\d{2}", line):
            continue
        row = split_bid_table_line(line)
        if row:
            rows.append(row)
    return rows


def render_markdown_table(headers: list[str], rows: list[list[str]]) -> str:
    normalized_rows = [headers, *rows]
    column_count = max(len(row) for row in normalized_rows)
    width_counter = Counter(len(row) for row in normalized_rows)
    target_width = max(width_counter, key=lambda key: (width_counter[key], key))
    column_count = max(column_count, target_width)
    padded_rows = [row + [""] * (column_count - len(row)) for row in normalized_rows]
    markdown_lines = [
        "| " + " | ".join(cell.replace("|", r"\|") or "-" for cell in padded_rows[0]) + " |",
        "| " + " | ".join(["---"] * column_count) + " |",
    ]
    for row in padded_rows[1:]:
        markdown_lines.append("| " + " | ".join(cell.replace("|", r"\|") or "-" for cell in row) + " |")
    return "\n".join(markdown_lines)


def build_bid_summary_metadata(source_path: Path, text: str) -> dict[str, Any] | None:
    if not looks_like_bid_summary_text(source_path, text):
        return None

    rows = parse_bid_summary_rows(text)
    if not rows:
        return None

    publish_dates = [row[0] for row in rows if row[0]]
    purchasers = {row[2] for row in rows if row[2]}
    budget_count = sum(1 for row in rows if row[3])
    title = next((line.strip() for line in text.splitlines() if line.strip()), source_path.stem)
    # 去掉 Markdown 标题符与附件前缀，避免把网页转存/手工整理痕迹带进正式标题字段。
    title = re.sub(r"^#+\s*", "", title)
    title = re.sub(r"^附件\s*\d+\s*", "", title).strip()
    table_markdown = render_markdown_table(list(BID_TABLE_HEADERS), rows)
    date_range = ""
    if publish_dates:
        date_range = f"{min(publish_dates)} 至 {max(publish_dates)}"

    return {
        "文档分类": "招投标/商机汇总",
        "推荐模板": "招投标汇总模板",
        "模板归属": "招投标汇总模板",
        "资料层级": "市场情报",
        "标题": title,
        "文件标题": title,
        "证据边界": "招投标汇总属于市场情报整理稿，项目状态、预算口径与中标结果需回查原始公告链接或正式采购平台。",
        "来源形态": "文本汇总文件",
        "招投标记录数": len(rows),
        "招投标发布日期范围": date_range,
        "招投标采购单位数": len(purchasers),
        "招投标预算样本数": budget_count,
        "招投标表格Markdown": table_markdown,
        "分流结果": "待审核",
        "是否适合直接入库": False,
    }


def build_bid_summary_text(original_text: str, table_markdown: str) -> str:
    return "\n\n".join(
        [
            "# 招投标结构化结果",
            table_markdown,
            "# 原始全文",
            original_text.strip(),
        ]
    ).strip()


# 直接解析 docx 的段落与表格，优先保留表格 Markdown，避免参数页被打散。
def extract_docx_text(source_path: Path) -> ExtractionResult:
    try:
        with zipfile.ZipFile(source_path) as archive:
            xml_bytes = archive.read("word/document.xml")
    except Exception as exc:
        return ExtractionResult(
            extractor_name="word:none",
            extraction_status="待人工复核",
            extracted_text="",
            preview_text="",
            text_length=0,
            page_count=None,
            source_encoding="",
            note=f"docx 读取失败：{exc}",
        )

    root = ET.fromstring(xml_bytes)
    body = root.find("w:body", DOCX_NS)
    if body is None:
        return ExtractionResult(
            extractor_name="word:docx",
            extraction_status="待人工复核",
            extracted_text="",
            preview_text="",
            text_length=0,
            page_count=None,
            source_encoding="docx",
            note="docx 未识别到正文节点。",
        )

    blocks: list[str] = []
    for child in body:
        if child.tag == f"{DOCX_TAG_PREFIX}p":
            paragraph = extract_docx_text_from_node(child)
            if paragraph:
                blocks.append(paragraph)
        elif child.tag == f"{DOCX_TAG_PREFIX}tbl":
            table_markdown = render_docx_table_as_markdown(child)
            if table_markdown:
                blocks.append(table_markdown)

    extracted_text = "\n\n".join(blocks).strip()
    if not extracted_text:
        return ExtractionResult(
            extractor_name="word:docx",
            extraction_status="待人工复核",
            extracted_text="",
            preview_text="",
            text_length=0,
            page_count=None,
            source_encoding="docx",
            note="docx 中未识别到有效正文内容。",
        )

    extra_metadata = build_bid_summary_metadata(source_path, extracted_text)
    if extra_metadata:
        extracted_text = build_bid_summary_text(extracted_text, extra_metadata["招投标表格Markdown"])
        note = "已完成 docx 提取，并识别为招投标汇总列表，已生成结构化表格。"
    else:
        note = "已完成 docx 正文与表格提取。"
    return ExtractionResult(
        extractor_name="word:docx",
        extraction_status="已提取文本",
        extracted_text=extracted_text,
        preview_text=normalize_preview(extracted_text),
        text_length=len(extracted_text),
        page_count=None,
        source_encoding="docx",
        note=note,
        extra_metadata=extra_metadata,
    )


# 复用 Word XML 文本节点抽取，保持段落与表格单元格的清洗规则一致。
def extract_docx_text_from_node(node: ET.Element) -> str:
    texts = [text_node.text or "" for text_node in node.findall(".//w:t", DOCX_NS)]
    return re.sub(r"\s+", " ", "".join(texts)).strip()


# 把表格单元格里的多段文本收敛成纯文本，避免在终稿中留下 <br> 这类 HTML 标签。
def merge_docx_table_cell_lines(cell_lines: list[str]) -> str:
    normalized_lines = [line.strip() for line in cell_lines if line.strip()]
    if not normalized_lines:
        return "-"

    merged = normalized_lines[0]
    for line in normalized_lines[1:]:
        if re.fullmatch(r"[0-9\-]{7,}", line):
            merged = f"{merged} {line}"
        elif (
            len(merged) <= 8
            and len(line) <= 8
            and re.fullmatch(r"[\u4e00-\u9fffA-Za-z]+", merged)
            and re.fullmatch(r"[\u4e00-\u9fffA-Za-z]+", line)
        ):
            merged = f"{merged}{line}"
        else:
            merged = f"{merged}；{line}"
    return merged.replace("|", r"\|")


# 读取 Word 单元格合并属性：gridSpan（横向合并列数）与 vMerge（纵向合并状态）。
def _docx_cell_grid_span(cell: ET.Element) -> int:
    tc_pr = cell.find("./w:tcPr", DOCX_NS)
    if tc_pr is None:
        return 1
    gs = tc_pr.find("./w:gridSpan", DOCX_NS)
    if gs is None:
        return 1
    val = gs.get(f"{DOCX_TAG_PREFIX}val") or "1"
    try:
        return max(int(val), 1)
    except ValueError:
        return 1


def _docx_cell_vmerge(cell: ET.Element) -> str:
    tc_pr = cell.find("./w:tcPr", DOCX_NS)
    if tc_pr is None:
        return ""
    vm = tc_pr.find("./w:vMerge", DOCX_NS)
    if vm is None:
        return ""
    val = vm.get(f"{DOCX_TAG_PREFIX}val")
    return "restart" if val == "restart" else "continue"


# 将 docx 表格转成 Markdown，便于后续参数抽取与人工复核。
def render_docx_table_as_markdown(table: ET.Element) -> str:
    rows_raw: list[list[tuple[str, str]]] = []
    for row in table.findall("./w:tr", DOCX_NS):
        row_cells: list[tuple[str, str]] = []
        for cell in row.findall("./w:tc", DOCX_NS):
            cell_lines: list[str] = []
            for paragraph in cell.findall("./w:p", DOCX_NS):
                paragraph_text = extract_docx_text_from_node(paragraph)
                if paragraph_text:
                    cell_lines.append(paragraph_text)
            cell_text = merge_docx_table_cell_lines(cell_lines)
            span = _docx_cell_grid_span(cell)
            vmerge = _docx_cell_vmerge(cell)
            # 横向合并：首列放文本，剩余列用空字符串占位，保证与未合并行列数一致
            row_cells.append((cell_text, vmerge))
            for _ in range(span - 1):
                row_cells.append(("", ""))
        rows_raw.append(row_cells)

    # 纵向合并：连续行如果是 continue 状态，文本继承自最近的 restart；为保持 markdown 可读性，
    # 直接留空占位，不回填文本（否则会出现相邻行文字重复）。
    rows: list[list[str]] = []
    for row in rows_raw:
        resolved = []
        for text, vmerge in row:
            resolved.append("" if vmerge == "continue" else text)
        rows.append(resolved)

    # 整行全空或仅为占位符 "-" 的装饰行过滤掉，避免形成"空 header + 分隔符"的破损表格
    rows = [r for r in rows if any(cell.strip() and cell.strip() != "-" for cell in r)]
    if not rows:
        return ""
    column_count = max(len(row) for row in rows)
    # 短于最长行的用空字符串补齐（历史用 "-" 会和 markdown 分隔符视觉冲突）
    normalized_rows = [row + [""] * (column_count - len(row)) for row in rows]

    header_raw = normalized_rows[0]
    header_is_empty = not any(cell.strip() and cell.strip() != "-" for cell in header_raw)
    if header_is_empty and len(normalized_rows) > 1 and any(cell.strip() and cell.strip() != "-" for cell in normalized_rows[1]):
        # 首行装饰性空行：使用第 2 行作为 header，后续行作为 body
        header = [cell if cell.strip() and cell.strip() != "-" else f"列{idx + 1}" for idx, cell in enumerate(normalized_rows[1])]
        body_rows = normalized_rows[2:]
    elif header_is_empty:
        # 所有行都没有明确 header，用列序号占位
        header = [f"列{idx + 1}" for idx in range(column_count)]
        body_rows = normalized_rows
    else:
        header = [cell if cell.strip() and cell.strip() != "-" else f"列{idx + 1}" for idx, cell in enumerate(header_raw)]
        body_rows = normalized_rows[1:]

    separator = ["---"] * column_count
    markdown_lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for row in body_rows:
        # 空 cell 用单个空格，避免 markdown 把 "||" 压成一条线
        safe_row = [cell if cell.strip() else " " for cell in row]
        markdown_lines.append("| " + " | ".join(safe_row) + " |")
    return "\n".join(markdown_lines)


def normalize_excel_cell(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, datetime):
        return value.strftime("%Y-%m-%d %H:%M:%S")
    if isinstance(value, date):
        return value.strftime("%Y-%m-%d")
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value).strip()


def trim_excel_row(values: tuple[Any, ...]) -> list[str]:
    cells = [normalize_excel_cell(value) for value in values]
    while cells and not cells[-1]:
        cells.pop()
    return cells


def render_excel_sheet_as_markdown(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    width_counter = Counter(len(row) for row in rows if row)
    column_count = max(width_counter, key=lambda key: (width_counter[key], key)) if width_counter else max(len(row) for row in rows)
    normalized_rows = [row + [""] * (column_count - len(row)) for row in rows]
    if normalized_rows and sum(1 for row in normalized_rows if row and not row[0]) / len(normalized_rows) >= 0.8:
        normalized_rows = [row[1:] for row in normalized_rows if len(row) > 1]
        column_count = max(len(row) for row in normalized_rows)
        normalized_rows = [row + [""] * (column_count - len(row)) for row in normalized_rows]
    header = [cell or f"列{index + 1}" for index, cell in enumerate(normalized_rows[0])]
    body_rows = normalized_rows[1:] if len(normalized_rows) > 1 else []
    return render_markdown_table(header, body_rows)


def extract_excel_text(source_path: Path) -> ExtractionResult:
    if source_path.suffix.lower() == ".xls":
        return extract_legacy_excel_text(source_path)

    try:
        from openpyxl import load_workbook
    except ImportError:
        return ExtractionResult(
            extractor_name="excel:none",
            extraction_status="待人工复核",
            extracted_text="",
            preview_text="",
            text_length=0,
            page_count=None,
            source_encoding="",
            note="openpyxl 不可用，当前无法读取 Excel。",
        )

    # read_only=True 下无法读取合并单元格信息；先用标准模式读一次 merged_cells 信息，再切回快速迭代。
    merged_ranges_per_sheet: dict[str, list[tuple[int, int, int, int]]] = {}
    try:
        meta_wb = load_workbook(source_path, read_only=False, data_only=True)
        for ws in meta_wb.worksheets:
            ranges = []
            for merged_range in list(ws.merged_cells.ranges):
                ranges.append((merged_range.min_row, merged_range.min_col, merged_range.max_row, merged_range.max_col))
            merged_ranges_per_sheet[ws.title] = ranges
        meta_wb.close()
    except Exception:
        merged_ranges_per_sheet = {}

    try:
        workbook = load_workbook(source_path, read_only=True, data_only=True)
    except Exception as exc:
        return ExtractionResult(
            extractor_name="excel:none",
            extraction_status="待人工复核",
            extracted_text="",
            preview_text="",
            text_length=0,
            page_count=None,
            source_encoding="",
            note=f"Excel 读取失败：{exc}",
        )

    try:
        blocks: list[str] = []
        effective_sheet_count = 0
        for worksheet in workbook.worksheets:
            raw_rows: list[list[str]] = [[normalize_excel_cell(v) for v in row] for row in worksheet.iter_rows(values_only=True)]
            # 合并单元格在 values_only 遍历里只有左上角保留了值，其余位置为空，这里把合并范围内的空位按左上角值补齐。
            merged_ranges = merged_ranges_per_sheet.get(worksheet.title, [])
            for (min_row, min_col, max_row, max_col) in merged_ranges:
                top_row_idx = min_row - 1
                top_col_idx = min_col - 1
                if top_row_idx < 0 or top_row_idx >= len(raw_rows):
                    continue
                origin_row = raw_rows[top_row_idx]
                if top_col_idx < 0 or top_col_idx >= len(origin_row):
                    continue
                origin_value = origin_row[top_col_idx]
                if not origin_value:
                    continue
                for rr in range(min_row - 1, max_row):
                    if rr >= len(raw_rows):
                        break
                    current_row = raw_rows[rr]
                    for cc in range(min_col - 1, max_col):
                        while cc >= len(current_row):
                            current_row.append("")
                        if rr == min_row - 1 and cc == min_col - 1:
                            continue
                        if not current_row[cc]:
                            current_row[cc] = origin_value
            rows = [trim_excel_row(tuple(row)) for row in raw_rows]
            rows = [row for row in rows if any(cell for cell in row)]
            if not rows:
                continue
            effective_sheet_count += 1
            blocks.append(f"# 工作表：{worksheet.title}")
            blocks.append(render_excel_sheet_as_markdown(rows))

        if not blocks:
            return ExtractionResult(
                extractor_name="excel:openpyxl",
                extraction_status="待人工复核",
                extracted_text="",
                preview_text="",
                text_length=0,
                page_count=0,
                source_encoding="xlsx",
                note="Excel 中未识别到有效内容。",
            )

        extracted_text = "\n\n".join(blocks)
        return ExtractionResult(
            extractor_name="excel:openpyxl",
            extraction_status="已提取文本",
            extracted_text=extracted_text,
            preview_text=normalize_preview(extracted_text),
            text_length=len(extracted_text),
            page_count=effective_sheet_count,
            source_encoding="xlsx",
            note=f"已完成 Excel 读取，共识别 {effective_sheet_count} 个有效工作表。",
        )
    finally:
        workbook.close()


# 旧版 xls 先通过 Windows Excel COM 保守导出为 xlsx，再复用现有 xlsx 提取逻辑；失败时保留待复核状态。
def extract_legacy_excel_text(source_path: Path) -> ExtractionResult:
    if pythoncom is None or win32_client is None:
        return ExtractionResult(
            extractor_name="excel:none",
            extraction_status="待人工复核",
            extracted_text="",
            preview_text="",
            text_length=0,
            page_count=None,
            source_encoding="",
            note="当前环境缺少 Windows Excel COM 依赖，无法读取 xls。",
        )

    temp_dir = tempfile.TemporaryDirectory()
    temp_xlsx_path = Path(temp_dir.name) / f"{source_path.stem}.xlsx"
    pythoncom.CoInitialize()
    excel = None
    workbook = None
    try:
        excel = win32_client.DispatchEx("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False
        workbook = excel.Workbooks.Open(str(source_path.resolve()))
        workbook.SaveAs(str(temp_xlsx_path), FileFormat=51)
    except Exception as exc:
        return ExtractionResult(
            extractor_name="excel:none",
            extraction_status="待人工复核",
            extracted_text="",
            preview_text="",
            text_length=0,
            page_count=None,
            source_encoding="",
            note=f"xls 转 xlsx 失败：{exc}",
        )
    finally:
        if workbook is not None:
            workbook.Close(False)
        if excel is not None:
            excel.Quit()
        pythoncom.CoUninitialize()

    try:
        result = extract_excel_text(temp_xlsx_path)
    finally:
        temp_dir.cleanup()
    if result.extraction_status == "已提取文本":
        return ExtractionResult(
            extractor_name="excel:com->openpyxl",
            extraction_status=result.extraction_status,
            extracted_text=result.extracted_text,
            preview_text=result.preview_text,
            text_length=result.text_length,
            page_count=result.page_count,
            source_encoding="xls",
            note="已通过 Excel COM 转存为 xlsx 后完成提取。",
            extra_metadata=result.extra_metadata,
        )
    return ExtractionResult(
        extractor_name="excel:com->openpyxl",
        extraction_status=result.extraction_status,
        extracted_text=result.extracted_text,
        preview_text=result.preview_text,
        text_length=result.text_length,
        page_count=result.page_count,
        source_encoding="xls",
        note=f"已通过 Excel COM 转存为 xlsx，但结果仍需复核：{result.note}",
        extra_metadata=result.extra_metadata,
    )


# 超大 PPT/PPTX 无法直接上传 MinerU 时，先用 PowerPoint COM 导出逐页图片，再复用 MinerU OCR 主链路。
def export_presentation_slides_to_images(source_path: Path, image_format: str = "PNG") -> tuple[tempfile.TemporaryDirectory, list[Path]]:
    if pythoncom is None or win32_client is None:
        raise RuntimeError("当前环境缺少 Windows PowerPoint COM 依赖，无法导出幻灯片图片。")

    temp_dir = tempfile.TemporaryDirectory()
    output_dir = Path(temp_dir.name)
    pythoncom.CoInitialize()
    powerpoint = None
    presentation = None
    try:
        powerpoint = win32_client.DispatchEx("PowerPoint.Application")
        # PowerPoint COM 不允许像 Excel 一样隐藏主窗口，强行设为 0 会直接报错。
        powerpoint.Visible = 1
        powerpoint.DisplayAlerts = 0
        presentation = powerpoint.Presentations.Open(str(source_path.resolve()), False, False, False)
        image_paths: list[Path] = []
        slide_count = int(presentation.Slides.Count)
        for index in range(1, slide_count + 1):
            output_path = output_dir / f"{index:04d}.{image_format.lower()}"
            presentation.Slides.Item(index).Export(str(output_path), image_format)
            image_paths.append(output_path)

        if not image_paths:
            raise RuntimeError("PowerPoint 未导出任何幻灯片图片。")
        return temp_dir, image_paths
    except Exception as exc:
        temp_dir.cleanup()
        raise RuntimeError(f"PPTX 导出分页图片失败：{exc}") from exc
    finally:
        if presentation is not None:
            presentation.Close()
        if powerpoint is not None:
            powerpoint.Quit()
        pythoncom.CoUninitialize()


# pptx 先直接读取文本框与表格内容，保持轻量保守，不尝试解析图片或复杂版式。
def extract_presentation_text(source_path: Path) -> ExtractionResult:
    try:
        from pptx import Presentation
    except ImportError:
        return ExtractionResult(
            extractor_name="presentation:none",
            extraction_status="待人工复核",
            extracted_text="",
            preview_text="",
            text_length=0,
            page_count=None,
            source_encoding="",
            note="python-pptx 不可用，当前无法读取演示文稿。",
        )

    try:
        presentation = Presentation(str(source_path))
    except Exception as exc:
        return ExtractionResult(
            extractor_name="presentation:none",
            extraction_status="待人工复核",
            extracted_text="",
            preview_text="",
            text_length=0,
            page_count=None,
            source_encoding="",
            note=f"演示文稿读取失败：{exc}",
        )

    blocks: list[str] = []
    effective_slide_count = 0
    for index, slide in enumerate(presentation.slides, start=1):
        slide_blocks: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text.strip())
                if text:
                    slide_blocks.append(text)
                continue
            if getattr(shape, "has_table", False):
                rows: list[list[str]] = []
                for row in shape.table.rows:
                    values: list[str] = []
                    for cell in row.cells:
                        # 合并单元格的"延续"部分直接置空，保证每行列数一致；文本由 origin 单元格承担。
                        if getattr(cell, "is_spanned", False):
                            values.append("")
                            continue
                        text = re.sub(r"\s+", " ", str(cell.text)).strip()
                        values.append(text)
                    if any(values):
                        rows.append(values)
                if rows:
                    slide_blocks.append(render_markdown_table([cell or f"列{idx + 1}" for idx, cell in enumerate(rows[0])], rows[1:]))
        if not slide_blocks:
            continue
        effective_slide_count += 1
        blocks.append(f"# 第{index}页")
        blocks.extend(slide_blocks)

    if not blocks:
        return ExtractionResult(
            extractor_name="presentation:pptx",
            extraction_status="待人工复核",
            extracted_text="",
            preview_text="",
            text_length=0,
            page_count=len(presentation.slides),
            source_encoding="pptx",
            note="演示文稿中未识别到有效文本或表格内容。",
        )

    extracted_text = "\n\n".join(blocks)
    return ExtractionResult(
        extractor_name="presentation:pptx",
        extraction_status="已提取文本",
        extracted_text=extracted_text,
        preview_text=normalize_preview(extracted_text),
        text_length=len(extracted_text),
        page_count=effective_slide_count,
        source_encoding="pptx",
        note=f"已完成演示文稿读取，共识别 {effective_slide_count} 页含文本内容的幻灯片。",
    )


def is_document_like_image(source_path: Path) -> bool:
    try:
        from PIL import Image
    except ImportError:
        return False

    try:
        with Image.open(source_path) as image:
            grayscale = image.convert("L")
            width, height = grayscale.size
            if width < 80 or height < 80:
                return False
            sample = grayscale.resize((min(256, width), min(256, height)))
            pixels = list(sample.getdata())
            total = len(pixels)
            if not total:
                return False
            white_ratio = sum(pixel >= 235 for pixel in pixels) / total
            dark_ratio = sum(pixel <= 90 for pixel in pixels) / total

            transition_count = 0
            horizontal_checks = 0
            for y in range(sample.height):
                for x in range(sample.width - 1):
                    current_pixel = sample.getpixel((x, y))
                    next_pixel = sample.getpixel((x + 1, y))
                    horizontal_checks += 1
                    if abs(current_pixel - next_pixel) >= 40:
                        transition_count += 1
            vertical_checks = 0
            for x in range(sample.width):
                for y in range(sample.height - 1):
                    current_pixel = sample.getpixel((x, y))
                    next_pixel = sample.getpixel((x, y + 1))
                    vertical_checks += 1
                    if abs(current_pixel - next_pixel) >= 40:
                        transition_count += 1
            transition_ratio = transition_count / max(horizontal_checks + vertical_checks, 1)
            return white_ratio >= 0.45 and dark_ratio >= 0.02 and transition_ratio >= 0.06
    except Exception as exc:
        logging.debug("is_document_like_image判定异常: %s", exc)
        return False


def should_skip_image_file(source_path: Path) -> bool:
    path_text = str(source_path).lower()
    stem = source_path.stem.lower()
    if any(keyword in path_text for keyword in DOCUMENT_IMAGE_KEYWORDS):
        return False
    if is_document_like_image(source_path):
        return False
    if any(keyword in path_text for keyword in WECHAT_IMAGE_KEYWORDS):
        return False
    if re.match(r"^\d{1,4}", stem):
        return False
    if any(keyword in stem for keyword in PAGED_IMAGE_NAME_KEYWORDS):
        return False
    return True



def has_document_like_ocr_text(text: str) -> bool:
    compact = re.sub(r"\s+", "", text or "")
    if has_meaningful_text(compact):
        return True
    return any(keyword in compact for keyword in DOCUMENT_IMAGE_KEYWORDS)



def finalize_image_result_after_ocr(source_path: Path, ocr_result: ExtractionResult) -> ExtractionResult:
    if has_document_like_ocr_text(ocr_result.extracted_text):
        return ocr_result
    if should_skip_image_file(source_path):
        return skip_image_as_photo(
            source_path,
            extra_note="已启用 OCR，但识别文本较少且未命中文档关键词，按纯照片跳过。",
        )
    return ocr_result



def skip_image_as_photo(source_path: Path, extra_note: str = "") -> ExtractionResult:
    # 纯照片默认直接跳过，不进入 OCR 或结构化正文提取。
    note = f"图片文件{source_path.name}当前按纯照片处理，直接跳过，不进入 OCR。"
    if extra_note:
        note = f"{note} {extra_note}"
    return ExtractionResult(
        extractor_name="skip:image_photo",
        extraction_status="跳过",
        extracted_text="",
        preview_text="",
        text_length=0,
        page_count=1,
        source_encoding="",
        note=note,
    )



def extract_image_with_ocr_placeholder(source_path: Path, extra_note: str = "") -> ExtractionResult:
    # 第三阶段先把文档型图片收口到 OCR 占位接口，后续直接替换成真实 OCR SDK。
    notes = [f"图片文件 {source_path.name} 当前进入 OCR 占位层。"]
    if extra_note:
        notes.append(extra_note)
    else:
        notes.append("当前未启用真实 OCR（缺少 MinerU token 或未开启 OCR 开关），仅保留待 OCR 占位结果。")
    return ExtractionResult(
        extractor_name="ocr:placeholder:image",
        extraction_status="待OCR",
        extracted_text="",
        preview_text="",
        text_length=0,
        page_count=1,
        source_encoding="",
        note="；".join(notes),
    )


def build_ocr_result(source_path: Path, text: str, *, page_count: int | None, note: str, extractor_name: str) -> ExtractionResult:
    cleaned_text = clean_mineru_markdown(text)
    preview = normalize_preview(cleaned_text)
    if has_meaningful_text(cleaned_text):
        return ExtractionResult(
            extractor_name=extractor_name,
            extraction_status="已提取文本",
            extracted_text=cleaned_text,
            preview_text=preview,
            text_length=len(cleaned_text),
            page_count=page_count,
            source_encoding="utf-8",
            note=note,
        )
    if cleaned_text.strip():
        return ExtractionResult(
            extractor_name=extractor_name,
            extraction_status="待审核",
            extracted_text=cleaned_text,
            preview_text=preview,
            text_length=len(cleaned_text),
            page_count=page_count,
            source_encoding="utf-8",
            note=note,
        )
    raise RuntimeError(f"{source_path.name} OCR 未返回可用文本。")


def _size_aware_max_polls(source_path: Path, poll_interval_seconds: int = 5) -> int:
    try:
        size_mb = source_path.stat().st_size / (1024 * 1024)
    except OSError:
        size_mb = 0.0
    base_polls = 60
    extra_polls = int((size_mb / 5.0) * 12)
    return min(max(base_polls, base_polls + extra_polls), 900)


def extract_image_with_mineru_ocr(source_path: Path, token: str) -> ExtractionResult:
    batch_result = run_mineru_batch([source_path], token, max_polls=_size_aware_max_polls(source_path))
    result = batch_result["results"][0]
    markdown = result.get("markdown", "")
    if result.get("state") != "done":
        raise RuntimeError(result.get("error") or f"OCR 状态异常：{result.get('state')}")
    return build_ocr_result(
        source_path,
        markdown,
        page_count=1,
        note=f"已通过 MinerU OCR 完成图片文本提取，批次号：{batch_result['batch_id']}。",
        extractor_name="ocr:mineru:image",
    )


def extract_pdf_with_mineru_ocr(source_path: Path, token: str, *, page_count: int | None) -> ExtractionResult:
    batch_result = run_mineru_batch([source_path], token, max_polls=_size_aware_max_polls(source_path))
    result = batch_result["results"][0]
    markdown = result.get("markdown", "")
    if result.get("state") != "done":
        raise RuntimeError(result.get("error") or f"OCR 状态异常：{result.get('state')}")
    return build_ocr_result(
        source_path,
        markdown,
        page_count=page_count,
        note=f"已通过 MinerU OCR 完成扫描型 PDF 文本提取，批次号：{batch_result['batch_id']}。",
        extractor_name="ocr:mineru:pdf",
    )


def extract_pdf_text(source_path: Path, *, enable_ocr: bool = False, ocr_token: str | None = None) -> ExtractionResult:
    pypdf_text, pypdf_pages, pypdf_note = extract_pdf_with_pypdf(source_path)
    plumber_text, plumber_pages, plumber_note = extract_pdf_with_pdfplumber(source_path)

    candidates = [
        ("pypdf", pypdf_text, pypdf_pages, pypdf_note),
        ("pdfplumber", plumber_text, plumber_pages, plumber_note),
    ]
    extractor_name, best_text, page_count, best_note = max(candidates, key=lambda item: len(item[1]))
    preview = normalize_preview(best_text)

    # 扫描/混合型 PDF 常见模式：只有封面/少量页有文字层，其余页是图像。
    # 只要用户显式开了 OCR，就按"每页平均字数"判断是否疏稀，疏稀时继续走 OCR 兜底，并把 OCR 结果与文字层拼接，保证扫描页内容不丢失。
    effective_pages = page_count or pypdf_pages or plumber_pages or 0
    avg_chars_per_page = (len(best_text) / effective_pages) if effective_pages else len(best_text)
    sparse_text = enable_ocr and ocr_token and effective_pages > 0 and avg_chars_per_page < 120

    if has_meaningful_text(best_text) and not sparse_text:
        return ExtractionResult(
            extractor_name=extractor_name,
            extraction_status="已提取文本",
            extracted_text=best_text,
            preview_text=preview,
            text_length=len(best_text),
            page_count=page_count,
            source_encoding="",
            note=best_note or f"已通过 {extractor_name} 完成 PDF 文本提取。",
        )

    if sparse_text:
        try:
            ocr_result = extract_pdf_with_mineru_ocr(source_path, ocr_token, page_count=effective_pages)
        except Exception as exc:
            if has_meaningful_text(best_text):
                return ExtractionResult(
                    extractor_name=extractor_name,
                    extraction_status="已提取文本",
                    extracted_text=best_text,
                    preview_text=preview,
                    text_length=len(best_text),
                    page_count=page_count,
                    source_encoding="",
                    note=f"{best_note or ''} PDF 文字层较稀疏，已尝试 OCR 但失败：{exc}".strip(),
                )
        else:
            # 文字层本来就是空的 → 直接返回 OCR 结果；否则把文字层与 OCR 拼接。
            if not best_text.strip():
                return ocr_result
            merged_text = best_text.strip() + "\n\n" + ocr_result.extracted_text
            return ExtractionResult(
                extractor_name=f"{extractor_name}+{ocr_result.extractor_name}",
                extraction_status="已提取文本",
                extracted_text=merged_text,
                preview_text=normalize_preview(merged_text),
                text_length=len(merged_text),
                page_count=effective_pages,
                source_encoding="utf-8",
                note=f"PDF 文字层稀疏（平均 {avg_chars_per_page:.0f} 字/页），已通过 OCR 补齐扫描页。{ocr_result.note}",
            )

    if best_text.strip():
        return ExtractionResult(
            extractor_name=extractor_name,
            extraction_status="待审核",
            extracted_text=best_text,
            preview_text=preview,
            text_length=len(best_text),
            page_count=page_count,
            source_encoding="",
            note="PDF 仅提取到少量文本，当前不足以直接进入稳定抽取链路。",
        )

    notes = [item for item in [pypdf_note, plumber_note] if item]
    notes.append("当前 PDF 未提取到可用文本，疑似扫描件或图片型 PDF，建议后续接 OCR。")
    if enable_ocr and ocr_token:
        try:
            return extract_pdf_with_mineru_ocr(source_path, ocr_token, page_count=pypdf_pages or plumber_pages)
        except Exception as exc:
            notes.append(f"已尝试真实 OCR，但调用失败：{exc}")
    else:
        notes.append("当前未启用真实 OCR（缺少 MinerU token 或未开启 OCR 开关），仅保留待 OCR 占位结果。")
    return ExtractionResult(
        extractor_name="ocr:placeholder:pdf",
        extraction_status="待OCR",
        extracted_text="",
        preview_text="",
        text_length=0,
        page_count=pypdf_pages or plumber_pages,
        source_encoding="",
        note="；".join(notes),
    )


def extract_pdf_with_pypdf(source_path: Path) -> tuple[str, int | None, str]:
    try:
        from pypdf import PdfReader
    except ImportError:
        return "", None, "pypdf 不可用。"

    try:
        reader = PdfReader(str(source_path))
        texts: list[str] = []
        for page in reader.pages:
            page_text = (page.extract_text() or "").strip()
            if page_text:
                texts.append(page_text)
        return "\n\n".join(texts), len(reader.pages), "已尝试使用 pypdf 提取。"
    except Exception as exc:
        return "", None, f"pypdf 提取失败：{exc}"


def extract_pdf_with_pdfplumber(source_path: Path) -> tuple[str, int | None, str]:
    try:
        import pdfplumber
    except ImportError:
        return "", None, "pdfplumber 不可用。"

    try:
        texts: list[str] = []
        with pdfplumber.open(str(source_path)) as pdf:
            for page in pdf.pages:
                page_text = (page.extract_text() or "").strip()
                if page_text:
                    texts.append(page_text)
            return "\n\n".join(texts), len(pdf.pages), "已尝试使用 pdfplumber 提取。"
    except Exception as exc:
        return "", None, f"pdfplumber 提取失败：{exc}"


def get_retry_after_seconds(response: requests.Response, default_seconds: int = MINERU_RATE_LIMIT_RETRY_DELAY_SECONDS) -> int:
    retry_after = response.headers.get("Retry-After", "").strip()
    if retry_after.isdigit():
        return max(int(retry_after), 1)
    return default_seconds


def _heartbeat(message: str) -> None:
    print(f"[MinerU] {message}", flush=True)


# 对 MinerU 单次 HTTP 请求做统一重试封装：支持网络异常、429 限流、5xx 服务端错误指数退避重试。
def mineru_request_with_retry(
    label: str,
    request_fn,
    *,
    max_attempts: int = MINERU_HTTP_MAX_RETRIES,
) -> requests.Response:
    delay = MINERU_HTTP_BACKOFF_SECONDS
    last_exc: Exception | None = None
    last_response: requests.Response | None = None

    for attempt in range(1, max_attempts + 1):
        try:
            response = request_fn()
        except requests.RequestException as exc:
            last_exc = exc
            last_response = None
            if attempt >= max_attempts:
                break
            _heartbeat(f"{label} 第 {attempt}/{max_attempts} 次失败：{exc}，{int(delay)}s 后重试")
            time.sleep(delay)
            delay *= 2
            continue

        if response.status_code == 429:
            last_response = response
            last_exc = None
            if attempt >= max_attempts:
                break
            wait = get_retry_after_seconds(response)
            _heartbeat(f"{label} 被限流 (429)，{wait}s 后重试 (第 {attempt}/{max_attempts})")
            time.sleep(wait)
            continue

        if 500 <= response.status_code < 600:
            last_response = response
            last_exc = None
            if attempt >= max_attempts:
                break
            _heartbeat(f"{label} 第 {attempt}/{max_attempts} 次失败：HTTP {response.status_code}，{int(delay)}s 后重试")
            time.sleep(delay)
            delay *= 2
            continue

        return response

    if last_exc is not None:
        raise last_exc
    assert last_response is not None
    last_response.raise_for_status()
    return last_response


# MinerU 限流时按响应头或默认间隔短暂重试，避免大批量图片链路被一次 429 直接打断。
def post_mineru_batch_upload_urls(
    http: requests.Session,
    api_base: str,
    headers: dict[str, str],
    payload: dict[str, Any],
) -> requests.Response:
    _heartbeat("请求批量上传链接...")
    return mineru_request_with_retry(
        "请求上传链接",
        lambda: http.post(
            f"{api_base}/file-urls/batch",
            headers=headers,
            json=payload,
            timeout=MINERU_REQUEST_TIMEOUT,
        ),
        max_attempts=MINERU_RATE_LIMIT_MAX_RETRIES + 1,
    )


# MinerU 先以独立批量接口接入，避免直接破坏当前最小闭环主链；后续样例跑稳后再切换默认主处理层。
def request_mineru_batch_upload_urls(
    token: str,
    files: list[Path],
    *,
    api_base: str = MINERU_API_BASE,
    model_version: str = "vlm",
    enable_ocr: bool = True,
    enable_formula: bool = True,
    enable_table: bool = True,
    session: requests.Session | None = None,
) -> tuple[str, list[dict[str, str]]]:
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {token}",
    }
    payload = {
        "files": [
            {"name": file_path.name, "data_id": f"file_{index}_{int(time.time())}"}
            for index, file_path in enumerate(files)
        ],
        "model_version": model_version,
        "is_ocr": enable_ocr,
        "enable_formula": enable_formula,
        "enable_table": enable_table,
    }
    http = session or create_mineru_session()
    response = post_mineru_batch_upload_urls(http, api_base, headers, payload)
    response.raise_for_status()
    result = response.json()
    if result.get("code") != 0:
        raise RuntimeError(f"MinerU 获取上传链接失败：{result.get('msg', '未知错误')}")
    data = result.get("data") or {}
    batch_id = data.get("batch_id")
    raw_urls = data.get("file_urls") or []
    normalized_urls: list[dict[str, str]] = []
    for item in raw_urls:
        if isinstance(item, str):
            normalized_urls.append({"upload_url": item})
        elif isinstance(item, dict) and item.get("upload_url"):
            normalized_urls.append({"upload_url": item["upload_url"]})
    if not batch_id or len(normalized_urls) != len(files):
        raise RuntimeError("MinerU 返回的批次或上传链接数量异常。")
    return batch_id, normalized_urls


def upload_mineru_file(upload_url: str, file_path: Path, *, session: requests.Session | None = None) -> None:
    http = session or create_mineru_session()
    _heartbeat(f"上传文件：{file_path.name}")

    def _do_put() -> requests.Response:
        with file_path.open("rb") as file_obj:
            return http.put(upload_url, data=file_obj, timeout=MINERU_UPLOAD_TIMEOUT)

    response = mineru_request_with_retry(f"上传 {file_path.name}", _do_put)
    if response.status_code not in {200, 201, 204}:
        raise RuntimeError(f"MinerU 文件上传失败：HTTP {response.status_code}")


def poll_mineru_batch_results(
    token: str,
    batch_id: str,
    *,
    api_base: str = MINERU_API_BASE,
    poll_interval_seconds: int = 5,
    max_polls: int = 60,
    session: requests.Session | None = None,
) -> list[dict[str, Any]]:
    headers = {"Authorization": f"Bearer {token}"}
    http = session or create_mineru_session()
    for attempt in range(1, max_polls + 1):
        _heartbeat(f"轮询结果 {attempt}/{max_polls}（批次 {batch_id}）")
        response = mineru_request_with_retry(
            f"轮询批次 {batch_id}",
            lambda: http.get(
                f"{api_base}/extract-results/batch/{batch_id}",
                headers=headers,
                timeout=MINERU_REQUEST_TIMEOUT,
            ),
        )
        response.raise_for_status()
        result = response.json()
        if result.get("code") != 0:
            raise RuntimeError(f"MinerU 查询结果失败：{result.get('msg', '未知错误')}")
        tasks = (result.get("data") or {}).get("extract_result") or []
        if tasks and all(task.get("state") in {"done", "failed"} for task in tasks):
            _heartbeat(f"轮询完成，任务已就绪（批次 {batch_id}）")
            return tasks
        time.sleep(poll_interval_seconds)
    raise TimeoutError(f"MinerU 批次 {batch_id} 在轮询次数内未完成。")


def _html_table_to_markdown(table_html: str) -> str:
    """
    将 MinerU 返回的单个 <table>...</table> 块转成 markdown pipe 表格。
    必须在剥离 <table>/<tr>/<td> 标签**之前**调用，否则表格结构会被压成一行文本。
    """
    row_matches = re.findall(r"<tr[^>]*>(.*?)</tr>", table_html, flags=re.IGNORECASE | re.DOTALL)
    rows: list[list[str]] = []
    for row_html in row_matches:
        cell_matches = re.findall(r"<t[dh][^>]*>(.*?)</t[dh]>", row_html, flags=re.IGNORECASE | re.DOTALL)
        cleaned_cells: list[str] = []
        for cell in cell_matches:
            text = re.sub(r"<br\s*/?>", " ", cell, flags=re.IGNORECASE)
            text = re.sub(r"<[^>]+>", "", text)
            text = re.sub(r"\s+", " ", text).strip()
            text = text.replace("|", r"\|")
            cleaned_cells.append(text if text else " ")
        if cleaned_cells:
            rows.append(cleaned_cells)

    if not rows:
        return ""

    column_count = max(len(row) for row in rows)
    normalized = [row + [" "] * (column_count - len(row)) for row in rows]
    header = normalized[0]
    body = normalized[1:] if len(normalized) > 1 else []

    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * column_count) + " |",
    ]
    for row in body:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def clean_mineru_markdown(markdown: str) -> str:
    # MinerU 对扫描件里的检测表格通常返回 HTML <table>；必须先把 HTML 表格转换成 markdown pipe 表格，
    # 否则下一步剥离 <table>/<tr>/<td> 会把表格压成一长串无分隔文本（检测报告场景高频踩坑点）。
    def _table_to_md(match: re.Match) -> str:
        converted = _html_table_to_markdown(match.group(0))
        return "\n\n" + converted + "\n\n" if converted else ""

    cleaned = re.sub(
        r"<table[^>]*>.*?</table>",
        _table_to_md,
        markdown,
        flags=re.IGNORECASE | re.DOTALL,
    )
    cleaned = re.sub(r"\[!\[.*?\]\(.*?\)\]\(.*?\)", "", cleaned, flags=re.DOTALL)
    cleaned = re.sub(r"!\[.*?\]\(.*?\)", "", cleaned, flags=re.DOTALL)
    cleaned = re.sub(r"<img[^>]*>", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"</?(?:table|tr|td|th|tbody|thead|p|div|span|b|strong|i|em|u)[^>]*>", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"</?br[^>]*>", "\n", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"<!--.*?-->", "", cleaned, flags=re.DOTALL)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def download_mineru_markdown(zip_url: str, *, session: requests.Session | None = None) -> str:
    http = session or create_mineru_session()
    _heartbeat("下载结果 ZIP 并解压 Markdown...")
    response = mineru_request_with_retry(
        "下载结果 ZIP",
        lambda: http.get(zip_url, timeout=MINERU_DOWNLOAD_TIMEOUT),
    )
    response.raise_for_status()
    zip_buffer = io.BytesIO(response.content)
    with zipfile.ZipFile(zip_buffer, "r") as zip_file:
        markdown_candidates = [name for name in zip_file.namelist() if name.endswith(".md")]
        if not markdown_candidates:
            raise RuntimeError("MinerU 结果包中未找到 Markdown 文件。")
        preferred_name = next((name for name in markdown_candidates if "full" in name.lower()), markdown_candidates[0])
        with zip_file.open(preferred_name) as file_obj:
            raw_markdown = file_obj.read().decode("utf-8")
    return clean_mineru_markdown(raw_markdown)


def run_mineru_batch(
    files: list[Path],
    token: str,
    *,
    api_base: str = MINERU_API_BASE,
    model_version: str = "vlm",
    enable_ocr: bool = True,
    enable_formula: bool = True,
    enable_table: bool = True,
    poll_interval_seconds: int = 5,
    max_polls: int = 60,
) -> dict[str, Any]:
    http = create_mineru_session()
    batch_id, upload_urls = request_mineru_batch_upload_urls(
        token,
        files,
        api_base=api_base,
        model_version=model_version,
        enable_ocr=enable_ocr,
        enable_formula=enable_formula,
        enable_table=enable_table,
        session=http,
    )
    for file_path, upload_info in zip(files, upload_urls):
        upload_mineru_file(upload_info["upload_url"], file_path, session=http)
    tasks = poll_mineru_batch_results(
        token,
        batch_id,
        api_base=api_base,
        poll_interval_seconds=poll_interval_seconds,
        max_polls=max_polls,
        session=http,
    )
    results: list[dict[str, Any]] = []
    for file_path, task in zip(files, tasks):
        state = task.get("state", "unknown")
        zip_url = task.get("full_zip_url", "")
        markdown = ""
        if state == "done" and zip_url:
            markdown = download_mineru_markdown(zip_url, session=http)
        results.append(
            {
                "file_path": str(file_path),
                "state": state,
                "markdown": markdown,
                "zip_url": zip_url,
                "error": task.get("err_msg", ""),
            }
        )
    return {"batch_id": batch_id, "results": results}


def extract_pdf_with_mineru(
    source_path: Path,
    token: str,
    *,
    api_base: str = MINERU_API_BASE,
    model_version: str = "vlm",
    enable_ocr: bool = True,
    enable_formula: bool = True,
    enable_table: bool = True,
    poll_interval_seconds: int = 5,
    max_polls: int | None = None,
) -> ExtractionResult:
    # MinerU 轮询预算按文件大小线性放大：每 5 MB 预留约 1 分钟，叠加 5 分钟基线，上限 75 分钟。
    # 避免上百兆检测报告 docx 被 60 次 (5 min) 固定上限强制超时。
    if max_polls is None:
        max_polls = _size_aware_max_polls(source_path, poll_interval_seconds)
        try:
            size_mb = source_path.stat().st_size / (1024 * 1024)
        except OSError:
            size_mb = 0.0
        _heartbeat(
            f"根据文件大小 {size_mb:.1f} MB 预留 {max_polls} 次轮询（约 {max_polls * poll_interval_seconds // 60} 分钟）"
        )
    batch_result = run_mineru_batch(
        [source_path],
        token,
        api_base=api_base,
        model_version=model_version,
        enable_ocr=enable_ocr,
        enable_formula=enable_formula,
        enable_table=enable_table,
        poll_interval_seconds=poll_interval_seconds,
        max_polls=max_polls,
    )
    result = batch_result["results"][0]
    markdown = clean_mineru_markdown(result.get("markdown", ""))
    preview = normalize_preview(markdown)
    if result.get("state") == "done" and has_meaningful_text(markdown):
        return ExtractionResult(
            extractor_name="mineru:batch",
            extraction_status="已提取文本",
            extracted_text=markdown,
            preview_text=preview,
            text_length=len(markdown),
            page_count=None,
            source_encoding="utf-8",
            note=f"已通过 MinerU 批量接口完成 Markdown 提取，批次号：{batch_result['batch_id']}。",
        )
    if markdown.strip():
        return ExtractionResult(
            extractor_name="mineru:batch",
            extraction_status="待审核",
            extracted_text=markdown,
            preview_text=preview,
            text_length=len(markdown),
            page_count=None,
            source_encoding="utf-8",
            note=f"MinerU 已返回少量 Markdown，当前需人工复核，批次号：{batch_result['batch_id']}。",
        )
    return ExtractionResult(
        extractor_name="mineru:batch",
        extraction_status="待人工复核",
        extracted_text="",
        preview_text="",
        text_length=0,
        page_count=None,
        source_encoding="",
        note=f"MinerU 未返回可用 Markdown，状态：{result.get('state')}，错误：{result.get('error') or '未注明'}，批次号：{batch_result['batch_id']}。",
    )
