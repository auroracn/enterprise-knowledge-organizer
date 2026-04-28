from __future__ import annotations

import argparse
import importlib.util
import json
import os
import requests
import sys
import tempfile
import unittest
from pathlib import Path
from unittest.mock import Mock, patch

import minimum_workflow.extractors as extractors_module
import minimum_workflow.sample_docx_extract_to_md as sample_docx_extract_to_md
from minimum_workflow.cli import (
    build_parser,
    build_review_markdown_name,
    build_scan_sample_id,
    build_scanned_sample,
    cleanup_review_outputs,
    collect_duplicate_scan_sample_ids,
    collect_scan_sources,
    collect_scan_sources_with_skips,
    is_image_directory_candidate,
    main,
    resolve_qwen_runtime,
    resolve_mineru_token,
    run_source_dir,
    sanitize_review_markdown_stem,
    select_preferred_scan_sources,
)
from minimum_workflow.contracts import ExtractionResult, SampleRecord, get_sample_by_id, load_contract
from minimum_workflow.extractors import build_bid_summary_metadata, clean_mineru_markdown, extract_legacy_excel_text, extract_pdf_text, extract_pdf_with_mineru, extract_text, request_mineru_batch_upload_urls
from minimum_workflow.field_extractors import extract_fields, normalize_policy_date
from minimum_workflow.qwen_client import enrich_payload_with_qwen, normalize_solution_summary_payload
from minimum_workflow.pipeline import (
    DIRECTORY_TYPE,
    build_effective_sample,
    build_markdown,
    build_structured_payload,
    decide_processing_route,
    detect_file_type,
    extract_with_strategy,
    merge_qwen_updates,
    resolve_qwen_runtime as resolve_pipeline_qwen_runtime,
    should_skip_qwen_for_sample,
)
from run_claude_output_workflow import (
    build_internal_output_root,
    build_review_output_dir,
    candidate_excerpt_lines,
    pick_acceptance_snippets,
    run_acceptance_for_item,
    run_directory,
    write_skipped_files_csv,
)


WEB_HELPER_PATH = Path(__file__).resolve().parents[1] / "知识整理助手.py"
WEB_HELPER_SPEC = importlib.util.spec_from_file_location("changfeng_web_helper", WEB_HELPER_PATH)
assert WEB_HELPER_SPEC and WEB_HELPER_SPEC.loader
knowledge_helper_ui = importlib.util.module_from_spec(WEB_HELPER_SPEC)
WEB_HELPER_SPEC.loader.exec_module(knowledge_helper_ui)


class MinimumWorkflowTest(unittest.TestCase):
    def test_docx_table_cell_lines_do_not_emit_html_breaks(self) -> None:
        self.assertEqual(sample_docx_extract_to_md.merge_table_cell_lines(["吴涛", "18523899765"]), "吴涛 18523899765")
        self.assertEqual(extractors_module.merge_docx_table_cell_lines(["吴涛", "18523899765"]), "吴涛 18523899765")
        self.assertNotIn("<br>", sample_docx_extract_to_md.merge_table_cell_lines(["吴涛", "18523899765"]))

    def test_wrapper_candidate_excerpt_lines_filters_metadata(self) -> None:
        text = """# 标题\n文档分类：方案/案例\n这是一个用于抽查的正文句子，长度足够。\n| 列1 | 列2 |\n另一条可以命中的正文短句，适合搜索。\n"""

        snippets = candidate_excerpt_lines(text)

        self.assertIn("这是一个用于抽查的正文句子，长度足够", snippets)
        self.assertIn("另一条可以命中的正文短句，适合搜索", snippets)
        self.assertTrue(all("文档分类" not in item for item in snippets))

    def test_wrapper_pick_acceptance_snippets_is_stable(self) -> None:
        text = "\n".join(
            [
                "第一条正文短句用于抽查命中。",
                "第二条正文短句用于抽查命中。",
                "第三条正文短句用于抽查命中。",
                "第四条正文短句用于抽查命中。",
            ]
        )

        first = pick_acceptance_snippets(text, seed="demo-seed")
        second = pick_acceptance_snippets(text, seed="demo-seed")

        self.assertEqual(first, second)
        self.assertEqual(len(first), 3)

    def test_sanitize_review_markdown_stem_replaces_webpage_fragment_with_parent_name(self) -> None:
        self.assertEqual(sanitize_review_markdown_stem("查看全文】", "全国通用航空企业名录（2026年1月版）"), "全国通用航空企业名录（2026年1月版）")
        self.assertEqual(build_review_markdown_name(Path("D:/资源库/全国通用航空企业名录（2026年1月版）/查看全文】.html")), "全国通用航空企业名录（2026年1月版）.md")

    def test_sanitize_review_markdown_stem_keeps_normal_title(self) -> None:
        self.assertEqual(sanitize_review_markdown_stem("贵州低空经济项目总体规划方案"), "贵州低空经济项目总体规划方案")
        self.assertEqual(build_review_markdown_name(Path("D:/资源库/贵州低空经济项目总体规划方案.pdf")), "贵州低空经济项目总体规划方案.md")

    def test_build_review_markdown_name_uses_inferred_title_for_weak_names(self) -> None:
        # 当源文件名是弱语义占位名时，优先使用推断的标题
        self.assertEqual(
            build_review_markdown_name(
                Path("D:/资源库/全国通用航空企业名录（2026年1月版）/查看全文】.md"),
                Path("D:/资源库"),
                inferred_title="无人驾驶航空器飞行管理暂行条例",
            ),
            "无人驾驶航空器飞行管理暂行条例.md",
        )
        # 当源文件名本身有语义时，不替换
        self.assertEqual(
            build_review_markdown_name(
                Path("D:/资源库/贵州低空经济项目总体规划方案.pdf"),
                inferred_title="某个推断标题",
            ),
            "贵州低空经济项目总体规划方案.md",
        )

    def test_cleanup_review_outputs_removes_legacy_web_fragment_name(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        review_output_root = Path(temp_dir.name) / "Claude输出"
        source_dir = Path("D:/资源库")
        source_path = source_dir / "全国通用航空企业名录（2026年1月版）" / "查看全文】.html"
        review_output_root.mkdir(parents=True)
        legacy_path = review_output_root / "查看全文】.md"
        cleaned_path = review_output_root / "全国通用航空企业名录（2026年1月版）.md"
        legacy_path.write_text("旧残片稿", encoding="utf-8")
        cleaned_path.write_text("新稿", encoding="utf-8")

        cleanup_review_outputs(
            review_output_root,
            source_dir,
            [{"source_path": str(source_path)}],
            set(),
        )

        self.assertFalse(legacy_path.exists())
        self.assertFalse(cleaned_path.exists())

    def test_wrapper_candidate_excerpt_lines_skips_attachment_titles_and_keeps_body_items(self) -> None:
        text = "\n".join(
            [
                "附件一：联系人信息表",
                "一、低空巡检平台覆盖山区输电走廊日常巡检",
                "二、防灾减灾应急指挥支持洪涝现场快速勘测",
                "附件二",
            ]
        )

        snippets = candidate_excerpt_lines(text)

        self.assertNotIn("附件一：联系人信息表", snippets)
        self.assertNotIn("附件二", snippets)
        self.assertIn("低空巡检平台覆盖山区输电走廊日常巡检", snippets)
        self.assertIn("防灾减灾应急指挥支持洪涝现场快速勘测", snippets)

    def test_wrapper_pick_acceptance_snippets_falls_back_when_body_lines_are_short(self) -> None:
        text = "\n".join(
            [
                "巡检覆盖周边农田地块",
                "应急支持重点水利区域",
                "保障山区森林防火巡查",
            ]
        )

        snippets = pick_acceptance_snippets(text, seed="short-lines")

        self.assertEqual(len(snippets), 3)
        self.assertEqual(snippets, ["巡检覆盖周边农田地块", "应急支持重点水利区域", "保障山区森林防火巡查"])

    def test_wrapper_run_acceptance_for_item_passes_when_markdown_reflows_symbols_and_line_breaks(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        internal_output_root = Path(temp_dir.name) / "internal"
        sample_id = "scan_reflow"
        sample_dir = internal_output_root / sample_id
        sample_dir.mkdir(parents=True)
        extracted_text = "\n".join(
            [
                "联系人及电话 贵州长风科技有限公司 0851-12345678 13900001111",
                "低空巡检平台覆盖山区输电走廊日常巡检",
                "防灾减灾应急指挥支持洪涝现场快速勘测",
            ]
        )
        (sample_dir / "extracted.txt").write_text(extracted_text, encoding="utf-8")
        markdown_path = Path(temp_dir.name) / "reflow.md"
        markdown_path.write_text(
            "# 标题\n\n- 联系人及电话\n- 贵州长风科技有限公司\n- 0851 12345678 / 13900001111\n- 低空巡检平台 覆盖 山区输电走廊 日常巡检\n- 防灾减灾应急指挥 支持 洪涝现场快速勘测",
            encoding="utf-8",
        )

        result = run_acceptance_for_item(
            {
                "sample_id": sample_id,
                "source_path": "D:/demo/reflow.pdf",
                "structured_markdown_path": str(markdown_path),
            },
            internal_output_root,
        )

        self.assertTrue(result.passed)
        self.assertEqual(result.passed_count, 3)
        self.assertEqual(result.missing_snippets, [])
        self.assertIn("归一化匹配通过部分短句", result.note)

    def test_wrapper_run_acceptance_for_item_passes_when_all_snippets_hit(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        internal_output_root = Path(temp_dir.name) / "internal"
        sample_id = "scan_demo"
        sample_dir = internal_output_root / sample_id
        sample_dir.mkdir(parents=True)
        extracted_text = "\n".join(
            [
                "第一条正文短句用于抽查命中。",
                "第二条正文短句用于抽查命中。",
                "第三条正文短句用于抽查命中。",
            ]
        )
        (sample_dir / "extracted.txt").write_text(extracted_text, encoding="utf-8")
        markdown_path = Path(temp_dir.name) / "demo.md"
        markdown_path.write_text("# 标题\n\n" + extracted_text, encoding="utf-8")

        result = run_acceptance_for_item(
            {
                "sample_id": sample_id,
                "source_path": "D:/demo/source.pdf",
                "structured_markdown_path": str(markdown_path),
            },
            internal_output_root,
        )

        self.assertTrue(result.passed)
        self.assertEqual(result.passed_count, 3)
        self.assertEqual(result.missing_snippets, [])

    def test_wrapper_run_acceptance_for_item_passes_for_table_restructured_markdown(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        internal_output_root = Path(temp_dir.name) / "internal"
        sample_id = "scan_table"
        sample_dir = internal_output_root / sample_id
        sample_dir.mkdir(parents=True)
        extracted_text = "\n".join(
            [
                "贵州低空经济项目近期招投标信息汇总。",
                "项目名称为低空巡检平台建设项目。",
                "招标人为某某发展集团有限公司。",
            ]
        )
        (sample_dir / "extracted.txt").write_text(extracted_text, encoding="utf-8")
        markdown_path = Path(temp_dir.name) / "table.md"
        markdown_path.write_text(
            "# 汇总表\n\n| 地区 | 项目名称 | 招标人 |\n| --- | --- | --- |\n| 贵州 | 低空巡检平台建设项目 | 某某发展集团有限公司 |\n| 备注 | 贵州低空经济项目近期招投标信息汇总 | 更新至2026年03月 |",
            encoding="utf-8",
        )

        result = run_acceptance_for_item(
            {
                "sample_id": sample_id,
                "source_path": "D:/demo/table.pdf",
                "structured_markdown_path": str(markdown_path),
            },
            internal_output_root,
        )

        self.assertTrue(result.passed)
        self.assertEqual(result.passed_count, 3)
        self.assertEqual(result.missing_snippets, [])
        self.assertIn("归一化匹配通过", result.note)

    def test_wrapper_write_skipped_files_csv_outputs_expected_rows(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name)
        report = {
            "items": [
                {"status": "skipped_duplicate", "source_path": "a.pdf", "preferred_source_path": "a.docx", "reason": "低优先级近似稿"},
                {"status": "success", "source_path": "b.pdf"},
                {"status": "failed", "source_path": "c.pdf", "error": "boom"},
            ]
        }

        csv_path = write_skipped_files_csv(source_dir, report)

        content = csv_path.read_text(encoding="utf-8-sig")
        self.assertIn("skipped_duplicate", content)
        self.assertIn("a.docx", content)
        self.assertIn("failed", content)
        self.assertIn("boom", content)
        self.assertNotIn("success,b.pdf", content)

    def test_wrapper_run_directory_writes_trace_and_returns_failure_on_acceptance_miss(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "source"
        output_root = Path(temp_dir.name) / "Claude输出"
        source_dir.mkdir()

        review_output_dir = build_review_output_dir(source_dir, output_root)
        review_output_dir.mkdir(parents=True)
        final_md = review_output_dir / "样本方案.md"
        final_md.write_text("# 样本方案\n\n只保留一条正文。", encoding="utf-8")

        internal_output_root = build_internal_output_root(source_dir)
        sample_dir = internal_output_root / "scan_样本方案"
        sample_dir.mkdir(parents=True, exist_ok=True)
        (sample_dir / "extracted.txt").write_text(
            "第一条正文短句用于抽查命中。\n第二条正文短句用于抽查命中。\n第三条正文短句用于抽查命中。",
            encoding="utf-8",
        )
        report_path = internal_output_root / "scan_report.json"
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(
            json.dumps(
                {
                    "total_count": 1,
                    "selected_count": 1,
                    "success_count": 1,
                    "failed_count": 0,
                    "skipped_duplicate_count": 0,
                    "skipped_photo_count": 0,
                    "skipped_non_source_count": 0,
                    "items": [
                        {
                            "status": "success",
                            "sample_id": "scan_样本方案",
                            "source_path": str(source_dir / "样本方案.pdf"),
                            "structured_markdown_path": str(final_md),
                        }
                    ],
                },
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )

        with patch("run_claude_output_workflow.run_source_dir", return_value=0):
            result = run_directory(
                source_dir,
                output_root=output_root,
                pdf_extractor="mineru",
                mineru_token=None,
                enable_ocr=False,
                enable_qwen=False,
                qwen_runtime={},
            )

        self.assertEqual(result, 1)
        self.assertTrue((source_dir / "链路说明.md").exists())
        self.assertTrue((source_dir / "skipped_files.csv").exists())
        trace_content = (source_dir / "链路说明.md").read_text(encoding="utf-8")
        self.assertIn("抽查未通过数：1", trace_content)
        self.assertIn("未命中短句", trace_content)
        self.assertIn("终稿 Markdown 已落在 Claude输出 子目录中。", trace_content)

    def test_wrapper_run_directory_trace_note_is_accurate_when_no_success_output(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "source"
        output_root = Path(temp_dir.name) / "Claude输出"
        source_dir.mkdir()

        internal_output_root = build_internal_output_root(source_dir)
        report_path = internal_output_root / "scan_report.json"
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(
            json.dumps(
                {
                    "total_count": 1,
                    "selected_count": 1,
                    "success_count": 0,
                    "failed_count": 1,
                    "skipped_duplicate_count": 0,
                    "skipped_photo_count": 0,
                    "skipped_non_source_count": 0,
                    "items": [
                        {
                            "status": "failed",
                            "sample_id": "scan_样本方案",
                            "source_path": str(source_dir / "样本方案.pdf"),
                            "error": "boom",
                        }
                    ],
                },
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )

        with patch("run_claude_output_workflow.run_source_dir", return_value=0):
            result = run_directory(
                source_dir,
                output_root=output_root,
                pdf_extractor="mineru",
                mineru_token=None,
                enable_ocr=False,
                enable_qwen=False,
                qwen_runtime={},
            )

        self.assertEqual(result, 0)
        trace_content = (source_dir / "链路说明.md").read_text(encoding="utf-8")
        self.assertIn("本次未生成终稿 Markdown。", trace_content)
        self.assertNotIn("终稿 Markdown 已落在 Claude输出 子目录中。", trace_content)

        source_path = Path("D:/tmp/智慧园区巡检技术方案.docx")
        batch_result = {
            "batch_id": "batch-docx-failed",
            "results": [
                {
                    "state": "failed",
                    "markdown": "",
                    "error": "mineru conversion failed",
                }
            ],
        }
        markitdown_result = ExtractionResult(
            extractor_name="markitdown:docx",
            extraction_status="已提取文本",
            extracted_text="# 智慧园区巡检技术方案\n\n项目背景\n\n巡检机器人",
            preview_text="# 智慧园区巡检技术方案",
            text_length=30,
            page_count=None,
            source_encoding="utf-8",
            note="已通过 MarkItDown 完成 DOCX 备用解析。",
        )
        local_blocks = [
            "智慧园区巡检技术方案",
            "| 参数 | 值 |\n| --- | --- |\n| 续航 | 3小时 |",
        ]

        with patch.object(sample_docx_extract_to_md, "resolve_mineru_token", return_value="token-demo"), patch.object(
            sample_docx_extract_to_md,
            "run_mineru_batch",
            return_value=batch_result,
        ), patch.object(
            sample_docx_extract_to_md,
            "try_extract_with_markitdown",
            return_value=markitdown_result,
            create=True,
        ), patch.object(
            sample_docx_extract_to_md,
            "extract_docx_blocks",
            return_value=local_blocks,
        ):
            result = sample_docx_extract_to_md.extract_source_content(source_path)

        metadata = dict(result["auto_metadata"])
        self.assertEqual(metadata["抽取器"], "markitdown:docx")
        self.assertEqual(metadata["转换状态"], "MarkItDown降级成功")
        self.assertIn("MarkItDown备用解析", metadata["处理链路"])
        self.assertIn("MinerU 未返回可用 Markdown", result["extraction_note"])
        self.assertIn("| 参数 | 值 |", result["extracted_text"])

    def test_extract_office_document_content_appends_local_docx_tables_after_mineru(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        docx_path = Path(temp_dir.name) / "机会清单.docx"
        docx_path.write_bytes(b"fake-docx")
        mineru_result = {
            "batch_id": "batch-123",
            "results": [{"state": "done", "markdown": "# 重庆市2025年第一批低空经济应用场景机会清单\n\n正文连续文本"}],
        }
        local_blocks = [
            "重庆市2025年第一批低空经济应用场景机会清单",
            "| 场景类型 | 序号 | 场景机会名称 | 单位名称 | 合作需求 | 联系方式 |\n| --- | --- | --- | --- | --- | --- |\n| 低空安全保障 | 1 | 全市无人机无线信标接收系统和重点区域无人机监测管控 | 重庆市公安局 | 联合建设 | 023-63961517 |",
        ]

        with patch("minimum_workflow.sample_docx_extract_to_md.resolve_mineru_token", return_value="token-demo"), patch(
            "minimum_workflow.sample_docx_extract_to_md.run_mineru_batch",
            return_value=mineru_result,
        ), patch(
            "minimum_workflow.sample_docx_extract_to_md.extract_docx_blocks",
            return_value=local_blocks,
        ):
            result = sample_docx_extract_to_md.extract_office_document_content(docx_path, "word")

        self.assertEqual(result["extraction_result"].extractor_name, "mineru:batch")
        self.assertIn("| 场景类型 | 序号 | 场景机会名称 | 单位名称 | 合作需求 | 联系方式 |", result["extracted_text"])
        self.assertIn("已补充本地 docx 表格结构提取", result["extraction_note"])
        self.assertTrue(any(block.lstrip().startswith("|") for block in result["blocks"]))

    def test_extract_document_image_goes_to_ocr_placeholder(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "设备检测报告.png"
        sample_path.write_bytes(b"fake-image")

        result = extract_text(sample_path, "image")

        self.assertEqual(result.extraction_status, "待OCR")
        self.assertEqual(result.extractor_name, "ocr:placeholder:image")
        self.assertEqual(result.page_count, 1)
        self.assertIn("OCR 占位层", result.note)
        self.assertIn("当前未启用真实 OCR", result.note)

    def test_extract_photo_image_is_skipped(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "现场照片.png"
        sample_path.write_bytes(b"fake-image")

        result = extract_text(sample_path, "image")

        self.assertEqual(result.extraction_status, "跳过")
        self.assertEqual(result.extractor_name, "skip:image_photo")
        self.assertEqual(result.page_count, 1)
        self.assertIn("纯照片处理", result.note)

    def test_extract_wechat_document_image_goes_to_ocr_placeholder(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "微信图片_20260409162847_364_69.jpg"
        sample_path.write_bytes(b"fake-image")

        result = extract_text(sample_path, "image")

        self.assertEqual(result.extraction_status, "待OCR")
        self.assertEqual(result.extractor_name, "ocr:placeholder:image")
        self.assertEqual(result.page_count, 1)
        self.assertIn("当前未启用真实 OCR", result.note)

    def test_extract_unlabeled_paged_image_goes_to_ocr_placeholder(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "001.jpg"
        sample_path.write_bytes(b"fake-image")

        result = extract_text(sample_path, "image")

        self.assertEqual(result.extraction_status, "待OCR")
        self.assertEqual(result.extractor_name, "ocr:placeholder:image")
        self.assertEqual(result.page_count, 1)
        self.assertIn("当前未启用真实 OCR", result.note)

    def test_extract_non_document_image_is_skipped(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "tyxx.jpg"
        sample_path.write_bytes(b"fake-image")

        result = extract_text(sample_path, "image")

        self.assertEqual(result.extraction_status, "跳过")
        self.assertEqual(result.extractor_name, "skip:image_photo")
        self.assertIn("纯照片处理", result.note)

    def test_extract_planning_notice_image_uses_ocr_placeholder(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "《百里杜鹃管理区百纳彝族乡总体规划（2020-2035）》公示.jpg"
        sample_path.write_bytes(b"fake-image")

        result = extract_text(sample_path, "image")

        self.assertEqual(result.extraction_status, "待OCR")
        self.assertEqual(result.extractor_name, "ocr:placeholder:image")
        self.assertIn("当前未启用真实 OCR", result.note)

    def test_extract_text_uses_markitdown_when_docx_local_parser_fails(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "sample.docx"
        sample_path.write_bytes(b"fake-docx")
        expected = ExtractionResult(
            extractor_name="markitdown:docx",
            extraction_status="已提取文本",
            extracted_text="# 标题\n\n正文",
            preview_text="# 标题 正文",
            text_length=8,
            page_count=None,
            source_encoding="utf-8",
            note="已通过 MarkItDown 完成 .docx 备用解析。",
        )
        failed_local = ExtractionResult(
            extractor_name="word:none",
            extraction_status="待人工复核",
            extracted_text="",
            preview_text="",
            text_length=0,
            page_count=None,
            source_encoding="",
            note="docx 读取失败：bad zip",
        )

        with patch("minimum_workflow.extractors.extract_docx_text", return_value=failed_local) as local_mock, patch(
            "minimum_workflow.sample_docx_extract_to_md.try_extract_with_markitdown",
            return_value=expected,
        ) as markitdown_mock:
            result = extract_text(sample_path, "word")

        local_mock.assert_called_once_with(sample_path)
        markitdown_mock.assert_called_once_with(sample_path, "word")
        self.assertEqual(result.extractor_name, "markitdown:docx")
        self.assertIn("docx 读取失败", result.note)

    def test_extract_text_reads_utf8_sig_text_file(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "sample.txt"
        sample_path.write_text(
            "文件名称：重庆低空经济政策\n第一条 支持低空示范场景。",
            encoding="utf-8-sig",
        )

        result = extract_text(sample_path, "txt")

        self.assertEqual(result.extraction_status, "已提取文本")
        self.assertEqual(result.extractor_name, "text:utf-8-sig")
        self.assertGreater(result.text_length, 0)
        self.assertIn("重庆低空经济政策", result.preview_text)

    def test_extract_text_cleans_webpage_shell_markdown(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "《2026 年低空经济全产业链薪酬调研报告》-航投人才.md"
        sample_path.write_text(
            "# 《2026 年低空经济全产业链薪酬调研报告》-航投人才\n"
            "原文: https://www.laernoc.com/newsinfo/demo\n\n"
            "行业搜索引擎\n首页\n新闻资讯\n国内资讯\n国际资讯\n政策法规\n法律法规\n国家政策\n地方政策\n行业标准\n产品库\n无人机\n应用场景\n低空报告\n行业报告\n会展赛事\n关于我们\n联系我们\n\n"
            "《2026 年低空经济全产业链薪酬调研报告》-航投人才\n"
            "创建时间： 2026-02-28 16:18:00         浏览量：1116\n"
            "航投人才发布《2026 年低空经济全产业链薪酬调研报告》，正文保留。\n"
            "来源：航投人才\n"
            " 来源： \n"
            "前一个：上一条\n"
            "免责声明：本站部分信息来源于互联网搜集。如有侵权，请告知，本站将立刻删除。\n"
            "快速链接\n",
            encoding="utf-8-sig",
        )

        result = extract_text(sample_path, "markdown")

        self.assertEqual(result.extraction_status, "已提取文本")
        self.assertTrue(result.extracted_text.startswith("# 《2026 年低空经济全产业链薪酬调研报告》-航投人才\n原文: https://www.laernoc.com/newsinfo/demo"))
        self.assertNotIn("行业搜索引擎", result.extracted_text)
        self.assertNotIn("创建时间", result.extracted_text)
        self.assertNotIn("前一个：", result.extracted_text)
        self.assertNotIn("免责声明", result.extracted_text)
        self.assertNotIn(" 来源： ", result.extracted_text)
        self.assertIn("航投人才发布《2026 年低空经济全产业链薪酬调研报告》", result.extracted_text)
        self.assertIn("来源：航投人才", result.extracted_text)
        self.assertIn("已清洗网页导航壳层", result.note)

    def test_extract_text_cleans_webpage_shell_markdown_with_frontmatter(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "webpage.md"
        sample_path.write_text(
            "---\n"
            "原文URL: https://example.com/article\n"
            "网页标题: 示范网页\n"
            "来源类型: html_url\n"
            "---\n\n"
            "# 示范网页\n"
            "原文: https://example.com/article\n\n"
            "行业搜索引擎\n首页\n新闻资讯\n国内资讯\n国际资讯\n政策法规\n法律法规\n国家政策\n地方政策\n行业标准\n产品库\n无人机\n应用场景\n低空报告\n行业报告\n会展赛事\n关于我们\n联系我们\n\n"
            "示范网页\n"
            "创建时间： 2026-02-28 16:18:00         浏览量：1116\n"
            "这里是真正正文。\n"
            "来源：示例站点\n"
            " 来源： \n"
            "免责声明：测试页脚\n",
            encoding="utf-8-sig",
        )

        result = extract_text(sample_path, "markdown")

        self.assertEqual(result.extraction_status, "已提取文本")
        self.assertIn("原文URL: https://example.com/article", result.extracted_text)
        self.assertIn("# 示范网页\n原文: https://example.com/article", result.extracted_text)
        self.assertIn("这里是真正正文。", result.extracted_text)
        self.assertNotIn("行业搜索引擎", result.extracted_text)
        self.assertNotIn("创建时间", result.extracted_text)
        self.assertNotIn("免责声明", result.extracted_text)
        self.assertIn("已清洗网页导航壳层", result.note)

        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "sample.pdf"
        sample_path.write_bytes(b"%PDF-1.4\n")

        # 这里模拟两个 PDF 提取器都拿不到可用文本，验证是否按规则降级到待 OCR。
        with patch(
            "minimum_workflow.extractors.extract_pdf_with_pypdf",
            return_value=("", 5, "已尝试使用 pypdf 提取。"),
        ), patch(
            "minimum_workflow.extractors.extract_pdf_with_pdfplumber",
            return_value=("", 5, "已尝试使用 pdfplumber 提取。"),
        ):
            result = extract_pdf_text(sample_path)

        self.assertEqual(result.extraction_status, "待OCR")
        self.assertEqual(result.extractor_name, "ocr:placeholder:pdf")
        self.assertEqual(result.page_count, 5)
        self.assertIn("建议后续接 OCR", result.note)
        self.assertIn("当前未启用真实 OCR", result.note)
        self.assertIn("当前未启用真实 OCR", result.note)

    def test_extract_pdf_text_uses_real_ocr_when_enabled(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "sample.pdf"
        sample_path.write_bytes(b"%PDF-1.4\n")
        ocr_result = ExtractionResult(
            extractor_name="ocr:mineru:pdf",
            extraction_status="已提取文本",
            extracted_text="中盛益华 AIoT 平台介绍与解决方案正文内容",
            preview_text="中盛益华 AIoT 平台介绍与解决方案正文内容",
            text_length=24,
            page_count=5,
            source_encoding="utf-8",
            note="已通过 MinerU OCR 完成扫描型 PDF 文本提取，批次号：batch-ocr-1。",
        )

        with patch(
            "minimum_workflow.extractors.extract_pdf_with_pypdf",
            return_value=("", 5, "已尝试使用 pypdf 提取。"),
        ), patch(
            "minimum_workflow.extractors.extract_pdf_with_pdfplumber",
            return_value=("", 5, "已尝试使用 pdfplumber 提取。"),
        ), patch(
            "minimum_workflow.extractors.extract_pdf_with_mineru_ocr",
            return_value=ocr_result,
        ) as ocr_mock:
            result = extract_pdf_text(sample_path, enable_ocr=True, ocr_token="token-demo")

        ocr_mock.assert_called_once_with(sample_path, "token-demo", page_count=5)
        self.assertEqual(result.extractor_name, "ocr:mineru:pdf")
        self.assertEqual(result.extraction_status, "已提取文本")
        self.assertGreater(result.text_length, 0)

    def test_extract_pdf_text_falls_back_to_placeholder_when_real_ocr_fails(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "sample.pdf"
        sample_path.write_bytes(b"%PDF-1.4\n")

        with patch(
            "minimum_workflow.extractors.extract_pdf_with_pypdf",
            return_value=("", 5, "已尝试使用 pypdf 提取。"),
        ), patch(
            "minimum_workflow.extractors.extract_pdf_with_pdfplumber",
            return_value=("", 5, "已尝试使用 pdfplumber 提取。"),
        ), patch(
            "minimum_workflow.extractors.extract_pdf_with_mineru_ocr",
            side_effect=RuntimeError("mock ocr error"),
        ):
            result = extract_pdf_text(sample_path, enable_ocr=True, ocr_token="token-demo")

        self.assertEqual(result.extractor_name, "ocr:placeholder:pdf")
        self.assertEqual(result.extraction_status, "待OCR")
        self.assertIn("已尝试真实 OCR，但调用失败：mock ocr error", result.note)

    def test_build_structured_payload_uses_mineru_pdf_strategy(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("policy_uav_regulation")
        extraction = ExtractionResult(
            extractor_name="mineru:batch",
            extraction_status="已提取文本",
            extracted_text="法规正文 Markdown",
            preview_text="法规正文 Markdown",
            text_length=20,
            page_count=None,
            source_encoding="utf-8",
            note="已通过 MinerU 批量接口完成 Markdown 提取，批次号：batch-003。",
        )

        with patch("minimum_workflow.pipeline.extract_with_strategy", return_value=extraction) as strategy_mock:
            payload, result = build_structured_payload(
                sample,
                contract,
                pdf_extractor="mineru",
                mineru_token="token-demo",
            )

        strategy_mock.assert_called_once()
        self.assertEqual(result.extractor_name, "mineru:batch")
        self.assertEqual(payload["处理路径"], "document_parse")
        self.assertEqual(payload["抽取器"], "mineru:batch")
        self.assertEqual(payload["抽取状态"], "已提取文本")

    def test_extract_document_image_uses_real_ocr_when_enabled(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "设备检测报告.png"
        sample_path.write_bytes(b"fake-image")
        ocr_result = ExtractionResult(
            extractor_name="ocr:mineru:image",
            extraction_status="已提取文本",
            extracted_text="设备检测报告正文内容",
            preview_text="设备检测报告正文内容",
            text_length=10,
            page_count=1,
            source_encoding="utf-8",
            note="已通过 MinerU OCR 完成图片文本提取，批次号：batch-img-1。",
        )

        with patch("minimum_workflow.extractors.extract_image_with_mineru_ocr", return_value=ocr_result) as ocr_mock:
            result = extract_text(sample_path, "image", enable_ocr=True, ocr_token="token-demo")

        ocr_mock.assert_called_once_with(sample_path, "token-demo")
        self.assertEqual(result.extractor_name, "ocr:mineru:image")
        self.assertEqual(result.extraction_status, "已提取文本")

    def test_extract_photo_image_skips_before_ocr_when_enabled(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "现场照片.png"
        sample_path.write_bytes(b"fake-image")

        with patch("minimum_workflow.extractors.extract_image_with_mineru_ocr") as ocr_mock:
            result = extract_text(sample_path, "image", enable_ocr=True, ocr_token="token-demo")

        ocr_mock.assert_not_called()
        self.assertEqual(result.extractor_name, "skip:image_photo")
        self.assertEqual(result.extraction_status, "跳过")
        self.assertIn("纯照片处理", result.note)

    def test_extract_document_image_falls_back_to_placeholder_when_real_ocr_fails(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "设备检测报告.png"
        sample_path.write_bytes(b"fake-image")

        with patch("minimum_workflow.extractors.extract_image_with_mineru_ocr", side_effect=RuntimeError("mock image ocr error")):
            result = extract_text(sample_path, "image", enable_ocr=True, ocr_token="token-demo")

        self.assertEqual(result.extractor_name, "ocr:placeholder:image")
        self.assertEqual(result.extraction_status, "待OCR")
        self.assertIn("已尝试真实 OCR，但调用失败：mock image ocr error", result.note)

    def test_build_structured_payload_marks_wait_review_when_pdf_waits_for_ocr(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("policy_uav_regulation")
        extraction = ExtractionResult(
            extractor_name="ocr:placeholder:pdf",
            extraction_status="待OCR",
            extracted_text="",
            preview_text="",
            text_length=0,
            page_count=41,
            source_encoding="",
            note="当前 PDF 未提取到可用文本，疑似扫描件或图片型 PDF，建议后续接 OCR。",
        )

        # 这里验证抽取层状态会继续传递到 pipeline 分流层，不会把待 OCR 文件误判成直接入库。
        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, result = build_structured_payload(sample, contract)

        self.assertEqual(result.extraction_status, "待OCR")
        self.assertEqual(payload["抽取状态"], "待OCR")
        self.assertEqual(payload["处理路径"], "ocr")
        self.assertEqual(payload["分流结果"], "待审核")
        self.assertFalse(payload["是否适合直接入库"])
        self.assertEqual(payload["页数"], 41)

    def test_build_structured_payload_keeps_ocr_route_for_real_ocr_result(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("supplier_zhongsheng_aiot_brochure")
        extraction = ExtractionResult(
            extractor_name="ocr:mineru:pdf",
            extraction_status="已提取文本",
            extracted_text="中盛益华 AIoT 平台介绍与解决方案正文内容，覆盖智慧园区与物联网管理能力。",
            preview_text="中盛益华 AIoT 平台介绍与解决方案正文内容",
            text_length=38,
            page_count=25,
            source_encoding="utf-8",
            note="已通过 MinerU OCR 完成扫描型 PDF 文本提取，批次号：batch-ocr-2。",
        )

        with patch("minimum_workflow.pipeline.extract_with_strategy", return_value=extraction):
            payload, result = build_structured_payload(sample, contract, enable_ocr=True, mineru_token="token-demo")

        self.assertEqual(result.extractor_name, "ocr:mineru:pdf")
        self.assertEqual(payload["处理路径"], "ocr")
        self.assertEqual(payload["抽取状态"], "已提取文本")
        self.assertEqual(payload["分流结果"], "待审核")
        self.assertGreater(payload["文本长度"], 0)

    def test_detect_file_type_and_route_cover_real_extensions(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        directory_path = Path(temp_dir.name) / "paged_dir"
        directory_path.mkdir()

        self.assertEqual(detect_file_type(Path("sample.json")), "json")
        self.assertEqual(detect_file_type(Path("sample.csv")), "csv")
        self.assertEqual(detect_file_type(Path("sample.log")), "log")
        self.assertEqual(detect_file_type(Path("sample.pptx")), "presentation")
        self.assertEqual(detect_file_type(Path("sample.webp")), "image")
        self.assertEqual(detect_file_type(directory_path), DIRECTORY_TYPE)

        self.assertEqual(decide_processing_route("json"), "text_direct")
        self.assertEqual(decide_processing_route("csv"), "text_direct")
        self.assertEqual(decide_processing_route("log"), "text_direct")
        self.assertEqual(decide_processing_route("presentation"), "document_parse")
        self.assertEqual(decide_processing_route(DIRECTORY_TYPE), "document_parse")

    def test_extract_presentation_text_reads_text_and_table(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        from pptx import Presentation

        sample_path = Path(temp_dir.name) / "sample.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(0, 0, 300, 120)
        textbox.text_frame.text = "广州红色国防科技教育培训基地方案"
        table = slide.shapes.add_table(2, 2, 0, 130, 300, 120).table
        table.cell(0, 0).text = "项目"
        table.cell(0, 1).text = "内容"
        table.cell(1, 0).text = "场景"
        table.cell(1, 1).text = "国防教育"
        presentation.save(sample_path)

        result = extract_text(sample_path, "presentation")

        self.assertEqual(result.extraction_status, "已提取文本")
        self.assertEqual(result.extractor_name, "presentation:pptx")
        self.assertIn("广州红色国防科技教育培训基地方案", result.extracted_text)
        self.assertIn("国防教育", result.extracted_text)

    def test_extract_excel_text_reads_legacy_xls_via_com_fallback(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        sample_path = Path(temp_dir.name) / "sample.xls"
        sample_path.write_bytes(b"fake-xls")
        bridged_result = ExtractionResult(
            extractor_name="excel:openpyxl",
            extraction_status="已提取文本",
            extracted_text="# 工作表：Sheet1\n\n| 列1 | 列2 |\n| --- | --- |\n| A | B |",
            preview_text="Sheet1 A B",
            text_length=42,
            page_count=1,
            source_encoding="xlsx",
            note="已完成 Excel 读取，共识别 1 个有效工作表。",
        )

        with patch("minimum_workflow.extractors.pythoncom") as pythoncom_mock, patch(
            "minimum_workflow.extractors.win32_client.DispatchEx",
        ) as dispatch_mock, patch(
            "minimum_workflow.extractors.extract_excel_text",
            return_value=bridged_result,
        ) as excel_mock:
            workbook = Mock()
            dispatch_mock.return_value.Workbooks.Open.return_value = workbook
            result = extract_legacy_excel_text(sample_path)

        pythoncom_mock.CoInitialize.assert_called_once()
        pythoncom_mock.CoUninitialize.assert_called_once()

        self.assertEqual(result.extraction_status, "已提取文本")
        self.assertEqual(result.extractor_name, "excel:com->openpyxl")
        self.assertIn("转存为 xlsx", result.note)
        workbook.SaveAs.assert_called_once()
        excel_mock.assert_called()

    def test_resolve_mineru_token_prefers_cli_token(self) -> None:
        with patch.dict(os.environ, {"MINERU_TOKEN": "env-token"}, clear=False):
            self.assertEqual(resolve_mineru_token("cli-token"), "cli-token")

    def test_resolve_mineru_token_falls_back_to_env(self) -> None:
        with patch.dict(os.environ, {"MINERU_TOKEN": "env-token"}, clear=False):
            self.assertEqual(resolve_mineru_token(None), "env-token")

    def test_resolve_mineru_token_falls_back_to_config(self) -> None:
        with patch.dict(os.environ, {}, clear=True), patch(
            "minimum_workflow.cli.load_runtime_settings",
            return_value={"mineru vlm大模型 用于转换md格式key": "cfg-token"},
        ):
            self.assertEqual(resolve_mineru_token(None), "cfg-token")

    def test_cli_parser_defaults_pdf_extractor_to_mineru(self) -> None:
        parser = build_parser()
        args = parser.parse_args(["--sample-id", "supplier_zhongsheng_aiot_brochure"])

        self.assertEqual(args.pdf_extractor, "mineru")

    def test_cli_parser_accepts_enable_ocr_flag(self) -> None:
        parser = build_parser()
        args = parser.parse_args(["--sample-id", "supplier_zhongsheng_aiot_brochure", "--enable-ocr", "--mineru-token", "token-demo"])

        self.assertTrue(args.enable_ocr)
        self.assertEqual(args.mineru_token, "token-demo")

    def test_cli_resolve_qwen_runtime_prefers_cli_values(self) -> None:
        with patch("minimum_workflow.cli.load_runtime_settings", return_value={"qwen_api_key": "cfg-key", "qwen api连接": "https://cfg", "qwen模型": "cfg-model"}):
            runtime = resolve_qwen_runtime(
                enable_qwen=True,
                cli_api_key="cli-key",
                cli_base_url="https://cli",
                cli_model="cli-model",
            )

        self.assertEqual(runtime["api_key"], "cli-key")
        self.assertEqual(runtime["base_url"], "https://cli")
        self.assertEqual(runtime["model"], "cli-model")

    def test_cli_resolve_qwen_runtime_falls_back_to_config(self) -> None:
        with patch("minimum_workflow.cli.load_runtime_settings", return_value={"qwen_api_key": "cfg-key", "qwen api连接": "https://cfg", "qwen模型": "cfg-model"}):
            runtime = resolve_qwen_runtime(
                enable_qwen=True,
                cli_api_key=None,
                cli_base_url=None,
                cli_model=None,
            )

        self.assertEqual(runtime["api_key"], "cfg-key")
        self.assertEqual(runtime["base_url"], "https://cfg")
        self.assertEqual(runtime["model"], "cfg-model")

    def test_cli_main_runs_source_dir(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name)
        output_dir = source_dir / "out"

        with patch.object(
            sys,
            "argv",
            ["minimum_workflow", "--source-dir", str(source_dir), "--output-dir", str(output_dir)],
        ), patch("minimum_workflow.cli.resolve_mineru_token", return_value=None), patch(
            "minimum_workflow.cli.resolve_qwen_runtime",
            return_value={},
        ), patch("minimum_workflow.cli.run_source_dir", return_value=0) as run_source_dir_mock:
            result = main()

        self.assertEqual(result, 0)
        run_source_dir_mock.assert_called_once_with(
            source_dir,
            output_dir=output_dir,
            pdf_extractor="mineru",
            mineru_token=None,
            enable_ocr=False,
            enable_qwen=False,
            qwen_runtime={},
            resume=False,
        )


    def test_run_source_dir_rejects_review_output_inside_source_dir(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "source"
        source_dir.mkdir()

        with self.assertRaisesRegex(ValueError, "输出目录不能位于源目录内部"):
            run_source_dir(
                source_dir,
                output_dir=source_dir / "Claude输出",
                pdf_extractor="mineru",
                mineru_token=None,
                enable_ocr=False,
                enable_qwen=False,
                qwen_runtime={},
            )

    def test_run_source_dir_rejects_internal_output_inside_source_dir(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "source"
        source_dir.mkdir()
        output_dir = Path(temp_dir.name) / "out"

        with self.assertRaisesRegex(ValueError, "内部输出目录不能位于源目录内部"):
            run_source_dir(
                source_dir,
                output_dir=output_dir,
                internal_output_dir=source_dir / ".omc" / "generated" / "directory_scan",
                pdf_extractor="mineru",
                mineru_token=None,
                enable_ocr=False,
                enable_qwen=False,
                qwen_runtime={},
            )

    def test_run_source_dir_reports_only_final_markdown_for_review(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "source"
        output_dir = Path(temp_dir.name) / "out"
        source_dir.mkdir()
        sample_path = source_dir / "样本方案.pdf"
        sample_path.write_bytes(b"fake-pdf")

        fake_result = Mock()
        fake_result.sample_id = "scan_样本方案"
        fake_result.output_dir = Path(temp_dir.name) / "internal" / "scan_样本方案"
        fake_result.output_dir.mkdir(parents=True, exist_ok=True)
        fake_result.structured_json_path = fake_result.output_dir / "structured.json"
        fake_result.structured_json_path.write_text(
            json.dumps({"抽取状态": "已提取文本"}, ensure_ascii=False),
            encoding="utf-8",
        )
        fake_result.structured_markdown_path = fake_result.output_dir / "structured.md"
        fake_result.structured_markdown_path.write_text("# 样本方案\n\n正文", encoding="utf-8")
        fake_result.status_path = fake_result.output_dir / "status.json"
        fake_result.extracted_text_path = fake_result.output_dir / "extracted.txt"

        with patch("minimum_workflow.cli.load_contract", return_value=Mock()), patch(
            "minimum_workflow.cli.collect_scan_sources",
            return_value=[sample_path],
        ), patch("minimum_workflow.cli.collect_duplicate_scan_sample_ids", return_value=set()), patch(
            "minimum_workflow.cli.run_pipeline",
            return_value=fake_result,
        ):
            result = run_source_dir(
                source_dir,
                output_dir=output_dir,
                pdf_extractor="mineru",
                mineru_token=None,
                enable_ocr=False,
                enable_qwen=False,
                qwen_runtime={},
            )

        self.assertEqual(result, 0)
        final_markdown_path = output_dir / "样本方案.md"
        self.assertTrue(final_markdown_path.exists())
        self.assertEqual(final_markdown_path.read_text(encoding="utf-8"), "# 样本方案\n\n正文")
        self.assertFalse((output_dir / "scan_report.json").exists())

        report_paths = list((Path.cwd() / ".omc" / "generated" / "directory_scan").glob("**/scan_report.json"))
        self.assertTrue(report_paths)
        report = json.loads(max(report_paths, key=lambda path: path.stat().st_mtime).read_text(encoding="utf-8"))
        self.assertEqual(report["success_count"], 1)
        self.assertEqual(report["items"][0]["structured_markdown_path"], str(final_markdown_path))
        self.assertNotIn("output_dir", report["items"][0])
        self.assertNotIn("structured_json_path", report["items"][0])
        self.assertNotIn("status_path", report["items"][0])
        self.assertNotIn("extracted_text_path", report["items"][0])

    def test_run_source_dir_blocks_waiting_ocr_pdf_from_review_output(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "source"
        output_dir = Path(temp_dir.name) / "out"
        source_dir.mkdir()
        pdf_path = source_dir / "扫描件.pdf"
        pdf_path.write_bytes(b"fake-pdf")

        fake_result = Mock()
        fake_result.sample_id = "scan_扫描件"
        fake_result.output_dir = Path(temp_dir.name) / "internal" / "scan_扫描件"
        fake_result.output_dir.mkdir(parents=True, exist_ok=True)
        fake_result.structured_json_path = fake_result.output_dir / "structured.json"
        fake_result.structured_json_path.write_text(
            json.dumps({"抽取状态": "待OCR"}, ensure_ascii=False),
            encoding="utf-8",
        )
        fake_result.structured_markdown_path = fake_result.output_dir / "structured.md"
        fake_result.structured_markdown_path.write_text("# 扫描件\n\nOCR 占位稿", encoding="utf-8")
        fake_result.status_path = fake_result.output_dir / "status.json"
        fake_result.extracted_text_path = fake_result.output_dir / "extracted.txt"

        with patch("minimum_workflow.cli.load_contract", return_value=Mock()), patch(
            "minimum_workflow.cli.run_pipeline",
            return_value=fake_result,
        ) as run_pipeline_mock:
            result = run_source_dir(
                source_dir,
                output_dir=output_dir,
                pdf_extractor="mineru",
                mineru_token=None,
                enable_ocr=False,
                enable_qwen=False,
                qwen_runtime={},
            )

        self.assertEqual(result, 1)
        run_pipeline_mock.assert_called_once()
        self.assertFalse((output_dir / "扫描件.md").exists())
        report_paths = list((Path.cwd() / ".omc" / "generated" / "directory_scan").glob("**/scan_report.json"))
        report = json.loads(max(report_paths, key=lambda path: path.stat().st_mtime).read_text(encoding="utf-8"))
        self.assertEqual(report["success_count"], 0)
        self.assertEqual(report["failed_count"], 1)
        self.assertIn("待OCR", report["items"][0]["error"])

    def test_run_source_dir_blocks_waiting_ocr_image_from_review_output(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "source"
        output_dir = Path(temp_dir.name) / "out"
        source_dir.mkdir()
        image_path = source_dir / "设备检测报告.png"
        image_path.write_bytes(b"fake-image")

        fake_result = Mock()
        fake_result.sample_id = "scan_设备检测报告"
        fake_result.output_dir = Path(temp_dir.name) / "internal" / "scan_设备检测报告"
        fake_result.output_dir.mkdir(parents=True, exist_ok=True)
        fake_result.structured_json_path = fake_result.output_dir / "structured.json"
        fake_result.structured_json_path.write_text(
            json.dumps({"抽取状态": "待OCR"}, ensure_ascii=False),
            encoding="utf-8",
        )
        fake_result.structured_markdown_path = fake_result.output_dir / "structured.md"
        fake_result.structured_markdown_path.write_text("# 设备检测报告\n\nOCR 占位稿", encoding="utf-8")
        fake_result.status_path = fake_result.output_dir / "status.json"
        fake_result.extracted_text_path = fake_result.output_dir / "extracted.txt"

        with patch("minimum_workflow.cli.load_contract", return_value=Mock()), patch(
            "minimum_workflow.cli.should_skip_image_file",
            return_value=False,
        ), patch("minimum_workflow.cli.run_pipeline", return_value=fake_result):
            result = run_source_dir(
                source_dir,
                output_dir=output_dir,
                pdf_extractor="mineru",
                mineru_token=None,
                enable_ocr=False,
                enable_qwen=False,
                qwen_runtime={},
            )

        self.assertEqual(result, 1)
        self.assertFalse((output_dir / "设备检测报告.md").exists())

    def test_run_source_dir_blocks_waiting_ocr_image_directory_from_review_output(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "低空智能网联体系参考架构（2024版）工信部装备工业发展中心2024-11-11"
        output_dir = Path(temp_dir.name) / "out"
        source_dir.mkdir()
        for name in ("001.jpg", "002.jpg", "003.jpg"):
            (source_dir / name).write_bytes(b"fake-image")

        fake_result = Mock()
        fake_result.sample_id = "scan_分页目录"
        fake_result.output_dir = Path(temp_dir.name) / "internal" / "scan_分页目录"
        fake_result.output_dir.mkdir(parents=True, exist_ok=True)
        fake_result.structured_json_path = fake_result.output_dir / "structured.json"
        fake_result.structured_json_path.write_text(
            json.dumps({"抽取状态": "待人工复核"}, ensure_ascii=False),
            encoding="utf-8",
        )
        fake_result.structured_markdown_path = fake_result.output_dir / "structured.md"
        fake_result.structured_markdown_path.write_text("# 分页目录\n\n未完成 OCR", encoding="utf-8")
        fake_result.status_path = fake_result.output_dir / "status.json"
        fake_result.extracted_text_path = fake_result.output_dir / "extracted.txt"

        with patch("minimum_workflow.cli.load_contract", return_value=Mock()), patch(
            "minimum_workflow.cli.run_pipeline",
            return_value=fake_result,
        ):
            result = run_source_dir(
                source_dir,
                output_dir=output_dir,
                pdf_extractor="mineru",
                mineru_token=None,
                enable_ocr=False,
                enable_qwen=False,
                qwen_runtime={},
            )

        self.assertEqual(result, 1)
        self.assertEqual(list(output_dir.glob("*.md")), [])

    def test_run_source_dir_allows_review_markdown_when_text_is_extracted_but_not_direct_import(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "source"
        output_dir = Path(temp_dir.name) / "out"
        source_dir.mkdir()
        doc_path = source_dir / "招标文件样例.pdf"
        doc_path.write_bytes(b"fake-pdf")

        fake_result = Mock()
        fake_result.sample_id = "scan_招标文件样例"
        fake_result.output_dir = Path(temp_dir.name) / "internal" / "scan_招标文件样例"
        fake_result.output_dir.mkdir(parents=True, exist_ok=True)
        fake_result.structured_json_path = fake_result.output_dir / "structured.json"
        fake_result.structured_json_path.write_text(
            json.dumps({"抽取状态": "已提取文本", "是否适合直接入库": False, "分流结果": "待审核"}, ensure_ascii=False),
            encoding="utf-8",
        )
        fake_result.structured_markdown_path = fake_result.output_dir / "structured.md"
        fake_result.structured_markdown_path.write_text("# 招标文件样例\n\n正文", encoding="utf-8")
        fake_result.status_path = fake_result.output_dir / "status.json"
        fake_result.extracted_text_path = fake_result.output_dir / "extracted.txt"

        with patch("minimum_workflow.cli.load_contract", return_value=Mock()), patch(
            "minimum_workflow.cli.run_pipeline",
            return_value=fake_result,
        ):
            result = run_source_dir(
                source_dir,
                output_dir=output_dir,
                pdf_extractor="mineru",
                mineru_token=None,
                enable_ocr=False,
                enable_qwen=False,
                qwen_runtime={},
            )

        self.assertEqual(result, 0)
        self.assertTrue((output_dir / "招标文件样例.md").exists())
        self.assertEqual((output_dir / "招标文件样例.md").read_text(encoding="utf-8"), "# 招标文件样例\n\n正文")

    def test_run_source_dir_removes_stale_review_markdown_when_publish_blocked(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "source"
        output_dir = Path(temp_dir.name) / "out"
        source_dir.mkdir()
        output_dir.mkdir()
        doc_path = source_dir / "招标文件样例.pdf"
        doc_path.write_bytes(b"fake-pdf")
        stale_output = output_dir / "招标文件样例.md"
        stale_output.write_text("旧终稿残留", encoding="utf-8")

        fake_result = Mock()
        fake_result.sample_id = "scan_招标文件样例"
        fake_result.output_dir = Path(temp_dir.name) / "internal" / "scan_招标文件样例"
        fake_result.output_dir.mkdir(parents=True, exist_ok=True)
        fake_result.structured_json_path = fake_result.output_dir / "structured.json"
        fake_result.structured_json_path.write_text(
            json.dumps({"抽取状态": "待OCR"}, ensure_ascii=False),
            encoding="utf-8",
        )
        fake_result.structured_markdown_path = fake_result.output_dir / "structured.md"
        fake_result.structured_markdown_path.write_text("# 招标文件样例\n\nOCR 占位稿", encoding="utf-8")
        fake_result.status_path = fake_result.output_dir / "status.json"
        fake_result.extracted_text_path = fake_result.output_dir / "extracted.txt"

        with patch("minimum_workflow.cli.load_contract", return_value=Mock()), patch(
            "minimum_workflow.cli.run_pipeline",
            return_value=fake_result,
        ):
            result = run_source_dir(
                source_dir,
                output_dir=output_dir,
                pdf_extractor="mineru",
                mineru_token=None,
                enable_ocr=False,
                enable_qwen=False,
                qwen_runtime={},
            )

        self.assertEqual(result, 1)
        self.assertFalse(stale_output.exists())

    def test_run_source_dir_removes_legacy_web_fragment_when_publishing_clean_name(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "source"
        output_dir = Path(temp_dir.name) / "out"
        nested_dir = source_dir / "全国通用航空企业名录（2026年1月版）"
        nested_dir.mkdir(parents=True)
        output_dir.mkdir()
        source_path = nested_dir / "查看全文】.md"
        source_path.write_text("原始网页转载内容", encoding="utf-8")

        legacy_output = output_dir / "查看全文】.md"
        legacy_output.write_text("旧残片稿", encoding="utf-8")

        fake_result = Mock()
        fake_result.sample_id = "scan_查看全文"
        fake_result.output_dir = Path(temp_dir.name) / "internal" / "scan_查看全文"
        fake_result.output_dir.mkdir(parents=True, exist_ok=True)
        fake_result.structured_json_path = fake_result.output_dir / "structured.json"
        fake_result.structured_json_path.write_text(
            json.dumps({"抽取状态": "已提取文本"}, ensure_ascii=False),
            encoding="utf-8",
        )
        fake_result.structured_markdown_path = fake_result.output_dir / "structured.md"
        fake_result.structured_markdown_path.write_text("# 未命名文档\n\n正文", encoding="utf-8")
        fake_result.status_path = fake_result.output_dir / "status.json"
        fake_result.extracted_text_path = fake_result.output_dir / "extracted.txt"

        with patch("minimum_workflow.cli.load_contract", return_value=Mock()), patch(
            "minimum_workflow.cli.run_pipeline",
            return_value=fake_result,
        ):
            result = run_source_dir(
                source_dir,
                output_dir=output_dir,
                pdf_extractor="mineru",
                mineru_token=None,
                enable_ocr=False,
                enable_qwen=False,
                qwen_runtime={},
            )

        self.assertEqual(result, 0)
        self.assertFalse(legacy_output.exists())
        self.assertTrue((output_dir / "全国通用航空企业名录（2026年1月版）.md").exists())

    def test_run_source_dir_skips_photo_images_from_review_output(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "source"
        output_dir = Path(temp_dir.name) / "out"
        source_dir.mkdir()
        photo_path = source_dir / "无人机飞行照片.png"
        doc_path = source_dir / "样本方案.pdf"
        photo_path.write_bytes(b"fake-image")
        doc_path.write_bytes(b"fake-pdf")

        fake_result = Mock()
        fake_result.sample_id = "scan_样本方案"
        fake_result.output_dir = Path(temp_dir.name) / "internal" / "scan_样本方案"
        fake_result.output_dir.mkdir(parents=True, exist_ok=True)
        fake_result.structured_json_path = fake_result.output_dir / "structured.json"
        fake_result.structured_json_path.write_text(
            json.dumps({"抽取状态": "已提取文本"}, ensure_ascii=False),
            encoding="utf-8",
        )
        fake_result.structured_markdown_path = fake_result.output_dir / "structured.md"
        fake_result.structured_markdown_path.write_text("# 样本方案\n\n正文", encoding="utf-8")
        fake_result.status_path = fake_result.output_dir / "status.json"
        fake_result.extracted_text_path = fake_result.output_dir / "extracted.txt"

        with patch("minimum_workflow.cli.load_contract", return_value=Mock()), patch(
            "minimum_workflow.cli.run_pipeline",
            return_value=fake_result,
        ) as run_pipeline_mock:
            result = run_source_dir(
                source_dir,
                output_dir=output_dir,
                pdf_extractor="mineru",
                mineru_token=None,
                enable_ocr=False,
                enable_qwen=False,
                qwen_runtime={},
            )

        self.assertEqual(result, 0)
        run_pipeline_mock.assert_called_once()
        self.assertTrue((output_dir / "样本方案.md").exists())
        self.assertFalse((output_dir / "无人机飞行照片.md").exists())
        self.assertFalse((output_dir / "链路说明.md").exists())
        report_paths = list((Path.cwd() / ".omc" / "generated" / "directory_scan").glob("**/scan_report.json"))
        report = json.loads(max(report_paths, key=lambda path: path.stat().st_mtime).read_text(encoding="utf-8"))
        self.assertEqual(report["skipped_photo_count"], 1)
        skipped_items = [item for item in report["items"] if item.get("status") == "skipped_photo"]
        self.assertEqual(len(skipped_items), 1)
        self.assertEqual(skipped_items[0]["source_path"], str(photo_path))

    def test_run_source_dir_cleans_stale_trace_markdown(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "source"
        output_dir = Path(temp_dir.name) / "out"
        source_dir.mkdir()
        output_dir.mkdir()
        trace_path = source_dir / "链路说明.md"
        doc_path = source_dir / "样本方案.pdf"
        trace_path.write_text("# 链路说明\n", encoding="utf-8")
        doc_path.write_bytes(b"fake-pdf")

        stale_trace_output = output_dir / "链路说明.md"
        stale_trace_output.write_text("旧污染内容", encoding="utf-8")

        fake_result = Mock()
        fake_result.sample_id = "scan_样本方案"
        fake_result.output_dir = Path(temp_dir.name) / "internal" / "scan_样本方案"
        fake_result.output_dir.mkdir(parents=True, exist_ok=True)
        fake_result.structured_json_path = fake_result.output_dir / "structured.json"
        fake_result.structured_json_path.write_text(
            json.dumps({"抽取状态": "已提取文本"}, ensure_ascii=False),
            encoding="utf-8",
        )
        fake_result.structured_markdown_path = fake_result.output_dir / "structured.md"
        fake_result.structured_markdown_path.write_text("# 样本方案\n\n正文", encoding="utf-8")
        fake_result.status_path = fake_result.output_dir / "status.json"
        fake_result.extracted_text_path = fake_result.output_dir / "extracted.txt"

        with patch("minimum_workflow.cli.load_contract", return_value=Mock()), patch(
            "minimum_workflow.cli.run_pipeline",
            return_value=fake_result,
        ) as run_pipeline_mock:
            result = run_source_dir(
                source_dir,
                output_dir=output_dir,
                pdf_extractor="mineru",
                mineru_token=None,
                enable_ocr=False,
                enable_qwen=False,
                qwen_runtime={},
            )

        self.assertEqual(result, 0)
        run_pipeline_mock.assert_called_once()
        self.assertFalse(stale_trace_output.exists())
        self.assertTrue((output_dir / "样本方案.md").exists())
        report_paths = list((Path.cwd() / ".omc" / "generated" / "directory_scan").glob("**/scan_report.json"))
        report = json.loads(max(report_paths, key=lambda path: path.stat().st_mtime).read_text(encoding="utf-8"))
        skipped_items = [item for item in report["items"] if item.get("status") == "skipped_non_source"]
        self.assertEqual(len(skipped_items), 1)
        self.assertEqual(skipped_items[0]["source_path"], str(trace_path))
        self.assertIn("过程说明文件", skipped_items[0]["reason"])


    def test_cli_main_rejects_output_dir_without_source_dir(self) -> None:
        with patch.object(sys, "argv", ["minimum_workflow", "--output-dir", "D:/tmp/out"]), patch(
            "minimum_workflow.cli.resolve_mineru_token",
            return_value=None,
        ), patch("minimum_workflow.cli.resolve_qwen_runtime", return_value={}):
            with self.assertRaises(SystemExit) as ctx:
                main()

        self.assertEqual(ctx.exception.code, 2)

    def test_cli_main_rejects_source_dir_with_sample_id(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        with patch.object(
            sys,
            "argv",
            [
                "minimum_workflow",
                "--source-dir",
                temp_dir.name,
                "--sample-id",
                "supplier_zhongsheng_aiot_brochure",
            ],
        ), patch("minimum_workflow.cli.resolve_mineru_token", return_value=None), patch(
            "minimum_workflow.cli.resolve_qwen_runtime",
            return_value={},
        ):
            with self.assertRaises(SystemExit) as ctx:
                main()

        self.assertEqual(ctx.exception.code, 2)

    def test_collect_scan_sources_includes_root_paged_image_directory(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "13、无人机图片"
        source_dir.mkdir()
        for name in ("11.png", "12.png", "13.png"):
            (source_dir / name).write_bytes(b"fake-image")

        sources = collect_scan_sources(source_dir)

        self.assertEqual(sources, [source_dir])
        self.assertTrue(is_image_directory_candidate(source_dir))

    def test_is_image_directory_candidate_rejects_mixed_document_directory(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "天翼航空相关资料"
        source_dir.mkdir()
        for name in ("tyxx.jpg", "tyxx1.jpg", "方案.pdf", "说明.docx"):
            (source_dir / name).write_bytes(b"fake-content")

        sources = collect_scan_sources(source_dir)

        self.assertFalse(is_image_directory_candidate(source_dir))
        self.assertCountEqual(sources, [source_dir / "方案.pdf", source_dir / "说明.docx"])

    def test_is_image_directory_candidate_ignores_process_artifacts(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "分页图片目录"
        source_dir.mkdir()
        for name in ("01封面.png", "02正文.png"):
            (source_dir / name).write_bytes(b"fake-image")
        (source_dir / "链路说明.md").write_text("# 链路说明\n", encoding="utf-8")
        (source_dir / "skipped_files.csv").write_text("status,source_path\n", encoding="utf-8")

        self.assertTrue(is_image_directory_candidate(source_dir))

    def test_build_scanned_sample_appends_suffix_for_duplicate_sample_ids(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "西安天翼"
        source_dir.mkdir()
        pdf_path = source_dir / "低空经济数智化教育解决方案.pdf"
        doc_path = source_dir / "低空经济数智化教育解决方案.doc"
        pdf_path.write_bytes(b"fake-pdf")
        doc_path.write_bytes(b"fake-doc")

        scan_sources = [doc_path, pdf_path]
        duplicate_sample_ids = collect_duplicate_scan_sample_ids(scan_sources, source_dir)

        self.assertEqual(duplicate_sample_ids, {build_scan_sample_id(Path("低空经济数智化教育解决方案.pdf"))})
        self.assertEqual(
            build_scanned_sample(doc_path, source_dir, duplicate_sample_ids=duplicate_sample_ids).sample_id,
            "scan_低空经济数智化教育解决方案_doc",
        )
        self.assertEqual(
            build_scanned_sample(pdf_path, source_dir, duplicate_sample_ids=duplicate_sample_ids).sample_id,
            "scan_低空经济数智化教育解决方案_pdf",
        )

    def test_select_preferred_scan_sources_prefers_office_over_pdf_in_same_directory(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "源目录"
        source_dir.mkdir()
        pdf_path = source_dir / "无人机设备合同.pdf"
        docx_path = source_dir / "无人机设备合同.docx"
        pdf_path.write_bytes(b"fake-pdf")
        docx_path.write_bytes(b"fake-docx")

        selected_sources, skipped_items = select_preferred_scan_sources(source_dir, [pdf_path, docx_path])

        self.assertEqual(selected_sources, [docx_path])
        self.assertEqual(len(skipped_items), 1)
        self.assertEqual(skipped_items[0]["status"], "skipped_duplicate")
        self.assertEqual(skipped_items[0]["source_path"], str(pdf_path))
        self.assertEqual(skipped_items[0]["preferred_source_path"], str(docx_path))

    def test_run_source_dir_reports_skipped_duplicate_source_with_reason(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "source"
        output_dir = Path(temp_dir.name) / "out"
        source_dir.mkdir()
        pdf_path = source_dir / "样本方案.pdf"
        docx_path = source_dir / "样本方案.docx"
        pdf_path.write_bytes(b"fake-pdf")
        docx_path.write_bytes(b"fake-docx")

        fake_result = Mock()
        fake_result.sample_id = "scan_样本方案"
        fake_result.output_dir = Path(temp_dir.name) / "internal" / "scan_样本方案"
        fake_result.output_dir.mkdir(parents=True, exist_ok=True)
        fake_result.structured_json_path = fake_result.output_dir / "structured.json"
        fake_result.structured_json_path.write_text(
            json.dumps({"抽取状态": "已提取文本"}, ensure_ascii=False),
            encoding="utf-8",
        )
        fake_result.structured_markdown_path = fake_result.output_dir / "structured.md"
        fake_result.structured_markdown_path.write_text("# 样本方案\n\n正文", encoding="utf-8")
        fake_result.status_path = fake_result.output_dir / "status.json"
        fake_result.extracted_text_path = fake_result.output_dir / "extracted.txt"

        with patch("minimum_workflow.cli.load_contract", return_value=Mock()), patch(
            "minimum_workflow.cli.collect_scan_sources",
            return_value=[pdf_path, docx_path],
        ), patch("minimum_workflow.cli.run_pipeline", return_value=fake_result) as run_pipeline_mock:
            result = run_source_dir(
                source_dir,
                output_dir=output_dir,
                pdf_extractor="mineru",
                mineru_token=None,
                enable_ocr=False,
                enable_qwen=False,
                qwen_runtime={},
            )

        self.assertEqual(result, 0)
        run_pipeline_mock.assert_called_once()
        final_markdown_path = output_dir / "样本方案.md"
        self.assertTrue(final_markdown_path.exists())
        report_paths = list((Path.cwd() / ".omc" / "generated" / "directory_scan").glob("**/scan_report.json"))
        report = json.loads(max(report_paths, key=lambda path: path.stat().st_mtime).read_text(encoding="utf-8"))
        self.assertEqual(report["selected_count"], 1)
        self.assertEqual(report["skipped_duplicate_count"], 1)
        skipped_items = [item for item in report["items"] if item.get("status") == "skipped_duplicate"]
        self.assertEqual(len(skipped_items), 1)
        self.assertEqual(skipped_items[0]["source_path"], str(pdf_path))
        self.assertEqual(skipped_items[0]["preferred_source_path"], str(docx_path))
        self.assertIn("按优先级", skipped_items[0]["reason"])

    def test_collect_scan_sources_with_skips_skips_photo_images(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "图片资料"
        source_dir.mkdir()
        photo_path = source_dir / "无人机飞行照片.png"
        doc_path = source_dir / "方案.pdf"
        photo_path.write_bytes(b"fake-image")
        doc_path.write_bytes(b"fake-pdf")

        sources, skipped_items = collect_scan_sources_with_skips(source_dir)

        self.assertEqual(sources, [doc_path])
        self.assertEqual(len(skipped_items), 1)
        self.assertEqual(skipped_items[0]["status"], "skipped_photo")
        self.assertEqual(skipped_items[0]["source_path"], str(photo_path))
        self.assertIn("纯照片处理", skipped_items[0]["reason"])

    def test_collect_scan_sources_with_skips_still_skips_photo_images_when_ocr_enabled(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "图片资料"
        source_dir.mkdir()
        photo_path = source_dir / "无人机飞行照片.png"
        doc_path = source_dir / "方案.pdf"
        photo_path.write_bytes(b"fake-image")
        doc_path.write_bytes(b"fake-pdf")

        sources, skipped_items = collect_scan_sources_with_skips(source_dir, enable_ocr=True)

        self.assertEqual(sources, [doc_path])
        self.assertEqual(len(skipped_items), 1)
        self.assertEqual(skipped_items[0]["status"], "skipped_photo")
        self.assertEqual(skipped_items[0]["source_path"], str(photo_path))
        self.assertIn("纯照片处理", skipped_items[0]["reason"])

    def test_collect_scan_sources_with_skips_skips_trace_markdown(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "资料目录"
        source_dir.mkdir()
        trace_path = source_dir / "链路说明.md"
        doc_path = source_dir / "方案.pdf"
        trace_path.write_text("# 链路说明\n", encoding="utf-8")
        doc_path.write_bytes(b"fake-pdf")

        sources, skipped_items = collect_scan_sources_with_skips(source_dir)

        self.assertEqual(sources, [doc_path])
        self.assertEqual(len(skipped_items), 1)
        self.assertEqual(skipped_items[0]["status"], "skipped_non_source")
        self.assertEqual(skipped_items[0]["source_path"], str(trace_path))
        self.assertIn("过程说明文件", skipped_items[0]["reason"])

    def test_collect_scan_sources_with_skips_skips_skipped_files_csv(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_dir = Path(temp_dir.name) / "资料目录"
        source_dir.mkdir()
        skipped_csv_path = source_dir / "skipped_files.csv"
        doc_path = source_dir / "方案.pdf"
        skipped_csv_path.write_text("status,source_path\n", encoding="utf-8")
        doc_path.write_bytes(b"fake-pdf")

        sources, skipped_items = collect_scan_sources_with_skips(source_dir)

        self.assertEqual(sources, [doc_path])
        self.assertEqual(len(skipped_items), 1)
        self.assertEqual(skipped_items[0]["status"], "skipped_non_source")
        self.assertEqual(skipped_items[0]["source_path"], str(skipped_csv_path))
        self.assertIn("过程记录文件", skipped_items[0]["reason"])


    def test_pipeline_resolve_qwen_runtime_returns_empty_when_disabled(self) -> None:
        runtime = resolve_pipeline_qwen_runtime(enable_qwen=False)
        self.assertEqual(runtime, {})

    def test_should_skip_qwen_for_policy_sample(self) -> None:
        policy_sample = get_sample_by_id("policy_uav_regulation")
        supplier_sample = get_sample_by_id("supplier_yuefei_2025")

        self.assertTrue(should_skip_qwen_for_sample(policy_sample))
        self.assertFalse(should_skip_qwen_for_sample(supplier_sample))

    def test_should_skip_qwen_for_directory_template_outside_allowlist(self) -> None:
        supplier_sample = get_sample_by_id("supplier_yuefei_2025")
        payload = {"推荐模板": "参考架构/白皮书口径（当前按原文全量提取Markdown输出）"}

        self.assertTrue(should_skip_qwen_for_sample(supplier_sample, payload, file_type=DIRECTORY_TYPE))
        self.assertFalse(should_skip_qwen_for_sample(supplier_sample, payload, file_type="pdf"))

    def test_build_effective_sample_uses_directory_auto_payload(self) -> None:
        base_sample = get_sample_by_id("supplier_yuefei_2025")
        payload = {
            "推荐模板": "方案案例模板",
            "文档分类": "行业参考架构/指导材料",
            "标题": "低空智能网联体系参考架构（2024版）",
            "主体名称": "工业和信息化部装备工业发展中心",
            "单位名称": "工业和信息化部装备工业发展中心",
        }

        effective = build_effective_sample(base_sample, payload)

        self.assertEqual(effective.recommended_template, "方案案例模板")
        self.assertEqual(effective.document_category, "行业参考架构/指导材料")
        self.assertEqual(effective.title_hint, "低空智能网联体系参考架构（2024版）")
        self.assertEqual(effective.subject_name_hint, "工业和信息化部装备工业发展中心")
        self.assertEqual(effective.unit_name_hint, "工业和信息化部装备工业发展中心")

    def test_build_structured_payload_skips_qwen_for_policy_template(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("policy_uav_regulation")
        extraction = ExtractionResult(
            extractor_name="mineru:batch",
            extraction_status="已提取文本",
            extracted_text="现公布《无人驾驶航空器飞行管理暂行条例》，自2024年1月1日起施行。",
            preview_text="无人驾驶航空器飞行管理暂行条例",
            text_length=40,
            page_count=None,
            source_encoding="utf-8",
            note="已通过 MinerU 批量接口完成 Markdown 提取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction), patch(
            "minimum_workflow.pipeline.resolve_qwen_runtime",
            return_value={"api_key": "demo-key", "base_url": "https://demo", "model": "qwen-plus"},
        ), patch("minimum_workflow.pipeline.enrich_payload_with_qwen") as qwen_mock:
            payload, _ = build_structured_payload(sample, contract, enable_qwen=True)

        qwen_mock.assert_not_called()
        self.assertEqual(payload["文件标题"], "无人驾驶航空器飞行管理暂行条例")
        self.assertEqual(payload["一级分类"], "")

    def test_build_structured_payload_applies_qwen_updates(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("supplier_yuefei_2025")
        extraction = ExtractionResult(
            extractor_name="text:utf-8-sig",
            extraction_status="已提取文本",
            extracted_text="公司名称：重庆跃飞智能科技有限公司",
            preview_text="重庆跃飞智能科技有限公司",
            text_length=20,
            page_count=None,
            source_encoding="utf-8-sig",
            note="已完成文本文件读取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction), patch(
            "minimum_workflow.pipeline.resolve_qwen_runtime",
            return_value={"api_key": "demo-key", "base_url": "https://demo", "model": "qwen-plus"},
        ), patch(
            "minimum_workflow.pipeline.enrich_payload_with_qwen",
            return_value={
                "一级分类": "供应商企业与厂家资料库",
                "二级分类": "无人机厂家",
                "分类置信度": 0.93,
                "分类依据": "正文出现公司名称与核心产品",
                "推荐模板": "供应商企业模板",
                "企业名称": "重庆跃飞智能科技有限公司",
            },
        ):
            payload, _ = build_structured_payload(sample, contract, enable_qwen=True)

        self.assertEqual(payload["一级分类"], "供应商企业与厂家资料库")
        self.assertEqual(payload["二级分类"], "无人机厂家")
        self.assertEqual(payload["分类置信度"], 0.93)
        self.assertEqual(payload["分类依据"], "正文出现公司名称与核心产品")
        self.assertEqual(payload["企业名称"], "重庆跃飞智能科技有限公司")

    def test_build_structured_payload_reextracts_fields_after_qwen_template_change(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("supplier_yuefei_2025")
        extraction = ExtractionResult(
            extractor_name="text:utf-8-sig",
            extraction_status="已提取文本",
            extracted_text=(
                "项目名称：低空巡检平台解决方案\n"
                "客户名称：重庆机场集团\n"
                "解决问题：提升巡检效率\n"
                "实施方式：无人机自动巡检\n"
                "效果数据：巡检时长缩短50%"
            ),
            preview_text="低空巡检平台解决方案",
            text_length=64,
            page_count=None,
            source_encoding="utf-8-sig",
            note="已完成文本文件读取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction), patch(
            "minimum_workflow.pipeline.resolve_qwen_runtime",
            return_value={"api_key": "demo-key", "base_url": "https://demo", "model": "qwen-plus"},
        ), patch(
            "minimum_workflow.pipeline.enrich_payload_with_qwen",
            return_value={
                "一级分类": "方案案例与实施资料库",
                "二级分类": "巡检解决方案",
                "分类置信度": 0.91,
                "分类依据": "正文包含项目名称、客户名称和实施方式",
                "文档分类": "方案/案例",
                "推荐模板": "方案案例模板",
            },
        ):
            payload, _ = build_structured_payload(sample, contract, enable_qwen=True)

        self.assertEqual(payload["推荐模板"], "方案案例模板")
        self.assertEqual(payload["文档分类"], "方案/案例")
        self.assertEqual(payload["方案名称字段"], "低空巡检平台解决方案")
        self.assertEqual(payload["客户名称字段"], "重庆机场集团")
        self.assertEqual(payload["解决问题字段"], "提升巡检效率")
        self.assertEqual(payload["实施方式字段"], "无人机自动巡检")
        self.assertEqual(payload["效果数据字段"], "巡检时长缩短50%")
        self.assertFalse(payload.get("企业名称"))

        markdown = build_markdown(sample, payload)

        self.assertIn("## 九、字段提取结果", markdown)
        self.assertIn("- 方案名称：低空巡检平台解决方案", markdown)
        self.assertNotIn("- 企业名称：", markdown)

    def test_build_structured_payload_uses_directory_auto_payload(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        directory_path = Path(temp_dir.name) / "低空智能网联体系参考架构（2024版）目录"
        directory_path.mkdir()
        contract = load_contract()
        base_sample = get_sample_by_id("supplier_yuefei_2025")
        sample = base_sample.__class__(
            sample_id="directory_reference_architecture",
            source_path=str(directory_path),
            document_category=base_sample.document_category,
            recommended_template=base_sample.recommended_template,
            title_hint=base_sample.title_hint,
            subject_name_hint=base_sample.subject_name_hint,
            product_name_hint=base_sample.product_name_hint,
            unit_name_hint=base_sample.unit_name_hint,
            tags=base_sample.tags,
            risks=base_sample.risks,
            notes=base_sample.notes,
            evidence_level=base_sample.evidence_level,
            fallback_decision=base_sample.fallback_decision,
            split_required=base_sample.split_required,
            split_note=base_sample.split_note,
        )
        extraction = ExtractionResult(
            extractor_name="ocr:mineru:image_directory",
            extraction_status="已提取文本",
            extracted_text="# 低空智能网联体系参考架构（2024版）\n\n工业和信息化部装备工业发展中心\n\n2024年11月",
            preview_text="低空智能网联体系参考架构（2024版） 工业和信息化部装备工业发展中心 2024年11月",
            text_length=53,
            page_count=3,
            source_encoding="utf-8",
            note="目录OCR成功",
            extra_metadata={
                "文档分类": "行业参考架构/指导材料",
                "推荐模板": "参考架构/白皮书口径（当前按原文全量提取Markdown输出）",
                "模板归属": "参考架构/白皮书口径（当前按原文全量提取Markdown输出）",
                "标题": "低空智能网联体系参考架构（2024版）",
                "文件标题": "低空智能网联体系参考架构（2024版）",
                "主体名称": "工业和信息化部装备工业发展中心",
                "单位名称": "工业和信息化部装备工业发展中心",
                "资料层级": "行业参考材料",
                "发布时间": "2024年11月",
                "证据边界": "行业参考架构/白皮书类材料，非正式政策文件；发布口径、适用效力与正式出处仍需结合原始发布来源复核。",
                "来源形态": "分页扫描图片目录",
                "目录判定": "分页扫描文档目录",
                "判定依据": "分页命名图片 3/3；命中封面页",
                "OCR页数": "3",
                "OCR结果概况": "done=3",
                "OCR失败页": "无",
                "取舍说明": "该目录被判定为同一文档的分页扫描目录。",
                "分流结果": "待审核",
                "是否适合直接入库": False,
            },
        )

        with patch("minimum_workflow.pipeline.extract_with_strategy", return_value=extraction), patch(
            "minimum_workflow.pipeline.resolve_qwen_runtime",
            return_value={"api_key": "demo-key", "base_url": "https://demo", "model": "qwen-plus"},
        ), patch("minimum_workflow.pipeline.enrich_payload_with_qwen") as qwen_mock:
            payload, result = build_structured_payload(sample, contract, enable_qwen=True, mineru_token="token-demo")

        qwen_mock.assert_not_called()
        self.assertEqual(result.extractor_name, "ocr:mineru:image_directory")
        self.assertEqual(payload["文件类型"], "image_directory")
        self.assertEqual(payload["文档分类"], "行业参考架构/指导材料")
        self.assertEqual(payload["推荐模板"], "参考架构/白皮书口径（当前按原文全量提取Markdown输出）")
        self.assertEqual(payload["标题"], "低空智能网联体系参考架构（2024版）")
        self.assertEqual(payload["主体名称"], "工业和信息化部装备工业发展中心")
        self.assertEqual(payload["来源形态"], "分页扫描图片目录")
        self.assertEqual(payload["OCR页数"], "3")
        self.assertEqual(payload["分流结果"], "待审核")

    def test_build_structured_payload_skips_qwen_when_runtime_missing(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("supplier_yuefei_2025")
        extraction = ExtractionResult(
            extractor_name="text:utf-8-sig",
            extraction_status="已提取文本",
            extracted_text="公司名称：重庆跃飞智能科技有限公司",
            preview_text="重庆跃飞智能科技有限公司",
            text_length=20,
            page_count=None,
            source_encoding="utf-8-sig",
            note="已完成文本文件读取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction), patch(
            "minimum_workflow.pipeline.resolve_qwen_runtime",
            return_value={},
        ), patch("minimum_workflow.pipeline.enrich_payload_with_qwen") as qwen_mock:
            payload, _ = build_structured_payload(sample, contract, enable_qwen=True)

        qwen_mock.assert_not_called()
        self.assertEqual(payload["一级分类"], "")
        self.assertEqual(payload["二级分类"], "")

    def test_merge_qwen_updates_only_overwrites_classification_fields(self) -> None:
        payload = {
            "推荐模板": "政策官方文件模板",
            "一级分类": "",
            "企业名称": "南京跃飞智能科技有限公司",
            "核心产品": ["原规则产品"],
        }
        qwen_updates = {
            "推荐模板": "供应商企业模板",
            "一级分类": "供应商企业与厂家资料库",
            "企业名称": "被Qwen覆盖的企业名称",
            "核心产品": ["被Qwen覆盖的产品"],
            "联系人姓名": "王立辉",
        }

        merge_qwen_updates(payload, qwen_updates)

        self.assertEqual(payload["推荐模板"], "供应商企业模板")
        self.assertEqual(payload["一级分类"], "供应商企业与厂家资料库")
        self.assertEqual(payload["企业名称"], "南京跃飞智能科技有限公司")
        self.assertEqual(payload["核心产品"], ["原规则产品"])
        self.assertEqual(payload["联系人姓名"], "王立辉")

    def test_markdown_shows_qwen_classification_and_supplemental_fields(self) -> None:
        sample = get_sample_by_id("supplier_yuefei_2025")
        payload = {
            "标题": "跃飞智能画册2025",
            "原始文件名": "跃飞智能画册2025.pdf",
            "原始路径": "D:/长风资料 - 副本/跃飞智能画册2025.pdf",
            "文件格式": "pdf",
            "文档分类": "供应商/企业",
            "推荐模板": "供应商企业模板",
            "一级分类": "供应商企业与厂家资料库",
            "二级分类": "企业画册与产品目录",
            "分类置信度": 0.95,
            "分类依据": "企业官方宣传画册",
            "处理路径": "document_parse",
            "抽取状态": "已提取文本",
            "抽取器": "pypdf",
            "文本长度": 12690,
            "页数": 14,
            "文本编码": "",
            "证据等级": "L2",
            "分流结果": "待审核",
            "内容主题标签": ["厂家主档"],
            "文本预览": "企业摘要",
            "提取正文": "# 第5页\n发展历程\n荣获国家级专精特新小巨人\n# 第6页\n荣誉资质\n标准制定\n国家级标准\n# 第7页\n创始人简介\n朱秋国\n云深处科技创始人&CEO\n# 第10页\n全球布局\n50个国家和地区\n# 第23页\n产品生态\n绝影 X30\n山猫M20\n绝影 Lite3",
            "主营方向": "工业无人机系统生产与行业服务",
            "核心产品": ["MF120侦察无人机"],
            "核心能力": ["研发制造"],
            "企业类别": "厂家",
            "企业名称": "南京跃飞智能科技有限公司",
            "联系人姓名": "王立辉",
            "联系人角色字段": "首席技术专家/教授",
            "联系方式字段": "网址: www.tendrones.com",
            "去重主键": ["跃飞智能"],
            "是否需要拆分": False,
            "拆分说明": "",
            "抽取说明": "已尝试使用 pypdf 提取。",
            "是否适合直接入库": False,
            "风险说明": ["画册属性明显"],
            "备注": ["用于校准厂家主档到产品族索引的最小输出。"],
        }

        markdown = build_markdown(sample, payload)

        self.assertIn("## 一、企业摘要", markdown)
        self.assertIn("企业摘要", markdown)
        self.assertIn("- 主营方向：工业无人机系统生产与行业服务", markdown)
        self.assertIn("- 核心产品：MF120侦察无人机", markdown)
        self.assertIn("- 核心能力：研发制造", markdown)
        self.assertIn("- 企业名称：南京跃飞智能科技有限公司", markdown)
        self.assertIn("荣获国家级专精特新小巨人", markdown)
        self.assertIn("## 六、荣誉资质与标准线索", markdown)
        self.assertIn("国家级标准", markdown)
        self.assertIn("## 七、创始人与团队线索", markdown)
        self.assertIn("云深处科技创始人&CEO", markdown)
        self.assertIn("## 八、分支机构与布局覆盖", markdown)
        self.assertIn("50个国家和地区", markdown)
        self.assertIn("## 九、产品生态与应用线索", markdown)
        self.assertIn("绝影 X30", markdown)
        self.assertIn("## 原文全文", markdown)
        self.assertIn("# 第5页\n发展历程", markdown)

    def test_supplier_template_markdown_omits_empty_fields_and_noisy_preview(self) -> None:
        sample = get_sample_by_id("supplier_yuefei_2025")
        payload = {
            "标题": "十大低空经济企业",
            "原始文件名": "十大低空经济企业.docx",
            "原始路径": "D:/长风资料/11、行业内公司/十大低空经济企业.docx",
            "文件格式": "docx",
            "文档分类": "供应商/企业资料",
            "推荐模板": "供应商企业模板",
            "一级分类": "供应商企业与厂家资料库",
            "二级分类": "企业名录/行业盘点",
            "分类置信度": 0.78,
            "分类依据": "正文以行业企业名单盘点为主",
            "处理路径": "document_parse",
            "抽取状态": "已提取文本",
            "抽取器": "word:docx",
            "文本长度": 3200,
            "页数": 18,
            "文本编码": "utf-8",
            "证据等级": "L3",
            "分流结果": "待审核",
            "内容主题标签": ["目录扫描", "自动判型"],
            "文本预览": "大低空整机制造企业：深圳市大疆创新科技有限公司、深圳市道通智能航空技术股份有限公司、乐山天穹动力科技有限公司、成都纵横自动化技术股份有限公司、普宙科技有限公司、深圳市科卫泰实业发展有限公司、哈瓦国际航空技术（深圳）有限公司、深圳飞马机器人股份有限公司。",
            "提取正文": "# 十大低空经济企业\n\n十大低空整机制造企业\n\n十大低空物流企业",
            "企业类别": "厂家",
            "企业名称": "",
            "主营方向": "",
            "核心产品": [],
            "核心能力": [],
            "去重主键": ["十大低空经济企业"],
            "是否需要拆分": False,
            "拆分说明": "",
            "抽取说明": "已完成 docx 正文提取。",
            "是否适合直接入库": False,
            "风险说明": ["目录扫描自动生成样本，字段需结合原文复核。"],
            "备注": [],
        }

        markdown = build_markdown(sample, payload)

        self.assertNotIn("大低空整机制造企业：", markdown)
        self.assertNotIn("- 主营方向：未提取", markdown)
        self.assertNotIn("- 核心能力：未提取", markdown)
        self.assertNotIn("- 核心产品：未提取", markdown)
        self.assertNotIn("- 企业名称：", markdown)
        self.assertNotIn("## 五、发展历程与关键节点", markdown)
        self.assertNotIn("## 六、荣誉资质与标准线索", markdown)
        self.assertIn("- 企业类别：厂家", markdown)
        self.assertIn("## 原文全文", markdown)

    def test_education_training_template_markdown_contains_dedicated_sections(self) -> None:
        sample = SampleRecord(
            sample_id="scan_education_training_markdown",
            source_path="D:/长风资料/长风/教育培训方向/培训方案.docx",
            document_category="教育培训",
            recommended_template="教育培训模板",
            title_hint="低空教育培训方案",
            subject_name_hint="教育培训方向",
            product_name_hint="",
            unit_name_hint="长风职业培训中心",
            tags=["教育培训"],
            risks=[],
            notes=[],
            evidence_level="L2",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        payload = {
            "标题": "低空教育培训方案",
            "原始文件名": "培训方案.docx",
            "原始路径": "D:/长风资料/长风/教育培训方向/培训方案.docx",
            "文件格式": "docx",
            "文档分类": "教育培训",
            "推荐模板": "教育培训模板",
            "一级分类": "内部方法流程与话术库",
            "二级分类": "教育培训资料",
            "分类置信度": 0.88,
            "分类依据": "命中培训方案、课程体系、培训对象等关键词",
            "处理路径": "document_parse",
            "抽取状态": "已提取文本",
            "抽取器": "word:docx",
            "文本长度": 850,
            "页数": 6,
            "文本编码": "utf-8",
            "证据等级": "L2",
            "分流结果": "待审核",
            "内容主题标签": ["教育培训", "课程体系"],
            "文本预览": "低空教育培训方案摘要",
            "提取正文": "# 第33页\n产教融合\n打造集教学、实训、科研、竞赛于一体的综合性平台\n# 第38页\n配套产品介绍\n绝影Lite3系列实训台\n# 第41页\n案例介绍\n云深处科技-西科大校园巡检案例",
            "文件标题": "低空教育培训方案",
            "单位名称字段": "长风职业培训中心",
            "培训主题字段": "低空人才培养",
            "适用对象字段": "职业院校学生",
            "培训类型字段": "培训方案",
            "专业方向字段": "无人机应用技术",
            "课程体系字段": "无人机基础；飞行控制；行业实训",
            "实施方式字段": "理论授课+实训操作",
            "核心内容字段": "围绕低空应用开展职业教育培训",
            "去重主键": ["低空教育培训方案"],
            "是否需要拆分": False,
            "拆分说明": "",
            "抽取说明": "已完成 docx 正文提取。",
            "是否适合直接入库": False,
            "风险说明": [],
            "备注": [],
        }

        markdown = build_markdown(sample, payload)

        self.assertIn("## 二、培训主题", markdown)
        self.assertIn("## 五、专业方向/课程体系", markdown)
        self.assertIn("## 八、字段提取结果", markdown)
        self.assertIn("## 九、课程体系与合作模式原文", markdown)
        self.assertIn("## 原文全文", markdown)
        self.assertIn("# 第33页\n产教融合", markdown)
        self.assertIn("绝影Lite3系列实训台", markdown)
        self.assertIn("云深处科技-西科大校园巡检案例", markdown)
        self.assertNotIn("## 抽取说明", markdown)

    def test_education_training_template_skips_empty_placeholder_sections(self) -> None:
        sample = SampleRecord(
            sample_id="scan_education_training_empty_fields",
            source_path="D:/长风资料/11、行业内公司/北京大学低空经济与区域经济发展高级研修班（三期）.pdf",
            document_category="教育培训",
            recommended_template="教育培训模板",
            title_hint="北京大学低空经济与区域经济发展高级研修班（三期）",
            subject_name_hint="北京大学低空经济与区域经济发展高级研修班（三期）",
            product_name_hint="",
            unit_name_hint="北京大学",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L3",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        payload = {
            "标题": "北京大学低空经济与区域经济发展高级研修班（三期）",
            "原始文件名": "北京大学低空经济与区域经济发展高级研修班（三期）.pdf",
            "原始路径": "D:/长风资料/11、行业内公司/北京大学低空经济与区域经济发展高级研修班（三期）.pdf",
            "文件格式": "pdf",
            "文档分类": "教育培训",
            "推荐模板": "教育培训模板",
            "一级分类": "未识别",
            "二级分类": "未识别",
            "分类置信度": "未识别",
            "分类依据": "未提取",
            "处理路径": "document_parse",
            "抽取状态": "已提取文本",
            "抽取器": "mineru:batch",
            "文本长度": 2861,
            "页数": "未识别",
            "文本编码": "utf-8",
            "证据等级": "L3",
            "分流结果": "待审核",
            "内容主题标签": ["目录扫描", "自动判型"],
            "文本预览": "北京大学低空经济与区域经济发展高级研修班（三期）摘要",
            "提取正文": "# 北京大学\n# 低空经济与区域经济发展\n# 高级研修班（三期）",
            "文件标题": "北京大学低空经济与区域经济发展高级研修班（三期）",
            "单位名称字段": "北京大学",
            "培训主题字段": "北京大学低空经济与区域经济发展高级研修班（三期）",
            "适用对象字段": "从事低空经济相关产业的人员",
            "培训类型字段": "教育培训资料",
            "专业方向字段": "低空经济；区域经济",
            "课程体系字段": "发展路径：低空经济政策解读及重点省份发展落实情况",
            "实施方式字段": "",
            "核心内容字段": "未来低空经济产业人才培养模式探索",
            "去重主键": ["北京大学低空经济与区域经济发展高级研修班（三期）"],
            "是否需要拆分": False,
            "拆分说明": "",
            "抽取说明": "已完成 PDF 正文提取。",
            "是否适合直接入库": False,
            "风险说明": [],
            "备注": [],
        }

        markdown = build_markdown(sample, payload)

        self.assertNotIn("## 六、实施方式", markdown)
        self.assertNotIn("- 实施方式：", markdown)
        self.assertIn("## 五、专业方向/课程体系", markdown)
        self.assertIn("- 专业方向：低空经济；区域经济", markdown)
        self.assertIn("- 课程体系：发展路径：低空经济政策解读及重点省份发展落实情况", markdown)
        self.assertIn("## 七、核心内容", markdown)
        self.assertIn("## 八、字段提取结果", markdown)
        contract = load_contract()
        sample = get_sample_by_id("supplier_yuefei_2025")
        extraction = ExtractionResult(
            extractor_name="text:utf-8-sig",
            extraction_status="已提取文本",
            extracted_text="公司名称：重庆跃飞智能科技有限公司",
            preview_text="重庆跃飞智能科技有限公司",
            text_length=20,
            page_count=None,
            source_encoding="utf-8-sig",
            note="已完成文本文件读取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction), patch(
            "minimum_workflow.pipeline.resolve_qwen_runtime",
            return_value={"api_key": "demo-key", "base_url": "https://demo", "model": "qwen-plus"},
        ), patch(
            "minimum_workflow.pipeline.enrich_payload_with_qwen",
            side_effect=RuntimeError("mock qwen error"),
        ):
            payload, _ = build_structured_payload(sample, contract, enable_qwen=True)

        self.assertEqual(payload["企业名称"], "重庆跃飞智能科技有限公司")
        self.assertIn("Qwen补强未生效：mock qwen error", payload["抽取说明"])

    def test_enrich_payload_with_qwen_accepts_list_fields(self) -> None:
        sample = get_sample_by_id("supplier_yuefei_2025")
        extraction = ExtractionResult(
            extractor_name="text:utf-8-sig",
            extraction_status="已提取文本",
            extracted_text="公司名称：重庆跃飞智能科技有限公司\n核心产品：无人机平台、智能机库",
            preview_text="重庆跃飞智能科技有限公司",
            text_length=40,
            page_count=None,
            source_encoding="utf-8-sig",
            note="已完成文本文件读取。",
        )
        payload = {"文档分类": "供应商/企业资料", "推荐模板": "供应商企业模板"}

        with patch(
            "minimum_workflow.qwen_client.request_qwen_classification_and_fields",
            return_value={
                "一级分类": "供应商企业与厂家资料库",
                "二级分类": "无人机厂家",
                "分类置信度": 0.91,
                "分类依据": "正文出现企业名称和产品列表",
                "字段": {
                    "核心产品": ["无人机平台", "智能机库"],
                    "企业名称": "重庆跃飞智能科技有限公司",
                },
            },
        ):
            updates = enrich_payload_with_qwen(
                sample,
                extraction,
                payload,
                api_key="demo-key",
                base_url="https://demo",
                model="qwen-plus",
            )

        self.assertEqual(updates["一级分类"], "供应商企业与厂家资料库")
        self.assertEqual(updates["二级分类"], "无人机厂家")
        self.assertEqual(updates["分类置信度"], 0.91)
        self.assertEqual(updates["核心产品"], ["无人机平台", "智能机库"])
        self.assertEqual(updates["企业名称"], "重庆跃飞智能科技有限公司")

    def test_extract_with_strategy_uses_mineru_for_pdf(self) -> None:
        sample_path = Path("D:/tmp/sample.pdf")
        expected = ExtractionResult(
            extractor_name="mineru:batch",
            extraction_status="已提取文本",
            extracted_text="markdown body",
            preview_text="markdown body",
            text_length=13,
            page_count=None,
            source_encoding="utf-8",
            note="已通过 MinerU 完成提取。",
        )
        with patch("minimum_workflow.pipeline.extract_pdf_with_mineru", return_value=expected) as mineru_mock:
            result = extract_with_strategy(sample_path, "pdf", pdf_extractor="mineru", mineru_token="token-demo")

        mineru_mock.assert_called_once_with(sample_path, "token-demo")
        self.assertEqual(result.extractor_name, "mineru:batch")

    def test_extract_with_strategy_falls_back_to_local_when_mineru_token_missing(self) -> None:
        expected = ExtractionResult(
            extractor_name="pdf:pymupdf",
            extraction_status="已提取文本",
            extracted_text="local body",
            preview_text="local body",
            text_length=10,
            page_count=2,
            source_encoding="utf-8",
            note="本地 PDF 解析成功。",
        )
        with patch("minimum_workflow.pipeline.extract_text", return_value=expected) as extract_text_mock:
            result = extract_with_strategy(Path("D:/tmp/sample.pdf"), "pdf", pdf_extractor="mineru", mineru_token=None)

        extract_text_mock.assert_called_once()
        self.assertEqual(result.extractor_name, "pdf:pymupdf")
        self.assertIn("未提供 MinerU token，已回退本地解析", result.note)

    def test_extract_with_strategy_falls_back_to_local_when_mineru_fails(self) -> None:
        expected = ExtractionResult(
            extractor_name="pdf:pymupdf",
            extraction_status="已提取文本",
            extracted_text="local body",
            preview_text="local body",
            text_length=10,
            page_count=2,
            source_encoding="utf-8",
            note="本地 PDF 解析成功。",
        )
        with patch("minimum_workflow.pipeline.extract_pdf_with_mineru", side_effect=RuntimeError("mineru failed")), patch(
            "minimum_workflow.pipeline.extract_text", return_value=expected
        ) as extract_text_mock:
            result = extract_with_strategy(Path("D:/tmp/sample.pdf"), "pdf", pdf_extractor="mineru", mineru_token="token-demo")

        extract_text_mock.assert_called_once()
        self.assertEqual(result.extractor_name, "pdf:pymupdf")
        self.assertIn("MinerU 优先抽取失败，已回退本地解析：mineru failed", result.note)

    def test_extract_with_strategy_supports_directory_source(self) -> None:
        source_dir = Path("D:/tmp/paged-dir")
        directory_result = {
            "blocks": ["# 标题", "正文第一页", "正文第二页"],
            "extracted_text": "# 标题\n\n正文第一页\n\n正文第二页",
            "extraction_note": "目录OCR成功",
            "auto_payload": {
                "文档分类": "行业参考架构/指导材料",
                "推荐模板": "参考架构/白皮书口径（当前按原文全量提取Markdown输出）",
                "标题": "低空智能网联体系参考架构（2024版）",
                "主体名称": "工业和信息化部装备工业发展中心",
                "单位名称": "工业和信息化部装备工业发展中心",
                "发布时间": "2024年11月",
                "来源形态": "分页扫描图片目录",
                "目录判定": "分页扫描文档目录",
                "判定依据": "分页命名图片 3/3；命中封面页",
                "OCR页数": "3",
                "OCR结果概况": "done=3",
                "OCR失败页": "无",
                "分流结果": "待审核",
                "是否适合直接入库": False,
            },
        }

        with patch("minimum_workflow.pipeline.extract_image_directory_content", return_value=directory_result):
            result = extract_with_strategy(source_dir, DIRECTORY_TYPE, mineru_token="token-demo")

        self.assertEqual(result.extractor_name, "ocr:mineru:image_directory")
        self.assertEqual(result.extraction_status, "已提取文本")
        self.assertEqual(result.page_count, 3)
        self.assertEqual(result.preview_text, "# 标题  正文第一页  正文第二页")
        self.assertEqual(result.extra_metadata["来源形态"], "分页扫描图片目录")

    def test_post_mineru_batch_upload_urls_retries_on_429(self) -> None:
        first_response = Mock()
        first_response.status_code = 429
        first_response.headers = {"Retry-After": "2"}
        first_response.raise_for_status.side_effect = requests.HTTPError("429 Client Error")

        second_response = Mock()
        second_response.status_code = 200
        second_response.headers = {}

        session = Mock()
        session.post.side_effect = [first_response, second_response]

        with patch("minimum_workflow.extractors.time.sleep") as sleep_mock:
            response = extractors_module.post_mineru_batch_upload_urls(
                session,
                "https://mineru.net/api/v4",
                {"Authorization": "Bearer token-demo"},
                {"files": [{"name": "sample.pdf", "data_id": "file_1"}]},
            )

        self.assertIs(response, second_response)
        self.assertEqual(session.post.call_count, 2)
        sleep_mock.assert_called_once_with(2)

    def test_post_mineru_batch_upload_urls_raises_after_retry_limit(self) -> None:
        response = Mock()
        response.status_code = 429
        response.headers = {}
        response.raise_for_status.side_effect = requests.HTTPError("429 Client Error")

        session = Mock()
        session.post.return_value = response

        with patch("minimum_workflow.extractors.time.sleep") as sleep_mock:
            with self.assertRaises(requests.HTTPError):
                extractors_module.post_mineru_batch_upload_urls(
                    session,
                    "https://mineru.net/api/v4",
                    {"Authorization": "Bearer token-demo"},
                    {"files": [{"name": "sample.pdf", "data_id": "file_1"}]},
                )

        self.assertEqual(session.post.call_count, 4)
        self.assertEqual(sleep_mock.call_count, 3)


    def test_extract_pdf_with_mineru_returns_markdown_result(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        sample_path = Path(temp_dir.name) / "sample.pdf"
        sample_path.write_bytes(b"%PDF-1.4\n")

        with patch(
            "minimum_workflow.extractors.run_mineru_batch",
            return_value={
                "batch_id": "batch-002",
                "results": [
                    {
                        "file_path": str(sample_path),
                        "state": "done",
                        "markdown": "# 标题\n\n![](images/a.jpg)\n<div>这是一份 MinerU 提取结果，包含足够的正文内容用于进入后续链路。</div>",
                        "zip_url": "https://download.example/full.zip",
                        "error": "",
                    }
                ],
            },
        ):
            result = extract_pdf_with_mineru(sample_path, "token-demo")

        self.assertEqual(result.extractor_name, "mineru:batch")
        self.assertEqual(result.extraction_status, "已提取文本")
        self.assertNotIn("![](", result.extracted_text)
        self.assertNotIn("<div>", result.extracted_text)
        self.assertIn("MinerU 提取结果", result.extracted_text)
        self.assertIn("批次号：batch-002", result.note)

    def test_policy_template_markdown_contains_policy_sections(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("policy_chongqing_txt_review")
        extraction = ExtractionResult(
            extractor_name="text:utf-8-sig",
            extraction_status="已提取文本",
            extracted_text="重庆市推动低空经济高质量发展若干政策措施",
            preview_text="重庆市推动低空经济高质量发展若干政策措施 支持创新打造低空示范场景",
            text_length=50,
            page_count=None,
            source_encoding="utf-8-sig",
            note="已完成文本文件读取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)
            markdown = build_markdown(sample, payload)

        self.assertIn("## 一、文件摘要", markdown)
        self.assertIn("## 二、核心要求", markdown)
        self.assertIn("## 三、与低空/长风业务相关的部分", markdown)
        self.assertIn("## 四、执行或应用意义", markdown)
        self.assertIn("## 五、时效与边界", markdown)
        self.assertIn("## 六、字段提取结果", markdown)

    def test_policy_fields_are_extracted_into_payload(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("policy_chongqing_txt_review")
        extraction = ExtractionResult(
            extractor_name="text:utf-8-sig",
            extraction_status="已提取文本",
            extracted_text=(
                "文件名称：《重庆市推动低空经济高质量发展若干政策措施》\n"
                "发文字号：渝府办发〔2025〕58号\n"
                "成文日期：2025年11月22日\n"
                "发文单位：重庆市人民政府办公厅\n"
                "一、支持创新打造低空示范场景\n"
                "二、支持低空产业创新发展\n"
                "本政策措施自印发之日起施行，有效期至2027年12月31日止。"
            ),
            preview_text="重庆市推动低空经济高质量发展若干政策措施 支持创新打造低空示范场景",
            text_length=120,
            page_count=None,
            source_encoding="utf-8-sig",
            note="已完成文本文件读取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)

        self.assertEqual(payload["文件标题"], "重庆市推动低空经济高质量发展若干政策措施")
        self.assertEqual(payload["发文字号"], "渝府办发〔2025〕58号")
        self.assertEqual(payload["发文单位"], "重庆市人民政府办公厅")
        self.assertEqual(payload["成文日期"], "2025-11-22")
        self.assertEqual(payload["生效状态"], "已生效")
        self.assertIn("一、支持创新打造低空示范场景", payload["核心任务"])

    def test_policy_mineru_fields_are_extracted_into_payload(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("policy_uav_regulation")
        extraction = ExtractionResult(
            extractor_name="mineru:batch",
            extraction_status="已提取文本",
            extracted_text=(
                "# 中华人民共和国国务院中华人民共和国中央军事委员会\n\n"
                "第761号\n\n"
                "现公布《无人驾驶航空器飞行管理暂行条例》，自2024年1月1日起施行。\n\n"
                "中央军委主席 习近平 国务院总理 李强\n\n"
                "2023年5月31日\n\n"
                "无人驾驶航空器飞行管理暂行条例\n\n"
                "# 第一章 总 则\n\n"
                "# 第二章 民用无人驾驶航空器及操控员管理\n"
            ),
            preview_text="无人驾驶航空器飞行管理暂行条例",
            text_length=200,
            page_count=None,
            source_encoding="utf-8",
            note="已通过 MinerU 批量接口完成 Markdown 提取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)

        self.assertEqual(payload["文件标题"], "无人驾驶航空器飞行管理暂行条例")
        self.assertEqual(payload["发文字号"], "国令第761号")
        self.assertEqual(payload["发文单位"], "国务院、中央军委")
        self.assertEqual(payload["成文日期"], "2023-05-31")
        self.assertEqual(payload["生效状态"], "已生效")
        self.assertIn("第一章 总 则", payload["核心任务"][0])

    def test_policy_title_does_not_use_procurement_qualification_clause(self) -> None:
        sample = get_sample_by_id("policy_uav_regulation")
        extraction = ExtractionResult(
            extractor_name="mineru:batch",
            extraction_status="已提取文本",
            extracted_text=(
                "第一章 采购公告\n"
                "1.2符合《中华人民共和国政府采购法》及其实施条例有关规定的投标供应商资格条件；\n"
            ),
            preview_text="第一章 采购公告",
            text_length=70,
            page_count=None,
            source_encoding="utf-8",
            note="已通过 MinerU 批量接口完成 Markdown 提取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["文件标题"], sample.title_hint)

    def test_supplier_fields_are_extracted_into_payload(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("supplier_yuefei_2025")
        extraction = ExtractionResult(
            extractor_name="text:utf-8-sig",
            extraction_status="已提取文本",
            extracted_text=(
                "公司名称：重庆跃飞智能科技有限公司\n"
                "主营业务：低空飞行器研发、制造、行业应用解决方案\n"
                "核心产品：无人机平台、智能机库、巡检系统\n"
                "核心能力：研发制造、行业集成、运营服务\n"
            ),
            preview_text="重庆跃飞智能科技有限公司 低空飞行器研发制造",
            text_length=80,
            page_count=None,
            source_encoding="utf-8-sig",
            note="已完成文本文件读取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)
            markdown = build_markdown(sample, payload)

        self.assertEqual(payload["企业名称"], "重庆跃飞智能科技有限公司")
        self.assertEqual(payload["企业类别"], "厂家")
        self.assertEqual(payload["主营方向"], "低空飞行器研发、制造、行业应用解决方案")
        self.assertIn("无人机平台", payload["核心产品"])
        self.assertIn("研发制造", payload["核心能力"])
        self.assertIn("## 四、字段提取结果", markdown)

    def test_supplier_brochure_fields_capture_product_lines_and_key_products(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("nw_supplier_huanuoxingkong")
        extraction = ExtractionResult(
            extractor_name="pypdf",
            extraction_status="已提取文本",
            extracted_text=(
                "华诺星空技术股份有限公司\n"
                "公司已形成以全域安防、 安全应急、 智能感知、 无人系统、 智慧交通等五大产品线为主的产品族群。\n"
                "公司拥有雷达、 光电、 人工智能、 多源感知融合等核心技术。\n"
                "华诺星空低空防御体系\n"
                "雷达探测设备 SC-R3000 | SC-R5000\n"
                "雷达探测设备 SC-R8000\n"
                "无人机探测定位一体设备 SC-P5000+\n"
                "分布式侦测定位系统\n"
                "箱组式反无人机系统\n"
                "车载式反无人机系统\n"
            ),
            preview_text="华诺星空技术股份有限公司 华诺星空低空防御体系",
            text_length=180,
            page_count=10,
            source_encoding="",
            note="已尝试使用 pypdf 提取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)

        self.assertEqual(payload["企业名称"], "华诺星空技术股份有限公司")
        self.assertEqual(payload["主营方向"], "全域安防；安全应急；智能感知；无人系统；智慧交通")
        self.assertIn("雷达探测设备 SC-R3000 | SC-R5000", payload["核心产品"])
        self.assertIn("分布式侦测定位系统", payload["核心产品"])
        self.assertNotIn("雷达探测设备", payload["核心产品"])
        self.assertTrue(all("�" not in item for item in payload["核心产品"]))
        self.assertIn("平台与算法", payload["核心能力"])
        self.assertIn("智能感知与低空防御", payload["核心能力"])

    def test_supplier_group_brochure_prefers_group_name_and_uav_products(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("supplier_united_aircraft_brochure")
        extraction = ExtractionResult(
            extractor_name="pypdf",
            extraction_status="已提取文本",
            extracted_text=(
                "联合飞机集团\n"
                "北京中航智科技有限公司\n"
                "联合飞机集团（简称联合飞机）是专业从事无人机等高端装备研发、生产和服务的企业集团。\n"
                "镧影R6000倾转旋翼飞行器\n"
                "TD550无人直升机\n"
                "Q100农业无人机\n"
                "共轴带尾推高速无人直升机\n"
                "自动飞行控制系统\n"
            ),
            preview_text="联合飞机集团 工业级无人机产品",
            text_length=120,
            page_count=8,
            source_encoding="",
            note="已尝试使用 pypdf 提取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)

        self.assertEqual(payload["企业名称"], "联合飞机集团")
        self.assertEqual(payload["主营方向"], "工业级无人机研发制造与行业应用服务")
        self.assertIn("镧影R6000倾转旋翼飞行器", payload["核心产品"])
        self.assertIn("TD550无人直升机", payload["核心产品"])
        self.assertIn("Q100农业无人机", payload["核心产品"])

    def test_contact_fields_are_extracted_into_payload(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("policy_chongqing_txt_review")
        sample = sample.__class__(
            sample_id="contact_demo",
            source_path=sample.source_path,
            document_category="单位/联系人",
            recommended_template="单位联系人模板",
            title_hint="重庆低空项目联系人",
            subject_name_hint="重庆低空项目联系人",
            product_name_hint="",
            unit_name_hint="重庆市人民政府办公厅",
            tags=sample.tags,
            risks=sample.risks,
            notes=sample.notes,
            evidence_level=sample.evidence_level,
            fallback_decision=sample.fallback_decision,
            split_required=sample.split_required,
            split_note=sample.split_note,
        )
        extraction = ExtractionResult(
            extractor_name="text:utf-8-sig",
            extraction_status="已提取文本",
            extracted_text=(
                "单位名称：重庆市人民政府办公厅\n"
                "联系人：张三\n"
                "联系人角色：项目经办人\n"
                "联系方式：13800000000\n"
            ),
            preview_text="重庆市人民政府办公厅 联系人张三",
            text_length=50,
            page_count=None,
            source_encoding="utf-8-sig",
            note="已完成文本文件读取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)
            markdown = build_markdown(sample, payload)

        self.assertEqual(payload["单位名称字段"], "重庆市人民政府办公厅")
        self.assertEqual(payload["单位类型字段"], "政府部门")
        self.assertEqual(payload["联系人姓名"], "张三")
        self.assertEqual(payload["联系人角色字段"], "项目经办人")
        self.assertEqual(payload["联系方式字段"], "13800000000")
        self.assertEqual(payload["对接线索字段"], "")
        self.assertIn("## 三、字段提取结果", markdown)
        self.assertIn("- 对接线索：未提取", markdown)

    def test_contact_sample_keeps_only_unit_and_interface_clues(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("contact_chongqing_public_data")
        extraction = ExtractionResult(
            extractor_name="text:utf-8-sig",
            extraction_status="已提取文本",
            extracted_text=(
                "重庆低空经济商机信息获取工具\n"
                "长风低空经济产业发展（深圳）有限公司\n"
                "https://localhost\n"
                "接口地址 https://data.cq.gov.cn/interface/athena/api/getSign\n"
                "支持格式 JSON\n"
                "请求方式 GET/POST\n"
                "接口描述 获取sign\n"
                "docker run -d --name n8n ...\n"
            ),
            preview_text="重庆低空经济商机信息获取工具 长风低空经济产业发展（深圳）有限公司 https://localhost",
            text_length=120,
            page_count=None,
            source_encoding="utf-8-sig",
            note="已完成文本文件读取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)
            markdown = build_markdown(sample, payload)

        self.assertEqual(payload["单位名称字段"], "重庆市相关公共数据开放接口提供方")
        self.assertEqual(payload["单位类型字段"], "政府数据接口/公共数据平台相关单位")
        self.assertEqual(payload["联系人姓名"], "")
        self.assertEqual(
            payload["对接线索字段"],
            "接口地址：https://data.cq.gov.cn/interface/athena/api/getSign；支持格式：JSON；请求方式：GET/POST；接口用途：获取sign",
        )
        self.assertIn("- 对接线索：接口地址：https://data.cq.gov.cn/interface/athena/api/getSign", markdown)
        self.assertNotIn("https://localhost", markdown)

    def test_solution_fields_are_extracted_into_payload(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("solution_zhongshan_medical")
        extraction = ExtractionResult(
            extractor_name="pypdf",
            extraction_status="已提取文本",
            extracted_text=(
                "区域检验低空转运方案\n"
                "客户/使用单位：钟山区医共体\n"
                "2025年\n"
                "低空医疗应急绿色生命通道\n"
                "禁飞区查询\n"
                "航线规划\n"
                "急救物流无人机\n"
                "自动化急救枢纽站\n"
                "无人机管理调度云平台\n"
                "成本预算\n"
                "时效提升\n"
                "分级诊疗\n"
                "无人机飞行不受地面环境制约，可用于医疗物资的运输，避免二次污染。\n"
            ),
            preview_text="区域检验低空转运方案 医疗应急绿色生命通道",
            text_length=180,
            page_count=10,
            source_encoding="",
            note="已尝试使用 pypdf 提取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)
            markdown = build_markdown(sample, payload)

        self.assertEqual(payload["方案名称字段"], "区域检验低空转运方案")
        self.assertEqual(payload["所属场景字段"], "医疗运输")
        self.assertEqual(payload["客户名称字段"], "钟山区医共体")
        self.assertEqual(payload["文件日期字段"], "2025年")
        self.assertIn("地面运输慢", payload["解决问题字段"])
        self.assertIn("自动化急救枢纽站", payload["产品能力字段"])
        self.assertEqual(payload["实施方式字段"], "先做禁飞区查询与航线规划，再执行运输、换电与运行保障")
        self.assertEqual(payload["预算组织字段"], "含预算/报价内容，需与方案正文拆层处理")
        self.assertEqual(payload["效果数据字段"], "无人机运输相较地面运输具备明显时效提升")
        self.assertIn("医疗检验转运", payload["可复用经验字段"])
        self.assertEqual(payload["证据类型字段"], "混合型资料")
        self.assertIn("## 九、字段提取结果", markdown)

    def test_solution_fields_ignore_company_intro_noise_for_yunshenchu_solution(self) -> None:
        contract = load_contract()
        sample = SampleRecord(
            sample_id="scan_yunshenchu_solution_noise",
            source_path="D:/长风资料/云深处公司产品资料/【云深处】企业园区巡检解决方案 20251017_V1.0.pptx",
            document_category="方案/案例",
            recommended_template="方案案例模板",
            title_hint="企业园区巡检解决方案",
            subject_name_hint="企业园区巡检解决方案",
            product_name_hint="",
            unit_name_hint="公司总部",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L3",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="presentation:pptx",
            extraction_status="已提取文本",
            extracted_text=(
                "企业园区巡检解决方案\n"
                "公司介绍\n"
                "杭州云深处科技股份有限公司\n"
                "公司总部\n"
                "应用赋能\n"
                "学校、景区、医院场景\n"
                "智能巡检机器人\n"
            ),
            preview_text="企业园区巡检解决方案",
            text_length=80,
            page_count=12,
            source_encoding="pptx",
            note="已完成演示文稿读取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)

        self.assertEqual(payload["方案名称字段"], "企业园区巡检解决方案")
        self.assertEqual(payload["所属场景字段"], "")
        self.assertEqual(payload["客户名称字段"], "")

    def test_solution_fields_ignore_verb_and_base_noise_for_customer(self) -> None:
        sample = SampleRecord(
            sample_id="scan_solution_customer_noise_only",
            source_path="D:/长风资料/云深处公司产品资料/公安解决方案20260120-V1.0.pptx",
            document_category="方案/案例",
            recommended_template="方案案例模板",
            title_hint="公安解决方案",
            subject_name_hint="公安解决方案",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L3",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="presentation:pptx",
            extraction_status="已提取文本",
            extracted_text=(
                "并在发现异常情况时及时向指挥中心推送信息\n"
                "云深处科技具身智能中试基地\n"
                "帮助指挥中心全面掌握现场风险状况\n"
            ),
            preview_text="公安解决方案",
            text_length=60,
            page_count=10,
            source_encoding="pptx",
            note="已完成演示文稿读取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["客户名称字段"], "")
        self.assertEqual(payload["所属场景字段"], "")

    def test_solution_template_markdown_contains_solution_sections(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("scenario_heilongjiang_2024")
        extraction = ExtractionResult(
            extractor_name="pypdf",
            extraction_status="已提取文本",
            extracted_text=(
                "2024年黑龙江省低空经济发展重点项目和应用场景清单\n"
                "2024年\n"
                "一、哈尔滨联合飞机大型无人机产业基地项目\n"
                "总投资 25 亿元\n"
                "八、自然灾害应急能力提升工程航空应急项目\n"
                "建设 5 套中型复合翼无人机救援平台\n"
                "打造 3 小时无人机服务圈\n"
            ),
            preview_text="黑龙江省低空经济重点项目和应用场景清单",
            text_length=120,
            page_count=3,
            source_encoding="",
            note="已尝试使用 pypdf 提取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)
            markdown = build_markdown(sample, payload)

        self.assertEqual(payload["方案名称字段"], "2024年黑龙江省低空经济发展重点项目和应用场景清单")
        self.assertEqual(payload["所属场景字段"], "低空应用场景")
        self.assertEqual(payload["文件日期字段"], "2024年")
        self.assertEqual(payload["预算组织字段"], "总投资 25 亿元")
        self.assertIn("中型复合翼无人机救援平台", payload["产品能力字段"])
        self.assertEqual(payload["效果数据字段"], "目标打造 3 小时无人机服务圈，提升区域航空应急能力")
        self.assertEqual(payload["可复用经验字段"], "可复用于区域项目机会清单、场景条目拆分和时间边界标注")
        self.assertEqual(payload["证据类型字段"], "清单型资料")
        self.assertIn("## 三、资料形态判断", markdown)
        self.assertIn("## 九、字段提取结果", markdown)
        self.assertIn("## 原文全文", markdown)

    def test_solution_scene_classification_xlsx_uses_catalog_guards(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("cf_solution_scene_classification_xlsx")
        extraction = ExtractionResult(
            extractor_name="excel:openpyxl",
            extraction_status="已提取文本",
            extracted_text=(
                "# 工作表：Sheet1\n\n"
                "| 长风低空---无人机应用场景分类 | 列2 | 列3 | 列4 | 列5 | 列6 | 列7 | 列8 | 列9 | 列10 | 列11 | 列12 |\n"
                "| --- | --- | --- | --- | --- | --- | --- | --- | --- | --- | --- | --- |\n"
                "| 序号 | 应用场景 | 场景简述 | 设备选型 | 设备参数 | 设备特点 | 制造商、产地、官网 | 应用收费 | 产品市场价 | 设备图片 | - | - |\n"
                "| 一、航拍摄影 | - | - | - | - | - | - | - | - | - | - | - |\n"
                "| 1 | 风光与自然景观拍摄 | 山川、湖泊、沙漠、森林等自然地貌的全景拍摄。 | 大疆Mavic 3 Pro | 最长飞行时间：43分钟 | 多焦段影像系统 | 深圳市大疆创新科技有限公司 | 单日租赁价格通常在100-518元 | 参考价格为13888元-31688元 | - | - | - |\n"
            ),
            preview_text="长风低空无人机场景应用分类",
            text_length=260,
            page_count=1,
            source_encoding="xlsx",
            note="已完成 Excel 读取，共识别 1 个有效工作表。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)

        self.assertEqual(payload["方案名称字段"], "长风低空[无人机场景应用分类]")
        self.assertEqual(payload["所属场景字段"], "低空应用场景")
        self.assertEqual(payload["客户名称字段"], "长风低空")
        self.assertEqual(payload["证据类型字段"], "清单型资料")
        self.assertNotIn("Parrot Bluegrass Fields", payload["方案名称字段"])
        self.assertNotIn("|", payload["方案名称字段"])

    def test_solution_recent_projects_xlsx_avoids_medical_false_hit(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("cf_solution_recent_projects_xlsx")
        extraction = ExtractionResult(
            extractor_name="excel:openpyxl",
            extraction_status="已提取文本",
            extracted_text=(
                "# 工作表：Sheet1\n\n"
                "| 贵阳贵安防汛抗旱指挥调度系统项目 | 列2 | 贵阳市水务管理局 | 438.52 | 319.28 | 列6 |\n"
                "| --- | --- | --- | --- | --- | --- |\n"
                "| 贵阳市低空飞行服务保障体系建设项目 | - | 贵阳市低空飞行服务保障体系建设项目：一是建立低空综合管理服务平台。二是在龙洞堡国际机场、重点景区营地、重要医疗机构、交通核心枢纽点建设为应用场景服务的各类起降场站，以及购买低空飞行器。 | 续建 | 贵阳市大数据局 | 贵阳市低空产业发展有限公司 |\n"
            ),
            preview_text="近期新建项目-长风可参与",
            text_length=220,
            page_count=1,
            source_encoding="xlsx",
            note="已完成 Excel 读取，共识别 1 个有效工作表。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)

        self.assertEqual(payload["方案名称字段"], "近期新建项目-长风可参与")
        self.assertEqual(payload["所属场景字段"], "低空项目机会")
        self.assertEqual(payload["客户名称字段"], "长风")
        self.assertEqual(payload["证据类型字段"], "清单型资料")
        self.assertNotEqual(payload["所属场景字段"], "医疗运输")

    def test_build_structured_payload_marks_skip_for_photo_image(self) -> None:
        contract = load_contract()
        base_sample = get_sample_by_id("policy_chongqing_txt_review")
        sample = base_sample.__class__(
            sample_id="photo_skip_demo",
            source_path="D:/长风资料 - 副本/现场照片.png",
            document_category="图片素材",
            recommended_template="产品设备模板",
            title_hint="现场照片",
            subject_name_hint="现场照片",
            product_name_hint="",
            unit_name_hint="",
            tags=["现场素材"],
            risks=[],
            notes=["用于验证纯照片跳过路径。"],
            evidence_level="L5",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="skip:image_photo",
            extraction_status="跳过",
            extracted_text="",
            preview_text="",
            text_length=0,
            page_count=1,
            source_encoding="",
            note="图片文件现场照片.png当前按纯照片处理，直接跳过，不进入 OCR。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)

        self.assertEqual(payload["处理路径"], "skip")
        self.assertEqual(payload["抽取状态"], "跳过")
        self.assertEqual(payload["分流结果"], "跳过")
        self.assertFalse(payload["是否适合直接入库"])

    def test_product_template_marks_wait_ocr_with_product_sections(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("nw_report_fy_q4")
        extraction = ExtractionResult(
            extractor_name="ocr:placeholder:pdf",
            extraction_status="待OCR",
            extracted_text="",
            preview_text="",
            text_length=0,
            page_count=7,
            source_encoding="",
            note="当前 PDF 未提取到可用文本，疑似扫描件或图片型 PDF，建议后续接 OCR。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)
            markdown = build_markdown(sample, payload)

        self.assertEqual(payload["处理路径"], "ocr")
        self.assertEqual(payload["抽取状态"], "待OCR")
        self.assertEqual(payload["产品名称字段"], "电子气象仪FY-Q4")
        self.assertEqual(payload["型号字段"], "FY-Q4")
        self.assertEqual(payload["产品类别字段"], "检测报告")
        self.assertEqual(payload["适用场景字段"], "气象监测")
        self.assertEqual(payload["搭配关系字段"], "未提取")
        self.assertEqual(payload["产品证据类型字段"], "检测报告")
        self.assertEqual(payload["报告文档类型字段"], "设备检测报告")
        self.assertEqual(payload["产品编号字段"], "")
        self.assertEqual(payload["检定依据字段"], "")
        self.assertEqual(payload["检定结果字段"], "")
        self.assertNotIn("## 二、核心用途", markdown)
        self.assertIn("## 五、适用场景", markdown)
        self.assertIn("## 七、字段提取结果", markdown)

    def test_product_template_skips_empty_placeholder_sections(self) -> None:
        contract = load_contract()
        sample = get_sample_by_id("nw_report_fy_q4")
        extraction = ExtractionResult(
            extractor_name="excel:openpyxl",
            extraction_status="已提取文本",
            extracted_text=(
                "# 工作表：Sheet1\n\n"
                "| 项目名称 | 项目内容 | 项目金额 |\n"
                "| --- | --- | --- |\n"
                "| 无人机医疗运输项目 | 医疗物资运输 | 100万 |\n"
            ),
            preview_text="# 工作表：Sheet1 | 项目名称 | 项目内容 | 项目金额 |",
            text_length=64,
            page_count=1,
            source_encoding="xlsx",
            note="已完成 Excel 读取，共识别 1 个有效工作表。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)

        payload.update(
            {
                "标题": "无人机医疗运输项目",
                "文件标题": "无人机医疗运输项目",
                "产品名称字段": "无人机医疗运输项目",
                "产品类别字段": "无人机设备",
                "产品证据类型字段": "说明书",
                "报告文档类型字段": "设备检验报告",
                "核心用途字段": "",
                "核心参数字段": "",
                "适用场景字段": "未提取",
                "搭配关系字段": "未提取",
                "型号字段": "",
                "供应商名称字段": "",
            }
        )

        markdown = build_markdown(sample, payload)

        self.assertNotIn("## 二、核心用途", markdown)
        self.assertNotIn("## 四、核心参数", markdown)
        self.assertNotIn("## 五、适用场景", markdown)
        self.assertNotIn("## 六、搭配、替代与挂载关系", markdown)
        self.assertNotIn("- 型号：", markdown)
        self.assertNotIn("- 供应商名称：", markdown)
        self.assertNotIn("- 核心用途：", markdown)
        self.assertNotIn("- 核心参数：", markdown)
        self.assertNotIn("- 适用场景：未提取", markdown)
        self.assertNotIn("- 搭配关系：未提取", markdown)
        self.assertIn("- 产品名称：无人机医疗运输项目", markdown)
        self.assertIn("- 产品类别：无人机设备", markdown)
        self.assertIn("- 证据类型：说明书", markdown)
        self.assertIn("- 报告文档类型：设备检验报告", markdown)

        contract = load_contract()
        sample = get_sample_by_id("nw_report_fy_q4")
        extraction = ExtractionResult(
            extractor_name="pypdf",
            extraction_status="已提取文本",
            extracted_text=(
                "产品名称：电子气象仪FY-Q4\n"
                "型号：FY-Q4\n"
                "生产单位：南京长风科技有限公司\n"
                "出厂编号：CF2026-0008\n"
                "检测依据：GB/T 12345-2026\n"
                "检测结果：合格\n"
                "检测机构：南京市计量测试研究院\n"
                "报告编号：BG-2026-021\n"
                "报告日期：2026年03月05日\n"
                "有效期至：2027年03月04日\n"
            ),
            preview_text="电子气象仪FY-Q4 检测报告",
            text_length=120,
            page_count=7,
            source_encoding="",
            note="已尝试使用 pypdf 提取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)
            markdown = build_markdown(sample, payload)

        self.assertEqual(payload["报告文档类型字段"], "设备检测报告")
        self.assertEqual(payload["产品编号字段"], "CF2026-0008")
        self.assertEqual(payload["检定依据字段"], "GB/T 12345-2026")
        self.assertEqual(payload["检定结果字段"], "合格")
        self.assertEqual(payload["检测机构字段"], "南京市计量测试研究院")
        self.assertEqual(payload["报告编号字段"], "BG-2026-021")
        self.assertEqual(payload["报告日期字段"], "2026-03-05")
        self.assertEqual(payload["有效期至字段"], "2027-03-04")
        self.assertIn("- 报告文档类型：设备检测报告", markdown)
        self.assertIn("- 检测机构：南京市计量测试研究院", markdown)

    def test_product_fields_fall_back_when_table_noise_enters_text(self) -> None:
        sample = SampleRecord(
            sample_id="scan_single_flight_moto",
            source_path="D:/长风资料/长风/长风-单人飞行摩托.docx",
            document_category="产品/设备",
            recommended_template="产品设备模板",
            title_hint="长风-单人飞行摩托",
            subject_name_hint="长风-单人飞行摩托",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L2",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="word:docx",
            extraction_status="已提取文本",
            extracted_text=(
                "# 工作表：Sheet1\n"
                "| 序号 | 产品名称 | 型号 | 备注 |\n"
                "| --- | --- | --- | --- |\n"
                "| 1 | 长风单人飞行摩托 | CF-01 | 演示样机 |\n"
            ),
            preview_text="长风单人飞行摩托",
            text_length=90,
            page_count=1,
            source_encoding="utf-8",
            note="已完成 docx 正文与表格提取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["产品名称字段"], "长风-单人飞行摩托")
        self.assertEqual(payload["型号字段"], "CF-01")

    def test_product_fields_avoid_table_header_as_product_name_or_model(self) -> None:
        sample = SampleRecord(
            sample_id="scan_cleaning_uav",
            source_path="D:/长风资料/长风/清洗无人机一览表.xlsx",
            document_category="产品/设备",
            recommended_template="产品设备模板",
            title_hint="清洗无人机一览表",
            subject_name_hint="清洗无人机一览表",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L2",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="excel:openpyxl",
            extraction_status="已提取文本",
            extracted_text=(
                "# 工作表：清单\n"
                "| 序号 | 名称 | 型号 | 规格参数 |\n"
                "| --- | --- | --- | --- |\n"
                "| 1 | 清洗无人机A | - | 载荷20kg |\n"
            ),
            preview_text="清洗无人机一览表",
            text_length=82,
            page_count=1,
            source_encoding="utf-8",
            note="已完成 Excel 读取，共识别 1 个有效工作表。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["产品名称字段"], "清洗无人机一览表")
        self.assertEqual(payload["型号字段"], "")

    def test_infer_document_title_skips_markdown_table_header(self) -> None:
        title = sample_docx_extract_to_md.infer_document_title(
            [
                "| 类别 | 参数项 | 参数值 |\n| --- | --- | --- |\n| 动力 | 载荷 | 20kg |",
                "清洗无人机一览表",
            ],
            "清洗无人机一览表.xlsx",
        )

        self.assertEqual(title, "清洗无人机一览表")

    def test_infer_document_title_skips_speaker_noise_for_education_training_pdf(self) -> None:
        title = sample_docx_extract_to_md.infer_document_title(
            [
                "主讲人：汪湖滨",
                "低空经济背景下的无人机专业建设及人才培养模式",
                "线上说明会",
            ],
            "低空经济背景下的无人机专业建设及人才培养模式.pdf",
        )

        self.assertEqual(title, "低空经济背景下的无人机专业建设及人才培养模式")

    def test_infer_document_title_prefers_explicit_policy_cover_title_over_section_heading(self) -> None:
        title = sample_docx_extract_to_md.infer_document_title(
            [
                "文件名称：《重庆市推动低空经济高质量发展若干政策措施》",
                "发文字号：渝府办发〔2025〕58号",
                "成文日期：2025年11月22日",
                "发文单位：重庆市人民政府办公厅",
                "三、加强低空保障体系建设",
            ],
            "重庆市推动低空经济高质量发展若干政策措施.txt",
        )

        self.assertEqual(title, "重庆市推动低空经济高质量发展若干政策措施")

    def test_infer_document_title_prefers_event_title_over_inner_forum_title(self) -> None:
        title = sample_docx_extract_to_md.infer_document_title(
            [
                "2026厦门国际低空经济暨无人机产业博览会（2026年4月23-25日）",
                "《低空安全监管技术应用交流会》",
                "时间:4月23-25日 地点:厦门国际会展中心",
            ],
            "2026厦门国际低空经济暨无人机产业博览会（2026年4月23-25日）.md",
        )

        self.assertEqual(title, "2026厦门国际低空经济暨无人机产业博览会")

    def test_infer_document_title_prefers_source_event_title_over_labeled_lines(self) -> None:
        title = sample_docx_extract_to_md.infer_document_title(
            [
                "展览场馆：上海汽车会展中心",
                "展会名称：第二届全球无人机与无人系统产业博览会（简称：GDE2026第二届国际无人机展）",
                "主办单位：国家技术转移东部中心",
            ],
            "GDE2026第二届国际无人机展（2026年4月23日-25日）.md",
        )

        self.assertEqual(title, "GDE2026第二届国际无人机展")

    def test_infer_document_title_ignores_numbered_event_agenda_heading(self) -> None:
        title = sample_docx_extract_to_md.infer_document_title(
            [
                "2026第二届低空经济暨无人系统应急应用技术创新展览会（2026年6月24日-26日）",
                "一、时间与地点",
                "四、论坛与配套活动安排",
                "1.低空经济生态产业链场景应用论坛",
            ],
            "2026第二届低空经济暨无人系统应急应用技术创新展览会（2026年6月24日-26日）.md",
        )

        self.assertEqual(title, "2026第二届低空经济暨无人系统应急应用技术创新展览会")

    def test_solution_name_falls_back_to_research_title(self) -> None:
        sample = SampleRecord(
            sample_id="scan_research_report",
            source_path="D:/长风资料/长风/低空产业（广州-深圳）入企调研、考察-陶袁华.pdf",
            document_category="方案/案例",
            recommended_template="方案案例模板",
            title_hint="低空产业（广州-深圳）入企调研、考察",
            subject_name_hint="低空产业（广州-深圳）入企调研、考察",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L2",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="pypdf",
            extraction_status="已提取文本",
            extracted_text="低空产业（广州-深圳）入企调研考察报告\n工业和信息化部等部门发布《通用航空装备创新应用实施方案（2024-2030年）》……",
            preview_text="低空产业（广州-深圳）入企调研考察报告",
            text_length=120,
            page_count=8,
            source_encoding="utf-8",
            note="已尝试使用 pypdf 提取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["方案名称字段"], "低空产业（广州-深圳）入企调研、考察")

    def test_education_training_fields_are_extracted_into_payload(self) -> None:
        sample = SampleRecord(
            sample_id="scan_education_training",
            source_path="D:/长风资料/长风/教育培训方向/职业教育专业目录.docx",
            document_category="教育培训",
            recommended_template="教育培训模板",
            title_hint="职业教育专业目录",
            subject_name_hint="教育培训方向",
            product_name_hint="",
            unit_name_hint="长风职业培训中心",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L2",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="word:docx",
            extraction_status="已提取文本",
            extracted_text=(
                "文件名称：职业教育专业目录\n"
                "培训单位：长风职业培训中心\n"
                "培训对象：职业院校学生、社会培训学员\n"
                "课程体系：无人机基础、飞行控制、行业应用实训\n"
                "教学内容：面向低空产业开展职业教育培训与课程实施。\n"
            ),
            preview_text="职业教育专业目录 长风职业培训中心",
            text_length=98,
            page_count=1,
            source_encoding="utf-8",
            note="已完成 docx 正文提取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["文件标题"], "职业教育专业目录")
        self.assertEqual(payload["单位名称字段"], "长风职业培训中心")
        self.assertEqual(payload["培训类型字段"], "职业教育")
        self.assertEqual(payload["适用对象字段"], "职业院校学生、社会培训学员")
        self.assertEqual(payload["课程体系字段"], "无人机基础、飞行控制、行业应用实训")

    def test_procurement_fields_are_extracted_into_payload(self) -> None:
        sample = SampleRecord(
            sample_id="scan_procurement_file",
            source_path="D:/长风资料/招采资料/1号标段采购文件.pdf",
            document_category="招标/采购文件",
            recommended_template="招标采购文件模板",
            title_hint="1号标段采购文件",
            subject_name_hint="招采资料",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L2",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="pypdf",
            extraction_status="已提取文本",
            extracted_text=(
                "项目名称：遵义市林业综合服务中心林业有害生物防治采购项目\n"
                "项目编号：GZYX 采-2026-35-1\n"
                "采购人：遵义市林业综合服务中心\n"
                "采购代理机构：贵州屹鑫工程管理有限公司\n"
                "采购方式：公开招标\n"
                "文件日期：2026年3月18日\n"
                "最高限价：120万元\n"
                "评标方法和评标标准\n"
                "本项目采用综合评分法。\n"
                "第四章 采购需求\n"
                "采购需求：完成林业有害生物监测、防治服务及成果提交。\n"
            ),
            preview_text="遵义市林业综合服务中心林业有害生物防治采购项目",
            text_length=180,
            page_count=12,
            source_encoding="utf-8",
            note="已尝试使用 pypdf 提取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["文件标题"], "遵义市林业综合服务中心林业有害生物防治采购项目")
        self.assertEqual(payload["项目编号字段"], "GZYX 采-2026-35-1")
        self.assertEqual(payload["采购人字段"], "遵义市林业综合服务中心")
        self.assertEqual(payload["采购代理机构字段"], "贵州屹鑫工程管理有限公司")
        self.assertEqual(payload["采购方式字段"], "公开招标")
        self.assertEqual(payload["文件日期字段"], "2026-03-18")
        self.assertEqual(payload["预算最高限价字段"], "120万元")
        self.assertEqual(payload["评分办法字段"], "综合评分法")
        self.assertEqual(payload["采购需求摘要字段"], "完成林业有害生物监测、防治服务及成果提交。")

    def test_procurement_invalid_file_date_returns_empty(self) -> None:
        sample = SampleRecord(
            sample_id="scan_procurement_invalid_date",
            source_path="D:/长风资料/招采资料/1号标段采购文件.pdf",
            document_category="招标/采购文件",
            recommended_template="招标采购文件模板",
            title_hint="1号标段采购文件",
            subject_name_hint="招采资料",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L2",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="pypdf",
            extraction_status="已提取文本",
            extracted_text=(
                "项目名称：遵义市林业综合服务中心林业有害生物防治采购项目\n"
                "采购方式：公开招标\n"
                "文件日期：2026年35月1日\n"
            ),
            preview_text="遵义市林业综合服务中心林业有害生物防治采购项目",
            text_length=64,
            page_count=2,
            source_encoding="utf-8",
            note="已尝试使用 pypdf 提取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["文件日期字段"], "")

    def test_bid_summary_metadata_cleans_heading_and_assigns_template(self) -> None:
        source_path = Path("D:/长风资料/低空资源/近期低空经济相关招投标信息汇总.txt")
        text = (
            "# 近期低空经济相关招投标信息汇总（更新至2026年03月02日）\n"
            "发布日期    项目名称    采购人    预算金额    截止时间\n"
            "2026-03-01    低空巡检服务项目    某市应急管理局    120万元    2026-03-20\n"
            "2026-03-02    无人机平台建设项目    某市交通运输局    260万元    2026-03-25\n"
            "2026-03-02    低空监管系统项目    某市城管局    310万元    2026-03-28\n"
        )

        metadata = build_bid_summary_metadata(source_path, text)

        self.assertIsNotNone(metadata)
        self.assertEqual(metadata["标题"], "近期低空经济相关招投标信息汇总（更新至2026年03月02日）")
        self.assertEqual(metadata["文件标题"], "近期低空经济相关招投标信息汇总（更新至2026年03月02日）")
        self.assertEqual(metadata["推荐模板"], "招投标汇总模板")
        self.assertEqual(metadata["模板归属"], "招投标汇总模板")
        self.assertEqual(metadata["招投标记录数"], 3)
        self.assertEqual(metadata["招投标发布日期范围"], "2026-03-01 至 2026-03-02")

    def test_normalize_policy_date_rejects_invalid_month(self) -> None:
        self.assertEqual(normalize_policy_date("2026年35月1日"), "")
        self.assertEqual(normalize_policy_date("2026年3月"), "2026-03")

    def test_procurement_template_markdown_contains_dedicated_sections(self) -> None:
        sample = SampleRecord(
            sample_id="scan_procurement_markdown",
            source_path="D:/长风资料/招采资料/1号标段采购文件.pdf",
            document_category="招标/采购文件",
            recommended_template="招标采购文件模板",
            title_hint="1号标段采购文件",
            subject_name_hint="招采资料",
            product_name_hint="",
            unit_name_hint="",
            tags=["招标采购"],
            risks=[],
            notes=[],
            evidence_level="L2",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        payload = {
            "标题": "1号标段采购文件",
            "原始文件名": "1号标段采购文件.pdf",
            "原始路径": "D:/长风资料/招采资料/1号标段采购文件.pdf",
            "文件格式": "pdf",
            "文档分类": "招标/采购文件",
            "推荐模板": "招标采购文件模板",
            "一级分类": "招采资料",
            "二级分类": "采购文件",
            "分类置信度": "高",
            "分类依据": "采购文件、采购人、采购方式、评标方法和评标标准等关键词命中。",
            "处理路径": "document_parse",
            "抽取状态": "已提取文本",
            "抽取器": "pypdf",
            "文本长度": 180,
            "页数": 12,
            "文本编码": "utf-8",
            "证据等级": "L2",
            "分流结果": "待审核",
            "内容主题标签": ["招标采购"],
            "文本预览": "遵义市林业综合服务中心林业有害生物防治采购项目",
            "文件标题": "遵义市林业综合服务中心林业有害生物防治采购项目",
            "项目编号字段": "GZYX 采-2026-35-1",
            "采购人字段": "遵义市林业综合服务中心",
            "采购代理机构字段": "贵州屹鑫工程管理有限公司",
            "采购方式字段": "公开招标",
            "文件日期字段": "2026-03-18",
            "预算最高限价字段": "120万元",
            "评分办法字段": "综合评分法",
            "采购需求摘要字段": "完成林业有害生物监测、防治服务及成果提交。",
            "风险说明": [],
            "备注": [],
            "去重主键": ["1号标段采购文件"],
            "抽取说明": "已按采购文件模板生成最小结构化结果。",
            "是否适合直接入库": False,
            "是否需要拆分": False,
            "拆分说明": "",
        }

        markdown = build_markdown(sample, payload)

        self.assertIn("## 二、采购方式与预算", markdown)
        self.assertIn("## 三、评分办法", markdown)
        self.assertIn("## 四、采购需求摘要", markdown)
        self.assertIn("## 五、字段提取结果", markdown)
        self.assertIn("- 项目编号：GZYX 采-2026-35-1", markdown)
        self.assertNotIn("## 抽取说明", markdown)

    def test_build_structured_payload_keeps_scan_title_hint_when_auto_title_is_section_heading(self) -> None:
        contract = load_contract()
        sample = SampleRecord(
            sample_id="scan_education_training_auto_title_fallback",
            source_path="D:/长风资料/天翼航空相关资料/高质量产教融合无人机专业五金建设方案.pdf",
            document_category="待判定资料",
            recommended_template="待人工补规则",
            title_hint="高质量产教融合无人机专业五金建设方案",
            subject_name_hint="",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L3",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="pypdf",
            extraction_status="已提取文本",
            extracted_text=(
                "模块一 Module 1 国家政策\n"
                "模块二 Module 2 企业介绍\n"
                "模块三 Module 3 建设方向\n"
                "模块四 Module 4 专业教学数字化转型\n"
                "模块五 Module 5 课程建设体系图\n"
                "模块六 Module 6 教学资源开发流程\n"
            ),
            preview_text="模块一 Module 1 国家政策",
            text_length=96,
            page_count=41,
            source_encoding="utf-8",
            note="已尝试使用 pypdf 提取。",
        )

        with patch("minimum_workflow.pipeline.extract_with_strategy", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)

        self.assertEqual(payload["标题"], "高质量产教融合无人机专业五金建设方案")
        self.assertEqual(payload["文件标题"], "高质量产教融合无人机专业五金建设方案")
        self.assertEqual(payload["培训主题字段"], "")
        self.assertEqual(payload["专业方向字段"], "")

    def test_education_training_topic_does_not_fall_back_to_direction(self) -> None:
        sample = SampleRecord(
            sample_id="scan_education_training_fallback",
            source_path="D:/长风资料/长风/教育培训方向/新建文本文档.txt",
            document_category="教育培训",
            recommended_template="教育培训模板",
            title_hint="",
            subject_name_hint="教育培训方向",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L2",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="text:txt",
            extraction_status="已提取文本",
            extracted_text=(
                "2025年职业教育专业目录\n"
                "序号 专业大类 专业类 专业名称\n"
                "1 装备制造大类 航空装备类 无人机应用技术\n"
            ),
            preview_text="2025年职业教育专业目录",
            text_length=52,
            page_count=1,
            source_encoding="utf-8",
            note="已完成 txt 正文提取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["培训主题字段"], "")
        self.assertEqual(payload["专业方向字段"], "")

    def test_education_training_long_paragraph_does_not_become_title_or_topic(self) -> None:
        sample = SampleRecord(
            sample_id="scan_education_training_long_paragraph",
            source_path="D:/长风资料/长风/教育培训方向/新建文本文档.txt",
            document_category="教育培训",
            recommended_template="教育培训模板",
            title_hint="新建文本文档",
            subject_name_hint="教育培训方向",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L2",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        long_line = (
            "按照 2025 年增补专业清单要求持续推进低空经济领域职业教育体系建设，"
            "围绕无人机应用技术、航空装备维护、飞行服务保障等方向统筹开展课程设计、"
            "实践教学、校企合作与人才培养工作，形成面向区域产业需求的专业建设方案。"
        )
        extraction = ExtractionResult(
            extractor_name="text:txt",
            extraction_status="已提取文本",
            extracted_text=long_line,
            preview_text=long_line[:40],
            text_length=len(long_line),
            page_count=1,
            source_encoding="utf-8",
            note="已完成 txt 正文提取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["文件标题"], "")
        self.assertEqual(payload["培训主题字段"], "")
        self.assertNotEqual(payload["文件标题"], long_line)
        self.assertNotEqual(payload["培训主题字段"], long_line)

    def test_product_model_ignores_battery_capacity_token(self) -> None:
        sample = SampleRecord(
            sample_id="scan_single_flight_moto_model",
            source_path="D:/长风资料/长风/长风-单人飞行摩托.docx",
            document_category="产品/设备",
            recommended_template="产品设备模板",
            title_hint="长风-单人飞行摩托",
            subject_name_hint="长风-单人飞行摩托",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L2",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="word:docx",
            extraction_status="已提取文本",
            extracted_text=(
                "长风-单人飞行摩托\n"
                "| 项目 | 参数 |\n"
                "| --- | --- |\n"
                "| 电池 | 24S36000毫安锂电池 |\n"
                "| 定位模式 | GPS+北斗 |\n"
                "| 续航时间 | 15min |\n"
            ),
            preview_text="长风-单人飞行摩托",
            text_length=70,
            page_count=1,
            source_encoding="utf-8",
            note="已完成 docx 正文与表格提取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["型号字段"], "")

    def test_normalize_solution_summary_payload_accepts_string_and_list_values(self) -> None:
        payload = {
            "主体名称": ["杭州迅蚁", "送吧航空"],
            "方案名称/案例名称": "六盘水医院无人机运输方案",
            "所属场景": "医疗物流",
            "客户/使用单位": "钟山区人民医院",
            "文件日期": "2025-12-08",
            "资料摘要": ["方案摘要A", "方案摘要B"],
            "资料形态判断": ["商业方案", "报价材料"],
            "应用背景": ["背景一", "背景二"],
            "解决的问题": "时效慢；冷链难",
            "投入的产品/设备/能力": ["TR9S", "温控箱"],
            "实施方式": "航线规划\n飞行监控",
            "预算、进度与组织方式": ["采购服务"],
            "结果与效果数据": ["时效提升"],
            "可复用经验": "医共体配送；标准化SOP",
            "入库与归档判断": "适合摘要入库",
            "备注": "",
        }

        normalized = normalize_solution_summary_payload(payload)

        self.assertEqual(normalized["主体名称"], "杭州迅蚁；送吧航空")
        self.assertEqual(normalized["资料摘要"], "方案摘要A；方案摘要B")
        self.assertEqual(normalized["资料形态判断"], "商业方案；报价材料")
        self.assertEqual(normalized["应用背景"], ["背景一", "背景二"])
        self.assertEqual(normalized["解决的问题"], ["时效慢", "冷链难"])
        self.assertEqual(normalized["实施方式"], ["航线规划", "飞行监控"])
        self.assertEqual(normalized["可复用经验"], ["医共体配送", "标准化SOP"])
        self.assertEqual(normalized["备注"], [])

    def test_classify_image_directory_detects_paged_document(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        source_dir = Path(temp_dir.name) / "低空智能网联体系参考架构（2024版）工信部装备工业发展中心2024-11-11"
        source_dir.mkdir()
        for name in ("00封面.png", "01.png", "02.jpeg"):
            (source_dir / name).write_bytes(b"fake-image")

        image_paths, reason = sample_docx_extract_to_md.classify_image_directory(source_dir)

        self.assertEqual([path.name for path in image_paths], ["00封面.png", "01.png", "02.jpeg"])
        self.assertIn("分页命名图片 3/3", reason)
        self.assertIn("命中封面页", reason)

    def test_classify_image_directory_rejects_photo_album(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        source_dir = Path(temp_dir.name) / "项目现场照片"
        source_dir.mkdir()
        for name in ("DJI_001.jpg", "微信图片_2024.jpg", "现场照片3.png"):
            (source_dir / name).write_bytes(b"fake-image")

        with self.assertRaises(RuntimeError) as ctx:
            sample_docx_extract_to_md.classify_image_directory(source_dir)

        self.assertIn("纯照片目录", str(ctx.exception))

    def test_classify_image_directory_accepts_wechat_screenshot_sequence(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        source_dir = Path(temp_dir.name) / "教育培训方向"
        source_dir.mkdir()
        for name in ("微信图片_202604120001.png", "微信图片_202604120002.png", "微信图片_202604120003.png"):
            (source_dir / name).write_bytes(b"fake-image")

        image_paths, reason = sample_docx_extract_to_md.classify_image_directory(source_dir)

        self.assertEqual(
            [path.name for path in image_paths],
            ["微信图片_202604120001.png", "微信图片_202604120002.png", "微信图片_202604120003.png"],
        )
        self.assertIn("分页命名图片 3/3", reason)

    def test_extract_source_content_supports_directory_source(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        source_dir = Path(temp_dir.name) / "低空智能网联体系参考架构（2024版）工信部装备工业发展中心2024-11-11"
        source_dir.mkdir()
        directory_result = {
            "blocks": ["第一页", "第二页"],
            "extracted_text": "第一页\n\n第二页",
            "is_heavy_pdf": False,
            "heavy_pdf_reason": "",
            "extraction_note": "目录OCR成功",
            "auto_metadata": [("文档分类", "行业参考架构/指导材料")],
            "selection_reason": "目录判定成功",
        }

        with patch.object(sample_docx_extract_to_md, "extract_image_directory_content", return_value=directory_result) as directory_mock:
            result = sample_docx_extract_to_md.extract_source_content(source_dir)

        directory_mock.assert_called_once_with(source_dir)
        self.assertEqual(result["extraction_note"], "目录OCR成功")
        self.assertEqual(result["auto_metadata"][0][0], "文档分类")

    def test_extract_source_content_uses_mineru_for_docx_and_builds_auto_metadata(self) -> None:
        source_path = Path("D:/tmp/云深处科技公司介绍（2026）.docx")
        batch_result = {
            "batch_id": "batch-docx-1",
            "results": [
                {
                    "state": "done",
                    "markdown": "# 云深处科技公司介绍\n\n杭州云深处科技有限公司\n\n公司简介\n\n四足机器人产品矩阵",
                    "error": "",
                }
            ],
        }

        with patch.object(sample_docx_extract_to_md, "resolve_mineru_token", return_value="token-demo"), patch.object(
            sample_docx_extract_to_md,
            "run_mineru_batch",
            return_value=batch_result,
        ) as run_mock, patch.object(
            sample_docx_extract_to_md,
            "try_extract_with_markitdown",
            create=True,
        ) as markitdown_mock, patch.object(
            sample_docx_extract_to_md,
            "extract_docx_blocks",
            return_value=[],
        ):
            result = sample_docx_extract_to_md.extract_source_content(source_path)

        run_mock.assert_called_once_with([source_path], "token-demo")
        markitdown_mock.assert_not_called()
        metadata = dict(result["auto_metadata"])
        self.assertEqual(metadata["抽取器"], "mineru:batch")
        self.assertEqual(metadata["转换状态"], "MinerU成功")
        self.assertEqual(metadata["唯一类别判定"], "供应商企业模板")
        self.assertIn("MinerU格式转换", metadata["处理链路"])
        self.assertIn("云深处科技公司介绍", result["extracted_text"])

    def test_extract_source_content_falls_back_to_local_docx_when_mineru_token_missing(self) -> None:
        source_path = Path("D:/tmp/智慧园区巡检技术方案.docx")

        with patch.object(sample_docx_extract_to_md, "resolve_mineru_token", return_value=None), patch.object(
            sample_docx_extract_to_md,
            "extract_docx_blocks",
            return_value=["智慧园区巡检技术方案", "项目背景", "巡检机器人"],
        ) as docx_mock:
            result = sample_docx_extract_to_md.extract_source_content(source_path)

        docx_mock.assert_called_once_with(source_path)
        metadata = dict(result["auto_metadata"])
        self.assertEqual(metadata["抽取器"], "word:docx")
        self.assertEqual(metadata["转换状态"], "回退本地解析")
        self.assertEqual(metadata["唯一类别判定"], "方案案例模板")
        self.assertIn("未提供 MinerU token", result["extraction_note"])

    def test_extract_source_content_uses_mineru_image_chain_for_large_pptx(self) -> None:
        source_path = Path("D:/tmp/云深处科技公司介绍（2026）.pptx")
        image_result = ExtractionResult(
            extractor_name="mineru:presentation_images",
            extraction_status="已提取文本",
            extracted_text="具身智能技术创新与行业应用引领者\n\n公司介绍\n\n四足机器人产品矩阵",
            preview_text="具身智能技术创新与行业应用引领者",
            text_length=40,
            page_count=42,
            source_encoding="utf-8",
            note="原始演示文稿超过 MinerU 单文件体积限制，已导出 42 页图片后按页序走 MinerU OCR；批次号：batch-pptx-image-1；OCR失败页：无。",
        )

        with patch.object(sample_docx_extract_to_md, "resolve_mineru_token", return_value="token-demo"), patch.object(
            Path,
            "stat",
            return_value=type("Stat", (), {"st_size": sample_docx_extract_to_md.MINERU_SINGLE_FILE_LIMIT_BYTES + 1})(),
        ), patch.object(
            sample_docx_extract_to_md,
            "run_mineru_batch",
        ) as run_mock, patch.object(
            sample_docx_extract_to_md,
            "extract_large_presentation_via_mineru_images",
            return_value=image_result,
        ) as image_mock:
            result = sample_docx_extract_to_md.extract_source_content(source_path)

        run_mock.assert_not_called()
        image_mock.assert_called_once_with(source_path, "token-demo")
        metadata = dict(result["auto_metadata"])
        self.assertEqual(metadata["抽取器"], "mineru:presentation_images")
        self.assertEqual(metadata["转换状态"], "MinerU成功")
        self.assertIn("MinerU格式转换", metadata["处理链路"])
        self.assertIn("公司介绍", result["extracted_text"])
        self.assertIn("页图片后按页序走 MinerU OCR", result["extraction_note"])

    def test_extract_source_content_falls_back_to_local_pptx_when_image_chain_fails(self) -> None:
        source_path = Path("D:/tmp/云深处科技-产教融合2026.pptx")
        batch_result = {
            "batch_id": "batch-pptx-oversize",
            "results": [
                {
                    "state": "failed",
                    "markdown": "",
                    "error": "file size exceeds limit (200MB), please choose a smaller file",
                }
            ],
        }
        local_result = ExtractionResult(
            extractor_name="presentation:pptx",
            extraction_status="已提取文本",
            extracted_text="产教融合\n\n打造集教学、实训、科研、竞赛于一体的综合性平台",
            preview_text="产教融合",
            text_length=32,
            page_count=20,
            source_encoding="pptx",
            note="已完成演示文稿读取。",
        )

        with patch.object(sample_docx_extract_to_md, "resolve_mineru_token", return_value="token-demo"), patch.object(
            sample_docx_extract_to_md,
            "run_mineru_batch",
            return_value=batch_result,
        ) as run_mock, patch.object(
            sample_docx_extract_to_md,
            "extract_large_presentation_via_mineru_images",
            side_effect=RuntimeError("PowerPoint COM not available"),
        ) as image_mock, patch.object(
            sample_docx_extract_to_md,
            "extract_text",
            return_value=local_result,
        ) as local_mock:
            result = sample_docx_extract_to_md.extract_source_content(source_path)

        run_mock.assert_called_once_with([source_path], "token-demo")
        image_mock.assert_called_once_with(source_path, "token-demo")
        local_mock.assert_called_once_with(source_path, "presentation")
        metadata = dict(result["auto_metadata"])
        self.assertEqual(metadata["抽取器"], "presentation:pptx")
        self.assertEqual(metadata["转换状态"], "回退本地解析")
        self.assertIn("分页图片 OCR 失败", result["extraction_note"])
        self.assertIn("产教融合", result["extracted_text"])

    def test_extract_source_content_falls_back_to_local_pdf_when_mineru_text_is_suspicious(self) -> None:
        source_path = Path("D:/tmp/sample.pdf")
        mineru_result = ExtractionResult(
            extractor_name="mineru:batch",
            extraction_status="已提取文本",
            extracted_text="中国一中红广通国\n无人机驾驶飞行器故障检测报告书\n正常片段",
            preview_text="中国一中红广通国",
            text_length=30,
            page_count=None,
            source_encoding="utf-8",
            note="已通过 MinerU 批量接口完成 Markdown 提取。",
        )
        fallback_result = ExtractionResult(
            extractor_name="pdf:pymupdf",
            extraction_status="已提取文本",
            extracted_text="智慧园区巡检技术方案\n\n项目背景\n\n巡检机器人",
            preview_text="智慧园区巡检技术方案",
            text_length=24,
            page_count=3,
            source_encoding="utf-8",
            note="已完成本地 PDF 解析。",
        )

        with patch.object(
            sample_docx_extract_to_md,
            "extract_office_document_content",
            return_value={
                "blocks": ["中国一中红广通国", "无人机驾驶飞行器故障检测报告书"],
                "extracted_text": mineru_result.extracted_text,
                "is_heavy_pdf": False,
                "heavy_pdf_reason": "",
                "extraction_note": mineru_result.note,
                "auto_metadata": [("抽取器", "mineru:batch")],
                "selection_reason": "初始 MinerU 结果。",
                "extraction_result": mineru_result,
            },
        ), patch.object(sample_docx_extract_to_md, "extract_pdf_text", return_value=fallback_result) as pdf_mock:
            result = sample_docx_extract_to_md.extract_source_content(source_path)

        pdf_mock.assert_called_once_with(source_path)
        metadata = dict(result["auto_metadata"])
        self.assertEqual(metadata["抽取器"], "pdf:pymupdf")
        self.assertEqual(metadata["转换状态"], "回退本地解析")
        self.assertIn("已回退本地解析", result["extraction_note"])
        self.assertIn("智慧园区巡检技术方案", result["extracted_text"])

    def test_infer_document_profile_marks_reference_architecture(self) -> None:
        blocks = [
            "# 低空智能网联体系参考架构（2024版）",
            "工业和信息化部装备工业发展中心",
            "2024年11月",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "低空智能网联体系参考架构（2024版）工信部装备工业发展中心2024-11-11",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "行业参考架构/指导材料")
        self.assertIn("参考架构/白皮书口径", profile["模板归属"])
        self.assertEqual(profile["文件标题"], "低空智能网联体系参考架构（2024版）")
        self.assertIn("装备工业发展中心", profile["主体名称"])
        self.assertEqual(profile["发布时间"], "2024年11月")

    def test_infer_document_profile_marks_research_report_as_solution(self) -> None:
        blocks = [
            "低空产业（广州-深圳）入企调研、考察",
            "调研背景：围绕低空产业链开展企业走访。",
            "应用场景：低空文旅、物流配送。",
            "合作方向：重点关注运营模式和项目落地。",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "低空产业（广州-深圳）入企调研、考察-陶袁华",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "方案/案例")
        self.assertEqual(profile["模板归属"], "方案案例模板")
        self.assertEqual(profile["文件标题"], "低空产业（广州-深圳）入企调研、考察")

    def test_infer_document_profile_marks_education_training_directory(self) -> None:
        blocks = [
            "职业教育专业目录",
            "培训对象：中职院校学生与社会培训学员。",
            "课程体系：无人机基础、飞行控制、行业应用实训。",
            "教学内容：围绕低空应用开展职业教育培训。",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "教育培训方向",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "教育培训")
        self.assertEqual(profile["模板归属"], "教育培训模板")
        self.assertEqual(profile["资料层级"], "教育培训资料")

    def test_infer_document_profile_marks_education_training_from_strong_content_only(self) -> None:
        blocks = [
            "2025年专业设置说明",
            "职业教育专业目录覆盖低空装备相关专业。",
            "培训对象：中职院校学生与社会培训学员。",
            "课程体系：无人机基础、飞行控制、行业应用实训。",
            "教学内容：围绕低空应用开展岗位训练。",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "新建文本文档",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "教育培训")
        self.assertEqual(profile["模板归属"], "教育培训模板")
        self.assertEqual(profile["资料层级"], "教育培训资料")

    def test_infer_document_profile_marks_education_training_for_wujin_construction_sample(self) -> None:
        blocks = [
            "高质量产教融合无人机专业五金建设方案",
            "专业教学数字化转型",
            "课程建设体系图",
            "教学资源开发流程",
            "部分建设案例",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "高质量产教融合无人机专业五金建设方案.pdf",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "教育培训")
        self.assertEqual(profile["模板归属"], "教育培训模板")

    def test_infer_document_profile_keeps_wechat_meeting_scan_out_of_policy_direct_ingest(self) -> None:
        blocks = [
            "迅蚁科技低空教育项目线上说明会",
            "主讲人：汪湖滨（迅蚁科技业务副总裁/教育事业部负责人）",
            "请勿外传，仅供内部交流",
            "教育部办公厅关于做好2026年职业教育拟招生专业设置管理工作的通知",
            "抓住政策红利窗口期·共享低空经济发展新机遇",
            "2026年4月8日",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "迅蚁2026-04-08线上会",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "方案/案例")
        self.assertEqual(profile["模板归属"], "方案案例模板")
        self.assertEqual(profile["是否适合直接入Dify"], "待审核")
        self.assertEqual(profile["发布时间"], "")

    def test_infer_document_profile_keeps_event_markdown_out_of_policy_template(self) -> None:
        blocks = [
            "2026厦门国际低空经济暨无人机产业博览会",
            "时间:4月23-25日 地点:厦门国际会展中心",
            "展会背景：为贯彻落实党中央关于发展通用航空和低空经济的战略部署，抢抓低空经济发展窗口期。",
            "近年来，我国政府高度重视低空经济的发展，出台了一系列政策措施，推动低空经济的规范化、规模化发展。",
            "同期活动：低空安全监管技术应用交流会。",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "2026厦门国际低空经济暨无人机产业博览会（2026年4月23-25日）.md",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "方案/案例")
        self.assertEqual(profile["模板归属"], "方案案例模板")
        self.assertEqual(profile["是否适合直接入Dify"], "待审核")
        self.assertEqual(profile["文件标题"], "2026厦门国际低空经济暨无人机产业博览会")
        self.assertEqual(profile["主体名称"], "")

    def test_infer_document_profile_prefers_reference_title_for_framework_guide(self) -> None:
        blocks = [
            "# 低空经济基础设施框架指引（2025版）",
            "# 中国民用机场协会",
            "为助力我国低空经济高质量发展，服务低空基础设施科学规划、有序建设，现研究编制《低空经济基础设施框架指引（2025版）》。",
            "《无人驾驶航空器飞行管理暂行条例》",
            "《民用无人驾驶航空器运行安全管理规则》",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "P01-P07《低空经济基础设施框架指引（2025年版）》.pdf",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "行业参考架构/指导材料")
        self.assertEqual(profile["模板归属"], "参考架构/白皮书口径（当前按原文全量提取Markdown输出）")
        self.assertEqual(profile["文件标题"], "低空经济基础设施框架指引（2025年版）")

    def test_infer_document_profile_marks_platform_solution_digest_as_solution(self) -> None:
        blocks = [
            "低空经济数字化监管服务平台解决方案汇总解读（1）——国家和省级平台",
            "低空数字化监管服务平台是必要条件之一，使得在“三维空间”的低空飞行看得见、呼得到、管得住、飞得稳。",
            "现就常见的低空经济数字化监管服务平台进行汇总，并对主要功能特点进行解读。",
            "国家级平台",
            "省级平台",
            "中国民航局主管的全国民用无人驾驶航空器低空飞行相关监管服务平台，主要功能包括：登记管理、适航管理、操控员资质、运营许可、市场管理、运行管理、系统管理。",
            "| 平台 | 建设单位 | 核心功能 |\n| --- | --- | --- |\n| 广东省低空飞行综合管理服务平台 | 广东省低空经济产业发展有限公司 | 数据整合与管理、设备连接与协同、智能分析与决策支持 |",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "低空经济数字化监管服务平台解决方案汇总解读（1）——国家和省级平台.docx",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "方案/案例")
        self.assertEqual(profile["模板归属"], "方案案例模板")
        self.assertEqual(profile["文件标题"], "低空经济数字化监管服务平台解决方案汇总解读（1）——国家和省级平台")

    def test_infer_document_profile_marks_base_plan_as_solution_not_product(self) -> None:
        blocks = [
            "XX城市低空经济示范基地方案",
            "一、基地定位",
            "总体定位：XX城市低空经济示范基地拟在市政府指导下建设。",
            "构建“一区一园一中枢”载体，覆盖低空空域管理、技术研发、行业应用、产业孵化、教育培训和物流运输。",
            "二、建设思路",
            "包括无人机全自动机场、无人机通信基站、eVTOL起降平台、无人机反制、导航定位系统及能源供应网络。",
            "| 产出 | 输出能力 | 目标 |\n| --- | --- | --- |\n| 低空智能计算中心 | 智能算力 | 10P |\n| 巡检无人机 | 套 | 100 |",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "低空经济示范基地方案.docx",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "方案/案例")
        self.assertEqual(profile["模板归属"], "方案案例模板")
        self.assertEqual(profile["文件标题"], "低空经济示范基地方案")

        blocks = [
            "文件名称：《重庆市推动低空经济高质量发展若干政策措施》",
            "发文字号：渝府办发〔2025〕58号",
            "成文日期：2025年11月22日",
            "发文单位：重庆市人民政府办公厅",
            "重庆市推动低空经济高质量发展若干政策措施",
            "三、加强低空保障体系建设",
            "支持低空飞行应用场景建设。",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "重庆市推动低空经济高质量发展若干政策措施.txt",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "政策/官方文件")
        self.assertEqual(profile["模板归属"], "政策官方文件模板")
        self.assertEqual(profile["文件标题"], "重庆市推动低空经济高质量发展若干政策措施")
        self.assertEqual(profile["主体名称"], "重庆市人民政府办公厅")
        self.assertEqual(profile["发布时间"], "2025年11月22日")
        self.assertEqual(profile["是否适合直接入Dify"], "是")

    def test_infer_document_profile_prefers_education_training_over_meeting_noise(self) -> None:
        blocks = [
            "主讲人：汪湖滨",
            "低空经济背景下的无人机专业建设及人才培养模式",
            "低空教育项目线上说明会",
            "课程体系：无人机基础、飞行控制、行业应用实训。",
            "产教融合推进专业建设与岗课赛证融通。",
            "专业建设及人才培养模式面向本科、高职、中职及技师院校。",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "低空经济背景下的无人机专业建设及人才培养模式.pdf",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "教育培训")
        self.assertEqual(profile["模板归属"], "教育培训模板")
        self.assertEqual(profile["文件标题"], "低空经济背景下的无人机专业建设及人才培养模式")

    def test_infer_document_profile_keeps_procurement_file_out_of_policy_template(self) -> None:
        blocks = [
            "采购文件",
            "项目编号：GZYX 采-2026-35-1",
            "采购人：遵义市林业综合服务中心",
            "采购代理机构：贵州屹鑫工程管理有限公司",
            "采购方式：公开招标",
            "第一章 采购公告",
            "第二章 供应商须知",
            "第三章 评标方法和评标标准",
            "第四章 采购需求",
            "第六章 合同主要条款（参考）",
            "第七章 响应文件格式",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "1号标段采购文件.pdf",
            blocks,
        )

        self.assertEqual(profile["模板归属"], "招标采购文件模板")
        self.assertEqual(profile["文档分类"], "招标/采购文件")
        self.assertEqual(profile["是否适合直接入Dify"], "待审核")
        self.assertEqual(profile["文件标题"], "1号标段采购文件")

    def test_infer_document_profile_marks_supplier_company_intro_ppt(self) -> None:
        blocks = [
            "具身智能技术创新与行业应用引领者",
            "公司介绍",
            "云深处科技，聚焦具身智能技术创新与行业应用，是专注于四足机器人、人形机器人及核心零部件的研发、生产、销售和服务的国家级高新技术企业。",
            "杭州云深处科技股份有限公司",
            "成立于2017年11月，位于杭州市西湖区。",
            "发展历程",
            "荣誉资质",
            "创始人简介",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "云深处科技公司介绍（2026）.pptx",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "供应商/企业资料")
        self.assertEqual(profile["模板归属"], "供应商企业模板")
        self.assertEqual(profile["是否适合直接入Dify"], "待审核")

    def test_infer_document_profile_marks_yunshenchu_solution_ppt(self) -> None:
        blocks = [
            "具身智能技术创新与行业应用引领者",
            "公司简介",
            "巡检方案",
            "应用案例",
            "智慧园区巡检解决方案",
            "技术方案 - 整体架构",
            "技术方案 - 巡检流程",
            "应用赋能 - 安防场景",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "【云深处】企业园区巡检解决方案 20251017_V1.0.pptx",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "方案/案例")
        self.assertEqual(profile["模板归属"], "方案案例模板")
        self.assertEqual(profile["是否适合直接入Dify"], "待审核")

    def test_infer_document_profile_marks_yunshenchu_industry_education_ppt(self) -> None:
        blocks = [
            "具身智能技术创新与行业应用引领者",
            "公司介绍",
            "打造集教学、实训、科研、竞赛于一体的综合性平台",
            "立足行业，赋能教育，促进产教融合发展！",
            "具身智能产教融合的方案探索",
            "联合实验室的搭建",
            "课程体系建设",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "云深处科技-产教融合2026.pptx",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "教育培训")
        self.assertEqual(profile["模板归属"], "教育培训模板")
        self.assertEqual(profile["是否适合直接入Dify"], "待审核")

    def test_infer_document_profile_marks_overall_intro_as_supplier_template(self) -> None:
        blocks = [
            "北京云圣智能总体介绍",
            "公司简介",
            "北京云圣智能科技有限责任公司专注于工业无人机与智能巡检。",
            "主营业务覆盖智慧巡检、应急消防与低空监管。",
            "发展历程",
            "荣誉资质",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "北京云圣智能总体介绍.pptx",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "供应商/企业资料")
        self.assertEqual(profile["模板归属"], "供应商企业模板")

    def test_infer_document_profile_keeps_master_plan_out_of_price_quote_template(self) -> None:
        blocks = [
            "百里杜鹃景区低空经济项目总体规划方案",
            "建设目标",
            "建设内容包括飞行服务、低空文旅、医疗运输与应急救援。",
            "投资计划",
            "总体架构",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "百里杜鹃景区低空经济项目总体规划方案2025-11-20.docx",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "方案/案例")
        self.assertEqual(profile["模板归属"], "方案案例模板")

    def test_infer_document_profile_marks_project_catalog_as_solution(self) -> None:
        blocks = [
            "无人机医疗运输项目清单",
            "| 项目名称 | 应用场景 | 项目简介 |",
            "| --- | --- | --- |",
            "| 山区血液运输项目 | 医疗运输 | 面向县域急救与样本转运 |",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "无人机医疗运输项目.xlsx",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "方案/案例")
        self.assertEqual(profile["模板归属"], "方案案例模板")

    def test_infer_document_profile_marks_excel_project_amount_table_as_solution(self) -> None:
        blocks = [
            "工作表：Sheet1",
            "| 列1 | 项目名称 | 项目内容 | 项目金额 |",
            "| --- | --- | --- | --- |",
            "| 1 | 2024年宝安区中心血站低空经济示范应用场景-无人机血液运输服务项目 | 无人机血液运输服务 | 1821000.0元 |",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "无人机医疗运输项目.xlsx",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "方案/案例")
        self.assertEqual(profile["模板归属"], "方案案例模板")

    def test_build_structured_payload_keeps_event_material_fields_conservative(self) -> None:
        contract = load_contract()
        sample = SampleRecord(
            sample_id="scan_xiamen_event_markdown",
            source_path="E:/低空经济资源库/应用场景/2026厦门国际低空经济暨无人机产业博览会（2026年4月23-25日）.md",
            document_category="待判定资料",
            recommended_template="待人工补规则",
            title_hint="2026厦门国际低空经济暨无人机产业博览会（2026年4月23-25日）",
            subject_name_hint="",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L3",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="text:utf-8-sig",
            extraction_status="已提取文本",
            extracted_text=(
                "# 2026厦门国际低空经济暨无人机产业博览会（2026年4月23-25日）\n\n"
                "2026厦门国际低空经济暨无人机产业博览会\n"
                "时间:4月23-25日 地点:厦门国际会展中心\n"
                "展会背景：围绕低空经济创新发展搭建交流合作平台。\n"
                "近年来，我国政府高度重视低空经济的发展，出台了一系列政策措施，推动低空经济的规范化、规模化发展。\n"
                "同期活动：\n"
                "《低空安全监管技术应用交流会》\n"
            ),
            preview_text="2026厦门国际低空经济暨无人机产业博览会",
            text_length=120,
            page_count=None,
            source_encoding="utf-8-sig",
            note="已完成文本文件读取。已清洗网页导航壳层。",
        )

        with patch("minimum_workflow.pipeline.extract_with_strategy", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)

        self.assertEqual(payload["推荐模板"], "方案案例模板")
        self.assertEqual(payload["标题"], "2026厦门国际低空经济暨无人机产业博览会")
        self.assertEqual(payload["文件标题"], "2026厦门国际低空经济暨无人机产业博览会")
        self.assertEqual(payload["主体名称"], "")
        self.assertEqual(payload["所属场景字段"], "")
        self.assertEqual(payload["客户名称字段"], "")
        self.assertEqual(payload["文件日期字段"], "")
        self.assertEqual(payload["分流结果"], "待审核")

    def test_infer_document_profile_marks_low_altitude_opportunity_catalog_as_solution(self) -> None:
        blocks = [
            "重庆市2025年第一批低空经济应用场景机会清单",
            "场景类型序号场景机会名称单位名称合作需求联系方式",
            "低空安全保障 全市无人机无线信标接收系统和重点区域无人机监测管控 重庆市公安局",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "附件1：重庆市2025年第一批低空经济应用场景机会清单.docx",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "方案/案例")
        self.assertEqual(profile["模板归属"], "方案案例模板")
        self.assertEqual(profile["资料层级"], "方案/案例资料")
        self.assertEqual(profile["是否适合直接入Dify"], "待审核")

    def test_infer_document_profile_marks_low_altitude_capability_catalog_as_solution(self) -> None:
        blocks = [
            "重庆市2025年第一批低空经济应用场景能力清单",
            "场景类型序号场景能力名称单位名称场景能力说明联系方式",
            "低空智联网与安全保障 基于卫星互联网的低空安全监管与运营服务",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "附件2：重庆市2025年第一批低空经济应用场景能力清单.docx",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "方案/案例")
        self.assertEqual(profile["模板归属"], "方案案例模板")
        self.assertEqual(profile["资料层级"], "方案/案例资料")
        self.assertEqual(profile["是否适合直接入Dify"], "待审核")

    def test_infer_document_title_prefers_supplier_cover_company_name_over_catalog_heading(self) -> None:
        title = sample_docx_extract_to_md.infer_document_title(
            [
                "杭州昊舜视讯科技有限公司",
                "公司介绍",
                "科技 创新 合作 共赢",
                "公司介绍\n产品介绍\n行业应用介绍\n项目案例",
            ],
            "公司介绍202504.pptx",
        )

        self.assertEqual(title, "杭州昊舜视讯科技有限公司")

    def test_infer_document_title_keeps_full_source_title_for_split_training_cover(self) -> None:
        title = sample_docx_extract_to_md.infer_document_title(
            [
                "北京大学",
                "低空经济与区域经济发展",
                "高级研修班（三期）",
                "北大培训20250862号",
                "课程安排",
            ],
            "北京大学低空经济与区域经济发展高级研修班（三期）.pdf",
        )

        self.assertEqual(title, "北京大学低空经济与区域经济发展高级研修班（三期）")

    def test_infer_document_profile_marks_procurement_scoring_table_as_procurement(self) -> None:
        blocks = [
            "| 1 | 项目概况描述 | 2分 | 投标供应商根据项目具体情况,对项目理解进行描述。 |",
            "| --- | --- | --- | --- |",
            "| 2 | 服务范围、内容 | 2分 | 投标供应商根据项目具体情况及项目概况理解对服务范围及内容进行分析描述。 |",
            "| 3 | 防治专项方案 | 3分 | 未提供方案或提供的方案不适用本项目的得0分。 |",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "微信图片_20260409162848_366_69.jpg",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "招标/采购文件")
        self.assertEqual(profile["模板归属"], "招标采购文件模板")

    def test_infer_document_profile_marks_low_altitude_weather_article_as_solution(self) -> None:
        blocks = [
            "低空经济+气象监测预警系统：3大核心能力、4类技术路线与5大区域落地实践",
            "2024年12月，中国气象局联合国家数据局率先部署。",
            "本文将从系统架构、技术路线与实战案例几个维度，全面解析低空气象监测预警系统。",
            "低空气象监测预警系统的核心能力，依赖感知层、算法层、应用层三大模块的协同运作。",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "低空经济+气象.docx",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "方案/案例")
        self.assertEqual(profile["模板归属"], "方案案例模板")

    def test_infer_document_profile_marks_enterprise_list_as_supplier_template(self) -> None:
        blocks = [
            "十大低空经济企业",
            "十大低空物流企业：丰翼科技(深圳)有限公司、杭州迅蚁网络科技有限公司、电鹰科技集团有限公司。",
            "十大低空应急救援企业：广州市华科尔科技股份有限公司、深圳联合飞机科技有限公司。",
            "十大低空运营服务企业：广州中科云图智能科技有限公司、星逻智能科技(苏州)有限公司。",
            "十大低空安全设备企业：北京历正科技有限责任公司、杰能科世智能安全科技(杭州)有限公司。",
            "十大低空经济链上企业：深圳市好盈科技股份有限公司、广州成至智能机器科技有限公司。",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "十大低空经济企业.docx",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "供应商/企业资料")
        self.assertEqual(profile["模板归属"], "供应商企业模板")
        self.assertEqual(profile["资料层级"], "供应商企业资料")
        self.assertEqual(profile["证据边界"], "企业名录/榜单类资料，企业入选范围、排序口径与真实性需结合原件复核。")
        self.assertEqual(profile["是否适合直接入Dify"], "待审核")

    def test_infer_document_profile_marks_peking_university_training_brochure(self) -> None:
        blocks = [
            "北京大学",
            "低空经济与区域经济发展",
            "高级研修班（三期）",
            "北大培训20250862号",
            "课程安排",
            "招生对象",
            "学习费用",
            "课程时间",
            "报名程序",
            "结业证书",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "北京大学低空经济与区域经济发展高级研修班（三期）.pdf",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "教育培训")
        self.assertEqual(profile["模板归属"], "教育培训模板")
        self.assertEqual(profile["文件标题"], "北京大学低空经济与区域经济发展高级研修班（三期）")
        self.assertEqual(profile["是否适合直接入Dify"], "待审核")

    def test_infer_document_profile_marks_wushan_plum_case_as_solution(self) -> None:
        blocks = [
            "2.3.3 无人机物流-应用案例",
            "(1) 重庆巫山脆李案例",
            "重庆巫山探索的‘即时响应+无人机转运+极速鲜’寄递物流体系。",
            "TD550最大载重200kg，任务半径200km，巡航速度可达120～140km/h。",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "重庆巫山脆李无人机运输案例.pdf",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "方案/案例")
        self.assertEqual(profile["模板归属"], "方案案例模板")
        self.assertEqual(profile["资料层级"], "方案/案例资料")
        self.assertEqual(profile["是否适合直接入Dify"], "待审核")

    def test_infer_document_profile_blanks_event_organizer_for_meeting_notice(self) -> None:
        blocks = [
            "2026年低空资源和经济发展大会暨低空经济发展工作委员会年度工作会议",
            "会议时间：2026年1月17日（星期六）9:00～18:00",
            "会议地点：北京市昌平区未来科学城万怡酒店",
            "主办单位",
            "中国AOPA低空经济发展工作委员会",
            "协办单位",
            "北京大学空天信息工程研究中心",
        ]

        profile = sample_docx_extract_to_md.infer_document_profile(
            "2026年低空资源和经济发展大会暨低空经济发展工作委员会年度工作会议.md",
            blocks,
        )

        self.assertEqual(profile["文档分类"], "方案/案例")
        self.assertEqual(profile["模板归属"], "方案案例模板")
        self.assertEqual(profile["主体名称"], "")

    def test_build_structured_payload_keeps_news_markdown_title_and_blank_solution_fields(self) -> None:
        contract = load_contract()
        sample = SampleRecord(
            sample_id="scan_ninghe_news_markdown",
            source_path="E:/低空经济资源库/应用场景/天津宁河区：低空巡检让宁河湿地保护“耳聪目明”.md",
            document_category="待判定资料",
            recommended_template="待人工补规则",
            title_hint="天津宁河区：低空巡检让宁河湿地保护“耳聪目明”",
            subject_name_hint="",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L3",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="text:utf-8-sig",
            extraction_status="已提取文本",
            extracted_text=(
                "# 天津宁河区：低空巡检让宁河湿地保护“耳聪目明”\n\n"
                "天津宁河区：低空巡检让宁河湿地保护“耳聪目明”\n"
                "创建时间：2025-05-22 17:11:04\n"
                "初夏的清晨，七里海湿地自然保护区内一架无人机划破天际。\n"
                "陈兆睿告诉记者，为实现对湿地全域常态化日常巡检及对生态系统的高效、精准监控。\n"
                "目前，在七里海湿地低空巡检技术已实现智能化、常态化应用，工作人员能更加全面细致地掌握湿地内的风吹草动。\n"
                "来源：中国环境APP\n"
            ),
            preview_text="天津宁河区：低空巡检让宁河湿地保护“耳聪目明”",
            text_length=160,
            page_count=None,
            source_encoding="utf-8-sig",
            note="已完成文本文件读取。已清洗网页导航壳层。",
        )

        with patch("minimum_workflow.pipeline.extract_with_strategy", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract)

        self.assertEqual(payload["标题"], "天津宁河区：低空巡检让宁河湿地保护“耳聪目明”")
        self.assertEqual(payload["文件标题"], "天津宁河区：低空巡检让宁河湿地保护“耳聪目明”")
        self.assertEqual(payload["方案名称字段"], "")
        self.assertEqual(payload["客户名称字段"], "")
        self.assertEqual(payload["文件日期字段"], "")
        self.assertEqual(payload["主体名称"], "")

    def test_education_training_fields_fallback_to_keyword_lines(self) -> None:
        sample = SampleRecord(
            sample_id="scan_education_training_keyword_fallback",
            source_path="D:/长风资料/天翼航空相关资料/低空经济背景下的无人机专业建设及人才培养模式.pdf",
            document_category="教育培训",
            recommended_template="教育培训模板",
            title_hint="低空经济背景下的无人机专业建设及人才培养模式",
            subject_name_hint="天翼航空相关资料",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L2",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="pypdf",
            extraction_status="已提取文本",
            extracted_text=(
                "主讲人：汪湖滨\n"
                "无人机专业建设及人才培养模式\n"
                "服务全国本科、高职、中职及技师院校共\n"
                "整体解决方案可为客户提供7大服务：课程体系、师资培训、人才服务、资格认证服务、赛事服务、定制服务、教学环境搭建。\n"
                "精品微课、实操视频、录屏课程、虚拟仿真、2D/3D动画、实训手册、任务工卡、电子教材、PPT、教学大纲、电子教案、考核题库、实际行业案例等。\n"
                "未来低空经济产业人才培养模式探索\n"
            ),
            preview_text="无人机专业建设及人才培养模式",
            text_length=140,
            page_count=12,
            source_encoding="utf-8",
            note="已尝试使用 pypdf 提取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["适用对象字段"], "本科、高职、中职及技师院校")
        self.assertEqual(payload["课程体系字段"], "整体解决方案可为客户提供7大服务：课程体系、师资培训、人才服务、资格认证服务、赛事服务、定制服务、教学环境搭建。")
        self.assertEqual(payload["核心内容字段"], "未来低空经济产业人才培养模式探索")

    def test_education_training_fields_ignore_yunshenchu_noise_lines(self) -> None:
        sample = SampleRecord(
            sample_id="scan_yunshenchu_education_noise",
            source_path="D:/长风资料/云深处公司产品资料/云深处科技-产教融合2026.pptx",
            document_category="教育培训",
            recommended_template="教育培训模板",
            title_hint="产教融合",
            subject_name_hint="产教融合",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L3",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="presentation:pptx",
            extraction_status="已提取文本",
            extracted_text=(
                "公司介绍\n"
                "产教融合\n"
                "杭州市B类高层次人才创新创业项目\n"
                "打造集教学、实训、科研、竞赛于一体的综合性平台\n"
                "XXX-云深处具身智能（某方向）联合实验室\n"
                "立足行业，赋能教育，促进产教融合发展！\n"
            ),
            preview_text="产教融合",
            text_length=100,
            page_count=20,
            source_encoding="pptx",
            note="已完成演示文稿读取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["文件标题"], "产教融合")
        self.assertEqual(payload["培训主题字段"], "产教融合")
        self.assertEqual(payload["适用对象字段"], "")
        self.assertEqual(payload["专业方向字段"], "")
        self.assertEqual(payload["课程体系字段"], "打造集教学、实训、科研、竞赛于一体的综合性平台")
        self.assertEqual(payload["核心内容字段"], "立足行业，赋能教育，促进产教融合发展！")

    def test_education_training_fields_ignore_generic_section_heading_and_fallback_title_hint(self) -> None:
        sample = SampleRecord(
            sample_id="scan_wujin_construction_heading_fallback",
            source_path="D:/长风资料/天翼航空相关资料/高质量产教融合无人机专业五金建设方案.pdf",
            document_category="教育培训",
            recommended_template="教育培训模板",
            title_hint="高质量产教融合无人机专业五金建设方案",
            subject_name_hint="天翼航空相关资料",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L2",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="pypdf",
            extraction_status="已提取文本",
            extracted_text=(
                "模块一 Module 1 国家政策\n"
                "模块二 Module 2 企业介绍\n"
                "模块三 Module 3 建设方向\n"
                "模块四 Module 4 专业教学数字化转型\n"
                "模块五 Module 5 课程建设体系图\n"
                "模块六 Module 6 教学资源开发流程\n"
                "模块七 Module 7 部分建设案例\n"
            ),
            preview_text="模块一 Module 1 国家政策",
            text_length=120,
            page_count=41,
            source_encoding="utf-8",
            note="已尝试使用 pypdf 提取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["文件标题"], "高质量产教融合无人机专业五金建设方案")
        self.assertEqual(payload["培训主题字段"], "")
        self.assertEqual(payload["专业方向字段"], "")
        self.assertEqual(payload["课程体系字段"], "课程建设体系图")

    def test_education_training_fields_prefer_main_title_over_inline_reference_title(self) -> None:
        sample = SampleRecord(
            sample_id="scan_education_training_realistic_pdf",
            source_path="D:/长风资料/天翼航空相关资料/低空经济背景下的无人机专业建设及人才培养模式.pdf",
            document_category="教育培训",
            recommended_template="教育培训模板",
            title_hint="低空经济背景下的无人机专业建设及人才培养模式",
            subject_name_hint="天翼航空相关资料",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L2",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="pypdf",
            extraction_status="已提取文本",
            extracted_text=(
                "无人机专业建设及人才培养模式\n"
                "低空经济背景下的\n"
                "主讲人：\n"
                "l 《无人机应用技术专业》国家级教学资源库参建单位\n"
                "服务全国本科、高职、中职及技师院校共\n"
                "整体解决方案可为客户提供7大服务：课程体系、师资培训、人才服务、资格认证服务、赛事服务、定制服务、教学环境搭建。7大\n"
                "l 未来低空经济产业人才培养模式探索\n"
                "l 产教融合下无人机创新教育课程体系建设\n"
                "l 基于智能飞行器应用技术的创新型专业师资培训\n"
            ),
            preview_text="无人机专业建设及人才培养模式",
            text_length=240,
            page_count=70,
            source_encoding="utf-8",
            note="已尝试使用 pypdf 提取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["文件标题"], "无人机专业建设及人才培养模式")
        self.assertEqual(payload["培训主题字段"], "")
        self.assertEqual(payload["专业方向字段"], "无人机应用技术专业")
        self.assertEqual(payload["适用对象字段"], "本科、高职、中职及技师院校")
        self.assertEqual(payload["课程体系字段"], "整体解决方案可为客户提供7大服务：课程体系、师资培训、人才服务、资格认证服务、赛事服务、定制服务、教学环境搭建。")
        self.assertEqual(payload["核心内容字段"], "未来低空经济产业人才培养模式探索；产教融合下无人机创新教育课程体系建设；基于智能飞行器应用技术的创新型专业师资培训")

    def test_education_training_fields_strip_markdown_heading_and_ignore_course_goal_noise(self) -> None:
        sample = SampleRecord(
            sample_id="scan_education_training_mineru_heading",
            source_path="D:/长风资料/天翼航空相关资料/低空经济背景下的无人机专业建设及人才培养模式.pdf",
            document_category="教育培训",
            recommended_template="教育培训模板",
            title_hint="低空经济背景下的无人机专业建设及人才培养模式",
            subject_name_hint="天翼航空相关资料",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L2",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="mineru:batch",
            extraction_status="已提取文本",
            extracted_text=(
                "# 低空经济背景下的无人机专业建设及人才培养模式\n"
                "主讲人：西安天翼智飞科技集团有限公司\n"
                "整体解决方案可为客户提供7大服务：课程体系、师资培训、人才服务、资格认证服务、赛事服务、定制服务、教学环境搭建。7大服务凝聚为4大产品。\n"
                "8.能不断积累经验，并从中归纳并找出夹性的能力；为以后的课程打下坚实的基础。\n"
                "未来低空经济产业人才培养模式探索\n"
            ),
            preview_text="# 低空经济背景下的无人机专业建设及人才培养模式",
            text_length=170,
            page_count=70,
            source_encoding="utf-8",
            note="已通过 MinerU 批量接口完成 Markdown 提取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["文件标题"], "低空经济背景下的无人机专业建设及人才培养模式")
        self.assertEqual(payload["培训主题字段"], "")
        self.assertEqual(payload["课程体系字段"], "整体解决方案可为客户提供7大服务：课程体系、师资培训、人才服务、资格认证服务、赛事服务、定制服务、教学环境搭建。")
        self.assertEqual(payload["核心内容字段"], "未来低空经济产业人才培养模式探索")

    def test_education_training_fields_keep_full_brochure_title_and_multiline_audience(self) -> None:
        sample = SampleRecord(
            sample_id="scan_peking_training_brochure_fields",
            source_path="D:/长风资料/11、行业内公司/北京大学低空经济与区域经济发展高级研修班（三期）.pdf",
            document_category="教育培训",
            recommended_template="教育培训模板",
            title_hint="北京大学低空经济与区域经济发展高级研修班（三期）",
            subject_name_hint="11、行业内公司",
            product_name_hint="",
            unit_name_hint="北京大学",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L3",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="mineru:batch",
            extraction_status="已提取文本",
            extracted_text=(
                "# 北京大学\n\n"
                "# 低空经济与区域经济发展\n\n"
                "# 高级研修班（三期）\n\n"
                "# 课程涵盖\n\n"
                "# 发展路径\n\n"
                "★战略新兴产业多维立体化发展方向\n"
                "★低空经济政策解读及重点省份发展落实情况\n"
                "★各地方低空经济产业发展重点分析\n\n"
                "# 产业生态\n\n"
                "★低空基础设施建设进展情况\n"
                "★低空应用场景与服务新业态解析\n"
                "★低空经济产业链条热点布局\n\n"
                "# 机遇与挑战\n\n"
                "★无人机、直升机等人才短板及培训\n"
                "★急需解决的挑战与应对方式\n\n"
                "# 低空经济名企参访\n\n"
                "★低空算力中心一曙光云计算集团股份有限公司\n\n"
                "课程安排\n\n"
                "《绿色低碳转型与战略新兴产业发展》\n"
                "《低空经济三个基础要素与引领国际标准的中国突破》\n"
                "《低空经济综合飞行管控技术应用》\n"
                "《低空路网基础设施建设的前沿科技进展》\n"
                "《区域低空经济产业发展与推进路径》\n"
                "《基于卫星数据链的低空飞行器主动监测系统》\n"
                "《低空经济-改变中国经济的新引擎》\n"
                "《低空经济发展和地方创新实践》\n\n"
                "# 招生对象\n\n"
                "从事低空经济相关产业的人员\n\n"
                "对低空经济以及创新创造感兴趣的企业家和个人\n\n"
                "# 课程时间\n\n"
                "2025年7月18日至21日\n"
            ),
            preview_text="# 北京大学 # 低空经济与区域经济发展 # 高级研修班（三期）",
            text_length=120,
            page_count=10,
            source_encoding="utf-8",
            note="已通过 MinerU 批量接口完成 Markdown 提取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["文件标题"], "北京大学低空经济与区域经济发展高级研修班（三期）")
        self.assertEqual(payload["培训主题字段"], "北京大学低空经济与区域经济发展高级研修班（三期）")
        self.assertEqual(payload["适用对象字段"], "从事低空经济相关产业的人员")
        self.assertEqual(payload["专业方向字段"], "低空经济；区域经济")
        self.assertEqual(
            payload["课程体系字段"],
            "发展路径：战略新兴产业多维立体化发展方向；发展路径：低空经济政策解读及重点省份发展落实情况；发展路径：各地方低空经济产业发展重点分析；产业生态：低空基础设施建设进展情况；产业生态：低空应用场景与服务新业态解析；产业生态：低空经济产业链条热点布局；机遇与挑战：无人机、直升机等人才短板及培训；机遇与挑战：急需解决的挑战与应对方式；低空经济名企参访：低空算力中心一曙光云计算集团股份有限公司",
        )
        self.assertEqual(
            payload["核心内容字段"],
            "绿色低碳转型与战略新兴产业发展；低空经济三个基础要素与引领国际标准的中国突破；低空经济综合飞行管控技术应用；低空路网基础设施建设的前沿科技进展；区域低空经济产业发展与推进路径；基于卫星数据链的低空飞行器主动监测系统；低空经济-改变中国经济的新引擎；低空经济发展和地方创新实践",
        )

    def test_education_training_fields_extract_topic_labels_from_single_line_brochure_title(self) -> None:
        sample = SampleRecord(
            sample_id="scan_training_single_line_topic_title",
            source_path="D:/长风资料/教育培训方向/北京航空航天大学低空经济与通用航空协同发展高级研修班（一期）.pdf",
            document_category="教育培训",
            recommended_template="教育培训模板",
            title_hint="北京航空航天大学低空经济与通用航空协同发展高级研修班（一期）",
            subject_name_hint="教育培训方向",
            product_name_hint="",
            unit_name_hint="北京航空航天大学",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L3",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="mineru:batch",
            extraction_status="已提取文本",
            extracted_text=(
                "# 北京航空航天大学低空经济与通用航空协同发展高级研修班（一期）\n\n"
                "# 课程涵盖\n\n"
                "# 核心价值\n\n"
                "★核心价值体系构建与产业升级路径\n\n"
                "# 招生对象\n\n"
                "低空经济相关从业人员\n"
            ),
            preview_text="# 北京航空航天大学低空经济与通用航空协同发展高级研修班（一期）",
            text_length=96,
            page_count=6,
            source_encoding="utf-8",
            note="已通过 MinerU 批量接口完成 Markdown 提取。",
        )

        payload = extract_fields(sample, extraction)

        self.assertEqual(payload["专业方向字段"], "低空经济；通用航空")
        self.assertEqual(payload["适用对象字段"], "低空经济相关从业人员")
        self.assertNotIn("核心价值", payload["专业方向字段"])
    def test_extract_image_directory_content_runs_ocr_and_builds_auto_metadata(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        source_dir = Path(temp_dir.name) / "低空智能网联体系参考架构（2024版）工信部装备工业发展中心2024-11-11"
        source_dir.mkdir()
        page_paths = []
        for name in ("00封面.png", "01.png", "02.jpeg"):
            path = source_dir / name
            path.write_bytes(b"fake-image")
            page_paths.append(path)

        batch_result = {
            "batch_id": "batch-dir-1",
            "results": [
                {"state": "done", "markdown": "# 低空智能网联体系参考架构（2024版）", "error": ""},
                {"state": "done", "markdown": "工业和信息化部装备工业发展中心", "error": ""},
                {"state": "done", "markdown": "2024年11月", "error": ""},
            ],
        }

        with patch.object(sample_docx_extract_to_md, "resolve_mineru_token", return_value="token-demo"), patch.object(
            sample_docx_extract_to_md,
            "run_mineru_batch",
            return_value=batch_result,
        ) as run_mock:
            result = sample_docx_extract_to_md.extract_image_directory_content(source_dir)

        run_mock.assert_called_once()
        self.assertEqual(result["blocks"][0], "# 低空智能网联体系参考架构（2024版）")
        self.assertIn(("来源形态", "分页扫描图片目录"), result["auto_metadata"])
        self.assertIn(("目录判定", "分页扫描文档目录"), result["auto_metadata"])
        self.assertEqual(result["auto_payload"]["标题"], "低空智能网联体系参考架构（2024版）")
        self.assertEqual(result["auto_payload"]["OCR页数"], "3")
        self.assertIn("batch-dir-1", result["extraction_note"])
        self.assertIn("分页扫描目录", result["selection_reason"])

    def test_extract_image_directory_content_marks_wechat_meeting_scan_for_review(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        source_dir = Path(temp_dir.name) / "迅蚁2026-04-08线上会"
        source_dir.mkdir()
        for name in ("微信图片_202604080001.png", "微信图片_202604080002.png", "微信图片_202604080003.png"):
            (source_dir / name).write_bytes(b"fake-image")

        batch_result = {
            "batch_id": "batch-dir-wechat-1",
            "results": [
                {"state": "done", "markdown": "# 迅蚁科技低空教育项目线上说明会", "error": ""},
                {"state": "done", "markdown": "主讲人：汪湖滨\n请勿外传，仅供内部交流", "error": ""},
                {"state": "done", "markdown": "教育部办公厅关于做好2026年职业教育拟招生专业设置管理工作的通知", "error": ""},
            ],
        }

        with patch.object(sample_docx_extract_to_md, "resolve_mineru_token", return_value="token-demo"), patch.object(
            sample_docx_extract_to_md,
            "run_mineru_batch",
            return_value=batch_result,
        ):
            result = sample_docx_extract_to_md.extract_image_directory_content(source_dir)

        self.assertEqual(result["auto_payload"]["目录判定"], "分页扫描文档目录")
        self.assertEqual(result["auto_payload"]["标题"], "迅蚁科技低空教育项目线上说明会")
        self.assertEqual(result["auto_payload"]["文件标题"], "迅蚁科技低空教育项目线上说明会")
        self.assertEqual(result["auto_payload"]["文档分类"], "方案/案例")
        self.assertEqual(result["auto_payload"]["推荐模板"], "方案案例模板")
        self.assertFalse(result["auto_payload"]["是否适合直接入库"])
        self.assertEqual(result["auto_payload"]["分流结果"], "待审核")
        self.assertEqual(result["auto_payload"]["发布时间"], "")

    def test_build_structured_payload_keeps_wechat_meeting_fields_conservative(self) -> None:
        contract = load_contract()
        sample = SampleRecord(
            sample_id="wechat_meeting_scan",
            source_path=r"D:\长风资料\迅蚁2026-04-08线上会",
            document_category="待判定资料",
            recommended_template="待人工补规则",
            title_hint="迅蚁2026-04-08线上会",
            subject_name_hint="",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=["目录扫描自动生成样本，模板归属与字段需结合原文复核。"],
            notes=["由目录扫描入口自动生成。"],
            evidence_level="L3",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
        )
        extraction = ExtractionResult(
            extractor_name="ocr:mineru:image_directory",
            extraction_status="已提取文本",
            extracted_text=(
                "迅蚁科技教育事业部·2026年4月8日\n"
                "# 迅蚁科技低空教育项目线上说明会\n"
                "主讲人：汪湖滨\n"
                "请勿外传，仅供内部交流\n"
                "教育部办公厅关于做好2026年职业教育拟招生专业设置管理工作的通知\n"
                "B6迅蚁无人机机场\n"
                "报价说明\n"
            ),
            preview_text="迅蚁科技低空教育项目线上说明会",
            text_length=120,
            page_count=3,
            source_encoding="utf-8",
            note="分页扫描目录已按页序完成 OCR；分页命名图片 3/3；文档型图片 3/3；批次号：batch-dir-wechat-1。",
            extra_metadata={
                "文档分类": "方案/案例",
                "推荐模板": "方案案例模板",
                "模板归属": "方案案例模板",
                "标题": "迅蚁科技低空教育项目线上说明会",
                "文件标题": "迅蚁科技低空教育项目线上说明会",
                "资料层级": "方案/案例资料",
                "证据边界": "会议/PPT/说明会类混合材料，夹带政策截图或引用页；不按单一政策官方文件直入，需结合原件与用途复核。",
                "来源形态": "分页扫描图片目录",
                "目录判定": "分页扫描文档目录",
                "判定依据": "分页命名图片 3/3；文档型图片 3/3",
                "OCR页数": "3",
                "OCR结果概况": "done=3",
                "OCR失败页": "无",
                "取舍说明": "该目录被判定为同一文档的分页扫描目录，已按页序统一 OCR 提取，不按纯照片目录跳过。",
                "分流结果": "待审核",
                "是否适合直接入库": False,
            },
        )

        with patch("minimum_workflow.pipeline.extract_with_strategy", return_value=extraction):
            payload, _ = build_structured_payload(sample, contract, enable_qwen=True, mineru_token="token-demo")

        self.assertEqual(payload["处理路径"], "ocr")
        self.assertEqual(payload["目录判定"], "分页扫描文档目录")
        self.assertEqual(payload["标题"], "迅蚁科技低空教育项目线上说明会")
        self.assertEqual(payload["文件标题"], "迅蚁科技低空教育项目线上说明会")
        self.assertEqual(payload["方案名称字段"], "")
        self.assertEqual(payload["客户名称字段"], "")
        self.assertEqual(payload["预算组织字段"], "")
        self.assertEqual(payload["推荐模板"], "方案案例模板")
        self.assertEqual(payload["分流结果"], "待审核")

    def test_sample_docx_main_merges_directory_auto_metadata(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        source_dir = Path(temp_dir.name) / "低空智能网联体系参考架构（2024版）工信部装备工业发展中心2024-11-11"
        source_dir.mkdir()
        output_path = Path(temp_dir.name) / "directory.md"
        args = argparse.Namespace(
            source=str(source_dir),
            output=str(output_path),
            pdf_output_mode="fulltext",
            qwen_api_key=None,
            qwen_base_url=None,
            qwen_model=None,
            discarded=[],
            meta=[],
            related_duplicates=[],
            related_policies=[],
            related_images=[],
            selection_reason="人工默认取舍说明。",
        )
        extraction = {
            "blocks": ["# 低空智能网联体系参考架构（2024版）", "正文第一段"],
            "extracted_text": "# 低空智能网联体系参考架构（2024版）\n\n正文第一段",
            "is_heavy_pdf": False,
            "heavy_pdf_reason": "",
            "extraction_note": "目录OCR成功",
            "auto_metadata": [
                ("文档分类", "行业参考架构/指导材料"),
                ("模板归属", "参考架构/白皮书口径（当前按原文全量提取Markdown输出）"),
                ("来源形态", "分页扫描图片目录"),
                ("目录判定", "分页扫描文档目录"),
            ],
            "selection_reason": "目录自动判定成功。",
        }

        with patch.object(sample_docx_extract_to_md, "parse_args", return_value=args), patch.object(
            sample_docx_extract_to_md,
            "extract_source_content",
            return_value=extraction,
        ):
            sample_docx_extract_to_md.main()

        output_text = output_path.read_text(encoding="utf-8")
        self.assertIn("文档分类: 行业参考架构/指导材料", output_text)
        self.assertIn("模板归属: 参考架构/白皮书口径（当前按原文全量提取Markdown输出）", output_text)
        self.assertIn("来源形态: 分页扫描图片目录", output_text)
        self.assertIn("目录判定: 分页扫描文档目录", output_text)
        self.assertIn("目录自动判定成功。", output_text)

    def test_build_structured_payload_markdown_preserves_existing_frontmatter_hints(self) -> None:
        contract = load_contract()
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        source_path = Path(temp_dir.name) / "已整理方案.md"
        source_path.write_text(
            "---\n"
            "推荐模板: 方案案例模板\n"
            "文档分类: 方案/案例\n"
            "自定义字段: 保留值\n"
            "---\n\n"
            "# 已整理方案\n\n"
            "客户/使用单位：贵州某医院\n"
            "正文内容保持不动。\n",
            encoding="utf-8",
        )
        sample = SampleRecord(
            sample_id="scan_existing_markdown_frontmatter",
            source_path=str(source_path),
            document_category="待判定资料",
            recommended_template="待人工补规则",
            title_hint="已整理方案",
            subject_name_hint="",
            product_name_hint="",
            unit_name_hint="",
            tags=["目录扫描", "自动判型"],
            risks=[],
            notes=[],
            evidence_level="L3",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
            relative_path_hint="已整理/已整理方案.md",
        )
        extraction = ExtractionResult(
            extractor_name="text:utf-8-sig",
            extraction_status="已提取文本",
            extracted_text=source_path.read_text(encoding="utf-8"),
            preview_text="已整理方案 贵州某医院",
            text_length=80,
            page_count=None,
            source_encoding="utf-8",
            note="已完成文本文件读取。",
        )

        with patch("minimum_workflow.pipeline.extract_text", return_value=extraction):
            payload, parsed_extraction = build_structured_payload(sample, contract)

        self.assertEqual(payload["文件类型"], "markdown")
        self.assertEqual(payload["推荐模板"], "方案案例模板")
        self.assertEqual(payload["文档分类"], "方案/案例")
        self.assertEqual(payload["原始Markdown元数据"]["自定义字段"], "保留值")
        self.assertNotIn("推荐模板: 方案案例模板", parsed_extraction.extracted_text)
        self.assertTrue(parsed_extraction.extracted_text.startswith("# 已整理方案"))

    def test_build_markdown_for_markdown_source_supplements_frontmatter_without_rewriting_body(self) -> None:
        sample = SampleRecord(
            sample_id="scan_markdown_preserve_body",
            source_path="D:/长风资料/已整理/企业方案.md",
            document_category="方案/案例",
            recommended_template="方案案例模板",
            title_hint="企业方案",
            subject_name_hint="长风科技",
            product_name_hint="",
            unit_name_hint="长风科技",
            tags=["目录扫描"],
            risks=[],
            notes=[],
            evidence_level="L3",
            fallback_decision="待审核",
            split_required=False,
            split_note="",
            relative_path_hint="已整理/企业方案.md",
        )
        payload = {
            "原始文件名": "企业方案.md",
            "原始路径": "D:/长风资料/已整理/企业方案.md",
            "文件格式": "md",
            "文件类型": "markdown",
            "处理路径": "text_direct",
            "文档分类": "方案/案例",
            "推荐模板": "方案案例模板",
            "主体名称": "长风科技",
            "单位名称": "长风科技",
            "标题": "企业方案",
            "文件标题": "企业方案",
            "提取正文": "# 企业方案\n\n这里是人工整理好的正文。\n- 第一条\n- 第二条\n",
            "方案名称字段": "企业方案",
            "客户名称字段": "贵阳某单位",
            "文件日期字段": "2026年04月",
            "生成时间": "2026-04-21T12:34:56",
            "原始Markdown元数据": {"自定义字段": "保留值"},
        }

        markdown = build_markdown(sample, payload)

        self.assertTrue(markdown.startswith("---\n"))
        self.assertIn("源文件名: 企业方案.md", markdown)
        self.assertIn("源文件相对路径: 已整理/企业方案.md", markdown)
        self.assertIn("文档类别: solution", markdown)
        self.assertIn("文档分类: 方案/案例", markdown)
        self.assertIn("推荐模板: 方案案例模板", markdown)
        self.assertIn("项目名称: 企业方案", markdown)
        self.assertIn("客户/牵头单位: 贵阳某单位", markdown)
        self.assertIn("编制单位: 长风科技", markdown)
        self.assertIn("编制日期: 2026年04月", markdown)
        self.assertIn("自定义字段: 保留值", markdown)
        self.assertIn("\n# 企业方案\n\n这里是人工整理好的正文。\n- 第一条\n- 第二条\n", markdown)
        self.assertNotIn("## 一、文件摘要", markdown)

    def test_build_scanned_sample_keeps_relative_path_hint(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        source_dir = Path(temp_dir.name) / "source"
        nested_dir = source_dir / "子目录"
        nested_dir.mkdir(parents=True)
        doc_path = nested_dir / "已整理文档.md"
        doc_path.write_text("# 内容", encoding="utf-8")

        sample = build_scanned_sample(doc_path, source_dir)

        self.assertEqual(sample.relative_path_hint, "子目录/已整理文档.md")

    def test_sample_docx_main_merges_directory_auto_metadata(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        source_dir = Path(temp_dir.name) / "低空智能网联体系参考架构（2024版）工信部装备工业发展中心2024-11-11"
        source_dir.mkdir()
        output_path = Path(temp_dir.name) / "directory.md"
        args = argparse.Namespace(
            source=str(source_dir),
            output=str(output_path),
            pdf_output_mode="fulltext",
            qwen_api_key=None,
            qwen_base_url=None,
            qwen_model=None,
            discarded=[],
            meta=[],
            related_duplicates=[],
            related_policies=[],
            related_images=[],
            selection_reason="人工默认取舍说明。",
        )
        extraction = {
            "blocks": ["# 低空智能网联体系参考架构（2024版）", "正文第一段"],
            "extracted_text": "# 低空智能网联体系参考架构（2024版）\n\n正文第一段",
            "is_heavy_pdf": False,
            "heavy_pdf_reason": "",
            "extraction_note": "目录OCR成功",
            "auto_metadata": [
                ("文档分类", "行业参考架构/指导材料"),
                ("模板归属", "参考架构/白皮书口径（当前按原文全量提取Markdown输出）"),
                ("来源形态", "分页扫描图片目录"),
                ("目录判定", "分页扫描文档目录"),
            ],
            "selection_reason": "目录自动判定成功。",
        }

        with patch.object(sample_docx_extract_to_md, "parse_args", return_value=args), patch.object(
            sample_docx_extract_to_md,
            "extract_source_content",
            return_value=extraction,
        ):
            sample_docx_extract_to_md.main()

        output_text = output_path.read_text(encoding="utf-8")
        self.assertIn("文档分类: 行业参考架构/指导材料", output_text)
        self.assertIn("模板归属: 参考架构/白皮书口径（当前按原文全量提取Markdown输出）", output_text)
        self.assertIn("来源形态: 分页扫描图片目录", output_text)
        self.assertIn("目录判定: 分页扫描文档目录", output_text)
        self.assertIn("目录自动判定成功。", output_text)

    def test_detect_heavy_pdf_layout_flags_watermark_and_parameter_markers(self) -> None:
        text = (
            "感谢您下载包图网平台上提供的PPT作品 ibaotu.com 请勿复制、传播、销售 "
            "环境温度 目标温度设为 货箱容积 适配机型 温控运输箱"
        )
        blocks = ["17.5L", "货箱容积", "适配机型", "环境温度"] * 8

        is_heavy, reason = sample_docx_extract_to_md.detect_heavy_pdf_layout(text, blocks)

        self.assertTrue(is_heavy)
        self.assertIn("命中重版式水印", reason)
        self.assertIn("命中图表/参数页特征词", reason)

    def test_should_use_pdf_summary_mode_handles_fulltext_summary_and_auto(self) -> None:
        pdf_path = Path("D:/tmp/heavy.pdf")
        docx_path = Path("D:/tmp/sample.docx")

        self.assertTrue(sample_docx_extract_to_md.should_use_pdf_summary_mode(pdf_path, "summary", False))
        self.assertTrue(sample_docx_extract_to_md.should_use_pdf_summary_mode(pdf_path, "auto", True))
        self.assertFalse(sample_docx_extract_to_md.should_use_pdf_summary_mode(pdf_path, "auto", False))
        self.assertFalse(sample_docx_extract_to_md.should_use_pdf_summary_mode(pdf_path, "fulltext", True))
        self.assertFalse(sample_docx_extract_to_md.should_use_pdf_summary_mode(docx_path, "summary", True))

    def test_has_meaningful_summary_payload_distinguishes_empty_and_non_empty_payloads(self) -> None:
        self.assertFalse(sample_docx_extract_to_md.has_meaningful_summary_payload({"资料摘要": "", "应用背景": []}))
        self.assertTrue(sample_docx_extract_to_md.has_meaningful_summary_payload({"资料摘要": "有摘要", "应用背景": []}))
        self.assertTrue(sample_docx_extract_to_md.has_meaningful_summary_payload({"资料摘要": "", "应用背景": ["背景"]}))

    def test_sample_docx_main_auto_summary_mode_writes_summary_markdown(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        source_path = Path(temp_dir.name) / "heavy.pdf"
        output_path = Path(temp_dir.name) / "summary.md"
        source_path.write_bytes(b"%PDF-1.4\n")
        args = argparse.Namespace(
            source=str(source_path),
            output=str(output_path),
            pdf_output_mode="auto",
            qwen_api_key=None,
            qwen_base_url=None,
            qwen_model=None,
            discarded=[],
            meta=[],
            related_duplicates=[],
            related_policies=[],
            related_images=[],
            selection_reason="保留重版式主文件。",
        )
        extraction = {
            "blocks": ["参数页碎片"],
            "extracted_text": "感谢您下载包图网平台上提供的PPT作品 环境温度 目标温度设为 货箱容积 适配机型",
            "is_heavy_pdf": True,
            "heavy_pdf_reason": "命中重版式水印：感谢您下载包图网平台上提供的PPT作品、ibaotu.com",
            "extraction_note": "已回退本地 PDF 解析。",
        }
        summary_payload = normalize_solution_summary_payload(
            {
                "主体名称": "杭州迅蚁/送吧航空",
                "方案名称/案例名称": "六盘水医院无人机运输方案",
                "所属场景": "医疗物流",
                "客户/使用单位": "钟山区人民医院",
                "文件日期": "2025-12-08",
                "资料摘要": "方案摘要",
                "资料形态判断": "商业方案",
                "应用背景": ["背景一"],
                "解决的问题": ["问题一"],
                "投入的产品/设备/能力": ["TR9S"],
                "实施方式": ["航线规划"],
                "预算、进度与组织方式": ["采购服务"],
                "结果与效果数据": ["时效提升"],
                "可复用经验": ["经验一"],
                "入库与归档判断": ["适合摘要入库"],
                "备注": ["需回看原PDF"],
            }
        )

        with patch.object(sample_docx_extract_to_md, "parse_args", return_value=args), patch.object(
            sample_docx_extract_to_md,
            "extract_source_content",
            return_value=extraction,
        ), patch.object(
            sample_docx_extract_to_md,
            "resolve_qwen_runtime",
            return_value={"api_key": "demo-key", "base_url": "https://demo", "model": "qwen-plus"},
        ), patch.object(
            sample_docx_extract_to_md,
            "summarize_solution_document_with_qwen",
            return_value=summary_payload,
        ) as summary_mock:
            sample_docx_extract_to_md.main()

        summary_mock.assert_called_once()
        output_text = output_path.read_text(encoding="utf-8")
        self.assertIn("输出类型: 模型摘要提取Markdown", output_text)
        self.assertIn("# 一、资料摘要", output_text)
        self.assertIn("- 方案摘要", output_text)
        self.assertIn("摘要触发方式: auto 模式命中重版式特征", output_text)
        self.assertIn("# 十、入库与归档判断", output_text)

    def test_sample_docx_main_fulltext_mode_skips_qwen_summary(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        source_path = Path(temp_dir.name) / "normal.pdf"
        output_path = Path(temp_dir.name) / "fulltext.md"
        source_path.write_bytes(b"%PDF-1.4\n")
        args = argparse.Namespace(
            source=str(source_path),
            output=str(output_path),
            pdf_output_mode="fulltext",
            qwen_api_key=None,
            qwen_base_url=None,
            qwen_model=None,
            discarded=[],
            meta=[],
            related_duplicates=[],
            related_policies=[],
            related_images=[],
            selection_reason="保留普通正文版。",
        )
        extraction = {
            "blocks": ["第一段正文", "第二段正文"],
            "extracted_text": "第一段正文\n\n第二段正文",
            "is_heavy_pdf": True,
            "heavy_pdf_reason": "命中图表特征",
            "extraction_note": "已完成本地解析。",
        }

        with patch.object(sample_docx_extract_to_md, "parse_args", return_value=args), patch.object(
            sample_docx_extract_to_md,
            "extract_source_content",
            return_value=extraction,
        ), patch.object(sample_docx_extract_to_md, "resolve_qwen_runtime") as runtime_mock, patch.object(
            sample_docx_extract_to_md,
            "summarize_solution_document_with_qwen",
        ) as summary_mock:
            sample_docx_extract_to_md.main()

        runtime_mock.assert_not_called()
        summary_mock.assert_not_called()
        output_text = output_path.read_text(encoding="utf-8")
        self.assertIn("输出类型: 原文全量提取Markdown", output_text)
        self.assertIn("# 原文全文", output_text)
        self.assertIn("第一段正文", output_text)
        self.assertNotIn("# 一、资料摘要", output_text)

    def test_web_helper_validate_source_url_rejects_private_host(self) -> None:
        with self.assertRaisesRegex(ValueError, "不允许抓取本机或内网地址"):
            knowledge_helper_ui._validate_source_url("http://127.0.0.1/internal")

        with self.assertRaisesRegex(ValueError, "网址中不能包含空格或换行"):
            knowledge_helper_ui._validate_source_url("https://example.org/demo\nnext")

        with patch.object(knowledge_helper_ui, "_resolve_hostname_addresses", return_value={"127.0.0.1"}):
            with self.assertRaisesRegex(ValueError, "不允许抓取本机或内网地址"):
                knowledge_helper_ui._validate_source_url("https://example.com/demo")

        self.assertEqual(
            knowledge_helper_ui._validate_source_url("https://example.org/demo"),
            "https://example.org/demo",
        )

    def test_web_helper_fetch_webpage_response_rejects_non_html_content(self) -> None:
        response = Mock()
        response.status_code = 200
        response.headers = {"Content-Type": "application/pdf"}
        response.raise_for_status = Mock()

        session = Mock()
        session.get.return_value = response

        with patch.object(knowledge_helper_ui.requests, "Session", return_value=session):
            with self.assertRaisesRegex(ValueError, "仅支持可直接访问的 HTML 页面"):
                knowledge_helper_ui._fetch_webpage_response("https://example.org/demo.pdf")

    def test_web_helper_build_upload_copy_plan_rejects_colliding_filenames(self) -> None:
        base_dir = Path("D:/sandbox/input")
        uploaded_files = [
            argparse.Namespace(name="D:/tmp/a/demo.txt", orig_name="demo.txt"),
            argparse.Namespace(name="D:/tmp/b/demo.txt", orig_name="demo.txt"),
        ]

        with self.assertRaisesRegex(ValueError, "检测到同名文件冲突"):
            knowledge_helper_ui._build_upload_copy_plan(base_dir, uploaded_files)

    def test_web_helper_build_upload_copy_plan_marks_flattened_structure(self) -> None:
        base_dir = Path("D:/sandbox/input")
        uploaded_files = [
            argparse.Namespace(name="D:/tmp/a/demo1.txt", orig_name="demo1.txt"),
            argparse.Namespace(name="D:/tmp/b/demo2.txt", orig_name="demo2.txt"),
        ]

        copy_plan, has_flattened_structure = knowledge_helper_ui._build_upload_copy_plan(base_dir, uploaded_files)

        self.assertEqual(
            copy_plan,
            [
                (Path("D:/tmp/a/demo1.txt"), base_dir / "demo1.txt"),
                (Path("D:/tmp/b/demo2.txt"), base_dir / "demo2.txt"),
            ],
        )
        self.assertTrue(has_flattened_structure)

    def test_web_helper_resolve_server_port_uses_available_port(self) -> None:
        with patch.object(knowledge_helper_ui, "_is_port_available", side_effect=[False, True]):
            with patch.dict(os.environ, {"GRADIO_SERVER_PORT": "7899"}, clear=False):
                self.assertEqual(knowledge_helper_ui._resolve_server_port(), 7900)

        with patch.object(knowledge_helper_ui, "_is_port_available", return_value=True):
            with patch.dict(os.environ, {"GRADIO_SERVER_PORT": "bad-port"}, clear=False):
                self.assertEqual(knowledge_helper_ui._resolve_server_port(), 7861)

    def test_web_helper_toggle_source_inputs_switches_visibility(self) -> None:
        upload_update, url_update, status_html = knowledge_helper_ui._toggle_source_inputs("url")

        self.assertFalse(upload_update["visible"])
        self.assertTrue(url_update["visible"])
        self.assertIn("网址模式", status_html["value"])

    def test_web_helper_render_dashboard_reports_invalid_dify_url_without_crashing(self) -> None:
        with patch.object(knowledge_helper_ui, "build_batch_choices", return_value=[]):
            updates = knowledge_helper_ui._render_dashboard(
                "",
                "ftp://192.168.110.78:17001/v1",
                "secret",
                "",
                True,
            )

        self.assertIn("HTTP", updates[1])
        self.assertEqual(updates[2], "未找到可用批次。")

    def test_web_helper_save_review_ignores_invalid_dify_url_for_local_review_save(self) -> None:
        with patch.object(knowledge_helper_ui, "save_manual_review") as save_mock, patch.object(
            knowledge_helper_ui,
            "_render_dashboard",
            return_value=("dashboard",),
        ) as render_mock:
            result = knowledge_helper_ui.save_review_and_refresh(
                "D:/demo/batch_001",
                "sample-1",
                "已有分类",
                "",
                ["dataset-1"],
                "ftp://192.168.110.78:17001/v1",
                "secret",
                "",
                True,
            )

        self.assertEqual(result, ("dashboard",))
        save_mock.assert_called_once()
        self.assertIsNone(save_mock.call_args.kwargs["runtime"])
        render_mock.assert_called_once()

    def test_web_helper_import_ready_rejects_invalid_dify_url(self) -> None:
        with patch.object(knowledge_helper_ui, "_render_dashboard", return_value=("dashboard",)) as render_mock, patch.object(
            knowledge_helper_ui,
            "import_ready_documents",
        ) as import_mock:
            result = knowledge_helper_ui.import_ready_and_refresh(
                "D:/demo/batch_001",
                "ftp://192.168.110.78:17001/v1",
                "secret",
                "",
                True,
            )

        self.assertEqual(result, ("dashboard",))
        import_mock.assert_not_called()
        render_mock.assert_called_once()
        self.assertIn("HTTP", render_mock.call_args.kwargs["status_message"])

    def test_web_helper_run_scan_timeout_fails_and_cleans_sandbox(self) -> None:
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)

        upload_path = Path(temp_dir.name) / "demo.txt"
        upload_path.write_text("demo", encoding="utf-8")
        sandbox_path = Path(temp_dir.name) / "sandbox"
        sandbox_path.mkdir()
        uploaded_file = argparse.Namespace(name=str(upload_path), orig_name="nested/demo.txt")

        class FakeQueue:
            def __init__(self, *args, **kwargs) -> None:
                pass

            def get(self, timeout: float | None = None):
                raise knowledge_helper_ui.queue.Empty

            def put(self, item) -> None:
                pass

        class FakeThread:
            def __init__(self, target=None, args=(), daemon: bool | None = None) -> None:
                self.target = target
                self.args = args
                self.daemon = daemon

            def start(self) -> None:
                return None

        class FakeProc:
            def __init__(self) -> None:
                self.stdout = []
                self.returncode = None
                self.terminated = False
                self.killed = False

            def poll(self):
                return self.returncode

            def terminate(self) -> None:
                self.terminated = True
                self.returncode = -15

            def wait(self, timeout: float | None = None):
                return self.returncode

            def kill(self) -> None:
                self.killed = True
                self.returncode = -9

        fake_proc = FakeProc()

        with patch.object(knowledge_helper_ui.tempfile, "mkdtemp", return_value=str(sandbox_path)), patch.object(
            knowledge_helper_ui.subprocess,
            "Popen",
            return_value=fake_proc,
        ), patch.object(knowledge_helper_ui.threading, "Thread", FakeThread), patch.object(
            knowledge_helper_ui.queue,
            "Queue",
            FakeQueue,
        ), patch.object(knowledge_helper_ui, "_build_download_archive") as archive_mock:
            events = list(
                knowledge_helper_ui.run_scan(
                    [uploaded_file],
                    "upload",
                    "",
                    "local",
                    False,
                    False,
                    "",
                    "",
                    "",
                    "",
                    False,
                )
            )

        archive_mock.assert_not_called()
        self.assertTrue(fake_proc.terminated)
        self.assertFalse(fake_proc.killed)
        self.assertFalse(sandbox_path.exists())
        self.assertIn("处理超时", events[-1][0])
        self.assertIn("未生成可下载结果", events[-1][0])


class MineruLargeFileSplitTest(unittest.TestCase):
    def _make_pdf(self, path: Path, pages: int) -> None:
        from pypdf import PdfWriter
        writer = PdfWriter()
        for _ in range(pages):
            writer.add_blank_page(width=100, height=100)
        with path.open("wb") as f:
            writer.write(f)

    def test_count_pdf_pages(self) -> None:
        from minimum_workflow.mineru_large_file import count_pdf_pages
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        pdf_path = Path(temp_dir.name) / "sample.pdf"
        self._make_pdf(pdf_path, 7)
        self.assertEqual(count_pdf_pages(pdf_path), 7)

    def test_should_use_split_strategy_by_pages(self) -> None:
        from minimum_workflow.mineru_large_file import should_use_split_strategy
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        pdf_path = Path(temp_dir.name) / "many.pdf"
        self._make_pdf(pdf_path, 150)
        self.assertTrue(should_use_split_strategy(pdf_path, "pdf"))

    def test_should_not_split_small_pdf(self) -> None:
        from minimum_workflow.mineru_large_file import should_use_split_strategy
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        pdf_path = Path(temp_dir.name) / "small.pdf"
        self._make_pdf(pdf_path, 10)
        self.assertFalse(should_use_split_strategy(pdf_path, "pdf"))

    def test_split_pdf_by_pages(self) -> None:
        from minimum_workflow.mineru_large_file import split_pdf_by_pages, count_pdf_pages
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        tmp_path = Path(temp_dir.name)
        pdf_path = tmp_path / "big.pdf"
        self._make_pdf(pdf_path, 250)
        chunks = split_pdf_by_pages(pdf_path, tmp_path / "out", pages_per_chunk=180)
        self.assertEqual(len(chunks), 2)
        self.assertEqual(count_pdf_pages(chunks[0]), 180)
        self.assertEqual(count_pdf_pages(chunks[1]), 70)
        self.assertTrue(chunks[0].name.endswith("_part01.pdf"))
        self.assertTrue(chunks[1].name.endswith("_part02.pdf"))

    def test_split_pdf_returns_original_when_under_threshold(self) -> None:
        from minimum_workflow.mineru_large_file import split_pdf_by_pages
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        tmp_path = Path(temp_dir.name)
        pdf_path = tmp_path / "small.pdf"
        self._make_pdf(pdf_path, 50)
        chunks = split_pdf_by_pages(pdf_path, tmp_path / "out", pages_per_chunk=180)
        self.assertEqual(chunks, [pdf_path])

    def test_extract_large_pdf_via_split_combines_md(self) -> None:
        from minimum_workflow.mineru_large_file import extract_large_file_via_split
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        tmp_path = Path(temp_dir.name)
        pdf_path = tmp_path / "report.pdf"
        self._make_pdf(pdf_path, 220)

        fake_batch = {
            "batch_id": "bt-123",
            "results": [
                {"state": "done", "markdown": "# Part One\n\n第一片正文内容足够长用于通过语义校验" * 5},
                {"state": "done", "markdown": "# Part Two\n\n第二片正文内容足够长用于通过语义校验" * 5},
            ],
        }
        with patch(
            "minimum_workflow.extractors.run_mineru_batch",
            return_value=fake_batch,
        ) as batch_mock:
            result = extract_large_file_via_split(pdf_path, "pdf", "token-demo")

        batch_mock.assert_called_once()
        call_args, call_kwargs = batch_mock.call_args
        chunks_passed = call_args[0]
        self.assertEqual(len(chunks_passed), 2)
        self.assertEqual(result.extractor_name, "mineru:split:pdf")
        self.assertEqual(result.extraction_status, "已提取文本")
        self.assertIn("Part One", result.extracted_text)
        self.assertIn("Part Two", result.extracted_text)
        self.assertIn("分片：", result.extracted_text)
        self.assertEqual(result.page_count, 220)
        self.assertIn("bt-123", result.note)


class CleanMineruMarkdownTableTest(unittest.TestCase):
    def test_html_table_converted_to_markdown_pipes(self) -> None:
        from minimum_workflow.extractors import clean_mineru_markdown
        html = (
            "<table><thead><tr><th>序号</th><th>检验项目</th><th>标准要求</th>"
            "<th>检验结果</th><th>备注</th></tr></thead>"
            "<tbody>"
            "<tr><td>1</td><td>额定流量(L/s)</td><td>≥16.0</td><td>16.04</td><td>吸深3m</td></tr>"
            "<tr><td>2</td><td>额定压力(MPa)</td><td>≥0.60</td><td>0.60</td><td> </td></tr>"
            "</tbody></table>"
        )
        result = clean_mineru_markdown(html)
        self.assertIn("| 序号 | 检验项目 | 标准要求 | 检验结果 | 备注 |", result)
        self.assertIn("| --- | --- | --- | --- | --- |", result)
        self.assertIn("| 1 | 额定流量(L/s) | ≥16.0 | 16.04 | 吸深3m |", result)
        self.assertIn("| 2 | 额定压力(MPa) | ≥0.60 | 0.60 |", result)

    def test_split_text_skips_html_comment_blocks(self) -> None:
        from minimum_workflow.document_profiles import split_text_to_blocks
        text = (
            "<!-- 分片：金沙_part01.pdf -->\n\n"
            "# 检验报告\n\n"
            "正文第一段内容"
        )
        blocks = split_text_to_blocks(text)
        self.assertEqual(len(blocks), 2)
        self.assertTrue(blocks[0].startswith("#"))
        self.assertIn("正文第一段", blocks[1])


if __name__ == "__main__":
    unittest.main()
